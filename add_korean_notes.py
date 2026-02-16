#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""PPTX에 한글 노트 추가 및 제목 슬라이드 생성"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pathlib import Path

if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8')

# 한글 요약 내용 (섹션별)
KOREAN_NOTES = [
    "n8n 마스터 가이드 - AI 통합 워크플로우 자동화 완벽 가이드",
    
    """n8n 개요:
- 페어 코드(fair-code) 라이선스 기반 워크플로우 자동화 도구
- API가 있는 모든 앱 연결 가능
- 셀프 호스팅으로 데이터 프라이버시 강화""",
    
    """호스팅 옵션:
- GCP 무료 티어: e2-micro 인스턴스 + 30GB 디스크
- 로컬 Docker: Cloudflare 터널링으로 외부 연결
- Railway: 월 $5로 24시간 운영""",
    
    """노드 종류:
- 트리거: 워크플로우 시작점 (번개 아이콘)
- 액션: 데이터 조작 및 외부 작업
- 코어 노드: If, Switch, HTTP Request 등
- 클러스터 노드: 복잡한 기능 수행""",
    
    """AI 통합 (MCP):
- MCP Server Trigger: 워크플로우를 AI 도구로 노출
- MCP Client Tool: 외부 도구 사용
- Claude Desktop과 연동 가능""",
    
    """LangChain 통합:
- LLM 체인, 요약 체인, 정보 추출기
- 벡터 스토어: Pinecone, Supabase, Weaviate
- RAG 시스템 구축 가능""",
    
    """보안 관리:
- 자격 증명 암호화 저장
- SSL, SSO, 2FA 설정
- RBAC 역할 기반 접근 제어""",
]

def add_korean_slides(input_path, output_path):
    """한글 제목 슬라이드 추가 및 노트 삽입"""
    prs = Presentation(input_path)
    
    # 첫 번째 슬라이드에 제목 오버레이 (노트로 추가)
    print("한글 노트 추가 중...")
    
    for i, slide in enumerate(prs.slides):
        if i < len(KOREAN_NOTES):
            notes = slide.notes_slide
            notes.notes_text_frame.text = KOREAN_NOTES[i]
        
        # 진행률 표시
        if (i + 1) % 10 == 0:
            print(f"  {i+1}/{len(prs.slides)} 슬라이드")
    
    prs.save(output_path)
    print(f"\n완료: {output_path}")
    print(f"총 {len(prs.slides)} 슬라이드")

def main():
    downloads = Path("downloads")
    input_pptx = downloads / "n8n_마스터_가이드_통합.pptx"
    output_pptx = downloads / "n8n_마스터_가이드_한글.pptx"
    
    add_korean_slides(input_pptx, output_pptx)

if __name__ == "__main__":
    main()
