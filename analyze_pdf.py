"""
analyze_pdf.py - 독립형 PDF/PPTX 분석 모듈 (Google Vision OCR)

웹앱 없이 단독 실행 가능. 다른 프로젝트에서도 재사용 가능.

사용법:
    from analyze_pdf import analyze_file

    result = analyze_file("path/to/file.pdf")
    # result = {
    #     "title": "...",
    #     "category": "disease",
    #     "tags": [...],
    #     "summary": "...",
    #     "content": "...",
    #     "analyzedBy": "Vision OCR" | "PDF 텍스트 추출" | "PPTX 텍스트 추출",
    #     "pageCount": 10,
    # }

CLI:
    python analyze_pdf.py "path/to/file.pdf"

환경변수:
    GOOGLE_VISION_API_KEY  - Google Vision API 키 (이미지 PDF일 때 필요)

의존성:
    pip install pdfplumber python-pptx pdf2image Pillow requests
    + Ghostscript 설치 필요 (이미지 PDF → OCR 시)
"""

import base64
import io
import json
import os
import re
import sys

import requests

# ── 설정 ────────────────────────────────────────────────────
VISION_API_URL = "https://vision.googleapis.com/v1/images:annotate"
MIN_TEXT_LENGTH = 100  # 이 이하면 이미지 기반 PDF로 판단
MAX_OCR_PAGES = 5

# Ghostscript 경로 (Windows)
GS_BIN = r"C:\Users\antigravity\scoop\apps\ghostscript\current\bin"
if os.path.exists(GS_BIN) and GS_BIN not in os.environ.get("PATH", ""):
    os.environ["PATH"] = GS_BIN + ";" + os.environ.get("PATH", "")

# 부위별 분류 체계 (웹앱 body-parts.ts와 동일)
BODY_PARTS = {
    "어깨": ["어깨", "견관절", "회전근개", "오십견", "석회성건염", "유착성피막염"],
    "팔꿈치": ["팔꿈치", "주관절", "테니스엘보", "골퍼엘보", "상과염"],
    "손/손목": ["손목", "손가락", "수근관", "방아쇠", "건초염", "터널증후군"],
    "고관절": ["고관절", "대퇴골두", "무혈성괴사", "대퇴골"],
    "무릎": ["무릎", "슬관절", "십자인대", "반월상", "연골"],
    "발/발목": ["발목", "족저근막", "아킬레스", "발바닥", "발뒤꿈치"],
    "척추": ["척추", "허리", "목", "디스크", "협착", "측만", "경추", "요추"],
}

# 전체 의학 키워드 (부위 키워드 + 일반 키워드)
GENERAL_KEYWORDS = [
    "관절", "예방", "운동", "건강", "치료", "증상", "원인",
    "진단", "수술", "재활", "통증", "염증", "근육", "뼈", "신경",
    "당뇨", "고혈압", "심장", "폐", "간", "위", "장", "피부",
    "스트레칭", "골절", "인대", "골다공증",
]

MEDICAL_KEYWORDS = list({kw for keywords in BODY_PARTS.values() for kw in keywords} | set(GENERAL_KEYWORDS))


# ── Google Vision OCR ───────────────────────────────────────
def _get_vision_api_key():
    """환경변수 또는 .env.local에서 Vision API 키 로드"""
    key = os.environ.get("GOOGLE_VISION_API_KEY", "")
    if key:
        return key

    # .env.local 파일 탐색 (현재 디렉토리 → 상위 디렉토리)
    search_dir = os.getcwd()
    for _ in range(5):
        env_path = os.path.join(search_dir, ".env.local")
        if os.path.exists(env_path):
            with open(env_path, encoding="utf-8") as f:
                for line in f:
                    if line.startswith("GOOGLE_VISION_API_KEY="):
                        return line.split("=", 1)[1].strip()
        parent = os.path.dirname(search_dir)
        if parent == search_dir:
            break
        search_dir = parent

    return ""


def ocr_image_with_vision(image_base64: str, api_key: str) -> str:
    """Google Vision API로 이미지 OCR"""
    payload = {
        "requests": [{
            "image": {"content": image_base64},
            "features": [{"type": "DOCUMENT_TEXT_DETECTION"}],
        }]
    }
    resp = requests.post(
        f"{VISION_API_URL}?key={api_key}",
        json=payload,
        timeout=30,
    )
    result = resp.json()

    if "error" in result:
        raise RuntimeError(f"Vision API: {result['error'].get('message', '')}")

    return result.get("responses", [{}])[0].get("fullTextAnnotation", {}).get("text", "")


# ── PDF 텍스트 추출 ─────────────────────────────────────────
def extract_text_from_pdf(file_path: str) -> tuple[str, int]:
    """pdfplumber로 PDF 텍스트 추출 (텍스트 기반 PDF용)"""
    import pdfplumber

    text_parts = []
    page_count = 0

    with pdfplumber.open(file_path) as pdf:
        page_count = len(pdf.pages)
        for page in pdf.pages:
            page_text = page.extract_text() or ""
            if page_text.strip():
                text_parts.append(page_text.strip())

    return "\n\n".join(text_parts), page_count


# ── PDF → 이미지 변환 (Ghostscript 사용) ────────────────────
def _pdf_pages_to_images(file_path: str, max_pages: int = MAX_OCR_PAGES) -> list[str]:
    """Ghostscript로 PDF 페이지를 PNG 이미지 파일로 변환. 경로 리스트 반환."""
    import subprocess
    import tempfile

    out_dir = tempfile.mkdtemp(prefix="nlm_ocr_")
    out_pattern = os.path.join(out_dir, "page_%03d.png")

    gs_cmd = [
        "gswin64c" if sys.platform == "win32" else "gs",
        "-dNOPAUSE", "-dBATCH", "-dSAFER", "-dQUIET",
        "-sDEVICE=png16m",
        "-r150",
        f"-dFirstPage=1",
        f"-dLastPage={max_pages}",
        f"-sOutputFile={out_pattern}",
        file_path,
    ]

    subprocess.run(gs_cmd, check=True, capture_output=True, timeout=120)

    # 생성된 파일 수집
    pages = sorted([
        os.path.join(out_dir, f)
        for f in os.listdir(out_dir)
        if f.endswith(".png")
    ])
    return pages


# ── PDF 이미지 OCR ──────────────────────────────────────────
def ocr_pdf_with_vision(file_path: str, api_key: str) -> str:
    """이미지 기반 PDF를 Ghostscript + Vision OCR로 텍스트 추출"""

    page_images = _pdf_pages_to_images(file_path)

    full_text = ""
    for i, img_path in enumerate(page_images):
        with open(img_path, "rb") as f:
            b64 = base64.b64encode(f.read()).decode("utf-8")

        page_text = ocr_image_with_vision(b64, api_key)
        if page_text:
            full_text += f"\n--- 페이지 {i + 1} ---\n{page_text}"

        # 임시 파일 삭제
        try:
            os.remove(img_path)
        except OSError:
            pass

    return full_text.strip()


# ── PPTX 텍스트 추출 ────────────────────────────────────────
def extract_text_from_pptx(file_path: str) -> tuple[str, int]:
    """python-pptx로 PPTX 텍스트 추출"""
    from pptx import Presentation

    prs = Presentation(file_path)
    text_parts = []

    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
        if slide_texts:
            text_parts.append(f"--- 슬라이드 {i} ---\n" + "\n".join(slide_texts))

    return "\n\n".join(text_parts), len(prs.slides)


# ── 텍스트 정리 ─────────────────────────────────────────────
def _clean_ocr_text(text: str) -> str:
    """OCR 텍스트에서 NotebookLM 워터마크 등 불필요한 텍스트 제거"""
    # NotebookLM 워터마크 제거
    text = re.sub(r"\b(Start\s+)?NotebookLM\b", "", text)
    # 페이지/슬라이드 구분자를 마크다운 구분선으로 변환
    text = re.sub(
        r"---\s*(페이지|슬라이드)\s*(\d+)\s*---",
        r"\n\n### \1 \2\n",
        text,
    )
    # 연속 빈줄 정리
    text = re.sub(r"\n{3,}", "\n\n", text)
    # Source: 줄 제거 (NotebookLM 출처 표시)
    text = re.sub(r"Source:.*$", "", text, flags=re.MULTILINE)
    return text.strip()


# ── 텍스트 파싱 ─────────────────────────────────────────────
def parse_extracted_text(text: str, file_name: str) -> dict:
    """추출된 텍스트에서 제목/요약/본문/태그 파싱"""
    cleaned = _clean_ocr_text(text)
    lines = [l.strip() for l in cleaned.split("\n") if l.strip()]

    # 제목: 3~50자 사이의 첫 의미있는 줄 (### 마크다운 제외)
    title = ""
    for line in lines:
        if line.startswith("#"):
            continue
        if 3 <= len(line) <= 50 and not re.match(r"^[\d\s.\-]+$", line):
            title = line
            break
    if not title:
        title = re.sub(r"\.[^/.]+$", "", file_name)

    # 요약: 처음 10줄에서 300자 (### 제외)
    content_lines = [l for l in lines if not l.startswith("#")]
    summary_text = " ".join(content_lines[:10])
    summary = summary_text[:300].strip() or "내용을 확인해주세요"

    # 본문 (정리된 텍스트)
    content = f"## 내용\n\n{cleaned[:15000]}"

    # 태그: 부위 대표 태그 + 의학 키워드 매칭
    # 매칭 횟수가 가장 많은 부위를 선택 (첫 매칭이 아닌 최다 매칭)
    lower_text = text.lower()
    body_part_tag = None
    best_count = 0
    for part_name, keywords in BODY_PARTS.items():
        count = sum(1 for kw in keywords if kw in lower_text)
        if count > best_count:
            best_count = count
            body_part_tag = part_name

    detail_tags = [kw for kw in MEDICAL_KEYWORDS if kw in lower_text]
    # 부위 대표 태그를 첫 번째에 삽입, 나머지는 상세 태그
    if body_part_tag:
        tags = [body_part_tag] + [t for t in detail_tags if t != body_part_tag][:4]
    else:
        tags = detail_tags[:5]
    if not tags:
        tags = ["자동추출"]

    return {
        "title": title,
        "summary": summary,
        "content": content,
        "tags": tags,
    }


# ── 메인 분석 함수 ──────────────────────────────────────────
def analyze_file(file_path: str) -> dict:
    """
    PDF 또는 PPTX 파일을 분석하여 제목/요약/본문/태그 추출.

    Returns:
        {
            "title": str,
            "category": "disease",
            "tags": list[str],
            "summary": str,
            "content": str,
            "analyzedBy": str,
            "pageCount": int,
        }
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"파일을 찾을 수 없습니다: {file_path}")

    ext = os.path.splitext(file_path)[1].lower()
    file_name = os.path.basename(file_path)

    if ext == ".pptx":
        text, page_count = extract_text_from_pptx(file_path)
        method = f"PPTX 텍스트 추출 ({page_count}슬라이드)"

    elif ext == ".pdf":
        # 1단계: 텍스트 추출 시도
        text, page_count = extract_text_from_pdf(file_path)

        if len(text.strip()) >= MIN_TEXT_LENGTH:
            method = "PDF 텍스트 추출"
        else:
            # 2단계: 이미지 PDF → Vision OCR
            api_key = _get_vision_api_key()
            if not api_key:
                raise RuntimeError(
                    "GOOGLE_VISION_API_KEY가 설정되지 않았습니다. "
                    "환경변수 또는 .env.local에 설정하세요."
                )
            text = ocr_pdf_with_vision(file_path, api_key)
            method = "Vision OCR"
    else:
        raise ValueError(f"지원하지 않는 파일 형식: {ext} (PDF, PPTX만 지원)")

    if not text.strip():
        return {
            "title": re.sub(r"\.[^/.]+$", "", file_name),
            "category": "disease",
            "tags": ["검토필요"],
            "summary": "텍스트 추출 실패",
            "content": "## 내용\n\n파일에서 텍스트를 추출할 수 없습니다.",
            "analyzedBy": method,
            "pageCount": page_count if ext == ".pdf" else 0,
        }

    parsed = parse_extracted_text(text, file_name)

    return {
        "title": parsed["title"],
        "category": "disease",
        "tags": parsed["tags"],
        "summary": parsed["summary"],
        "content": parsed["content"],
        "analyzedBy": method,
        "pageCount": page_count if ext != ".pptx" else 0,
    }


# ── CLI ─────────────────────────────────────────────────────
if __name__ == "__main__":
    if sys.platform == "win32" and sys.stdout.encoding != "utf-8":
        sys.stdout.reconfigure(encoding="utf-8")
        sys.stderr.reconfigure(encoding="utf-8")

    if len(sys.argv) < 2:
        print("사용법: python analyze_pdf.py <파일경로>")
        sys.exit(1)

    path = sys.argv[1]
    try:
        result = analyze_file(path)
        print(json.dumps(result, ensure_ascii=False, indent=2))
    except Exception as e:
        print(f"오류: {e}", file=sys.stderr)
        sys.exit(1)
