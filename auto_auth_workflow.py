#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
NotebookLM 자동 인증 워크플로우
- Google 앱 비밀번호로 로그인
- 쿠키 추출 및 저장
- nlm CLI 인증 설정
"""
import asyncio
import json
import sys
import os
import io
import re
import time
from pathlib import Path
from datetime import datetime

os.environ['PYTHONUNBUFFERED'] = '1'
if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8', line_buffering=True)

# 설정
AUTH_DIR = Path.home() / ".notebooklm-mcp-cli"
PROFILE_DIR = AUTH_DIR / "profiles" / "default"
DOWNLOAD_DIR = Path("G:/내 드라이브/notebooklm")

# 인증 정보 (환경변수에서 가져옴)
GOOGLE_EMAIL = os.environ.get("GOOGLE_EMAIL", "")
APP_PASSWORD = os.environ.get("GOOGLE_APP_PASSWORD", "")
APIFY_API_KEY = os.environ.get("APIFY_API_KEY", "")


def log(msg):
    print(msg, flush=True)


def safe_filename(name):
    """파일명에 사용할 수 없는 문자 제거"""
    return re.sub(r'[<>:"/\\|?*]', '', name).strip()


async def google_login_with_app_password(page, email: str, app_password: str) -> bool:
    """Google 앱 비밀번호로 로그인"""
    log("  [1] 이메일 입력...")

    # 이메일 입력
    email_input = await page.wait_for_selector('input[type="email"]', timeout=10000)
    await email_input.fill(email)
    await page.keyboard.press('Enter')
    await asyncio.sleep(3)

    # 비밀번호 입력 (앱 비밀번호)
    log("  [2] 앱 비밀번호 입력...")
    try:
        password_input = await page.wait_for_selector('input[type="password"]', timeout=10000)
        # 앱 비밀번호는 공백 제거
        clean_password = app_password.replace(' ', '')
        await password_input.fill(clean_password)
        await page.keyboard.press('Enter')
        await asyncio.sleep(5)
    except Exception as e:
        log(f"  ⚠️ 비밀번호 입력 오류: {e}")
        return False

    # 로그인 확인
    log("  [3] 로그인 확인...")
    for _ in range(10):
        await asyncio.sleep(2)
        if "notebooklm.google.com" in page.url and "accounts.google" not in page.url:
            log("  ✓ 로그인 성공!")
            return True

        # 2단계 인증 확인
        if "challenge" in page.url or "signin/v2/challenge" in page.url:
            log("  ⚠️ 2단계 인증 필요 - 수동 처리 대기...")
            for _ in range(30):  # 2.5분 대기
                await asyncio.sleep(5)
                if "notebooklm.google.com" in page.url:
                    return True

    return False


async def extract_and_save_cookies(context) -> dict:
    """쿠키 추출 및 저장"""
    cookies = await context.cookies()

    # 쿠키를 딕셔너리로 변환
    cookies_dict = {}
    for cookie in cookies:
        if 'google.com' in cookie.get('domain', ''):
            cookies_dict[cookie['name']] = cookie['value']

    # CSRF 토큰 생성
    sapisid = cookies_dict.get('SAPISID', cookies_dict.get('__Secure-3PAPISID', ''))
    csrf_token = f"{sapisid[:16]}:{int(time.time() * 1000)}" if sapisid else ""

    # auth.json 저장
    auth_data = {
        "cookies": cookies_dict,
        "csrf_token": csrf_token,
        "session_id": "",
        "extracted_at": time.time(),
        "auto_login": True
    }

    AUTH_DIR.mkdir(parents=True, exist_ok=True)
    PROFILE_DIR.mkdir(parents=True, exist_ok=True)

    # 루트 auth.json
    with open(AUTH_DIR / "auth.json", 'w', encoding='utf-8') as f:
        json.dump(auth_data, f, indent=2)

    # 프로필 디렉토리
    cookies_list = [
        {
            "name": name,
            "value": value,
            "domain": ".google.com",
            "path": "/",
            "expires": -1,
            "httpOnly": False,
            "secure": True,
            "sameSite": "Lax"
        }
        for name, value in cookies_dict.items()
    ]

    with open(PROFILE_DIR / "cookies.json", "w") as f:
        json.dump(cookies_list, f, indent=2)

    with open(PROFILE_DIR / "metadata.json", "w") as f:
        json.dump({
            "csrf_token": csrf_token,
            "session_id": "",
            "email": GOOGLE_EMAIL,
            "last_validated": datetime.now().isoformat()
        }, f, indent=2)

    with open(PROFILE_DIR / "auth.json", "w") as f:
        json.dump(auth_data, f, indent=2)

    log(f"  ✓ 쿠키 저장 완료 ({len(cookies_dict)}개)")
    return auth_data


def pdf_to_pptx(pdf_path, output_name=None):
    """PDF를 PPTX로 변환"""
    import fitz
    from pptx import Presentation
    from pptx.util import Inches

    pdf_path = Path(pdf_path)
    if output_name:
        output_path = DOWNLOAD_DIR / f"{safe_filename(output_name)}.pptx"
    else:
        output_path = pdf_path.with_suffix('.pptx')

    doc = fitz.open(pdf_path)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for page_num in range(len(doc)):
        page = doc[page_num]
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img_data = pix.tobytes("png")

        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            io.BytesIO(img_data),
            Inches(0), Inches(0),
            width=prs.slide_width,
            height=prs.slide_height
        )

    doc.close()
    prs.save(output_path)
    return output_path, len(prs.slides)


async def download_slides(page, keyword: str, display_name: str) -> dict:
    """노트북에서 슬라이드 다운로드"""
    result = {'name': display_name, 'pdf': None, 'pptx': None, 'slides': 0}

    try:
        # 홈으로 이동
        await page.goto("https://notebooklm.google.com/", wait_until='networkidle', timeout=30000)
        await asyncio.sleep(3)

        # 노트북 찾기
        log(f"  노트북 찾기: '{keyword}'...")
        notebook_link = None

        all_elements = await page.query_selector_all('a, div[role="button"], [class*="notebook"]')
        for el in all_elements:
            try:
                text = await el.inner_text()
                if keyword in text:
                    notebook_link = el
                    log(f"  → 발견!")
                    break
            except:
                continue

        if not notebook_link:
            log(f"  ❌ 노트북 찾기 실패")
            return result

        await notebook_link.click()
        await asyncio.sleep(5)

        # 다운로드 시도
        pdf_path = None

        # 메뉴 버튼 찾기
        menu_buttons = await page.query_selector_all(
            '[aria-haspopup="menu"], button[aria-label*="more"], button[aria-label*="More"]'
        )

        for menu_btn in menu_buttons[-10:]:
            try:
                await menu_btn.click(force=True)
                await asyncio.sleep(1.5)

                download_item = await page.query_selector(
                    '[role="menuitem"]:has-text("Download PDF"), '
                    '[role="menuitem"]:has-text("Download"), '
                    '[role="menuitem"]:has-text("다운로드")'
                )

                if download_item:
                    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                    filename = f"{safe_filename(display_name)}_{timestamp}.pdf"

                    async with page.expect_download(timeout=60000) as download_info:
                        await download_item.click()

                    download = await download_info.value
                    pdf_path = DOWNLOAD_DIR / filename
                    await download.save_as(str(pdf_path))
                    log(f"  ✓ PDF 다운로드: {filename}")
                    break
            except:
                await page.keyboard.press('Escape')

        # 좌표 기반 시도
        if not pdf_path:
            coordinates = [(1846, 365), (1800, 400), (1850, 350)]
            for x, y in coordinates:
                try:
                    await page.mouse.click(x, y)
                    await asyncio.sleep(1.5)

                    menu_items = await page.query_selector_all('[role="menuitem"]')
                    for item in menu_items:
                        text = await item.inner_text()
                        if 'Download' in text or '다운로드' in text:
                            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                            filename = f"{safe_filename(display_name)}_{timestamp}.pdf"

                            async with page.expect_download(timeout=60000) as download_info:
                                await item.click()

                            download = await download_info.value
                            pdf_path = DOWNLOAD_DIR / filename
                            await download.save_as(str(pdf_path))
                            break
                    if pdf_path:
                        break
                except:
                    await page.keyboard.press('Escape')

        if pdf_path and pdf_path.exists():
            result['pdf'] = str(pdf_path)
            pptx_path, slide_count = pdf_to_pptx(pdf_path, display_name)
            result['pptx'] = str(pptx_path)
            result['slides'] = slide_count
            log(f"  ✓ PPTX: {pptx_path.name} ({slide_count}장)")

    except Exception as e:
        log(f"  ❌ 오류: {e}")

    return result


async def main():
    from playwright.async_api import async_playwright

    log("=" * 60)
    log("NotebookLM 자동 인증 + 슬라이드 다운로드")
    log("=" * 60)
    log(f"  이메일: {GOOGLE_EMAIL}")
    log(f"  앱 비밀번호: {'*' * 16}")
    log("=" * 60)

    DOWNLOAD_DIR.mkdir(parents=True, exist_ok=True)

    # 처리할 노트북
    notebooks = [
        ("족저근막염", "족저근막염"),
        ("족모지외반증", "족모지외반증"),
        ("발톱무좀", "발톱무좀"),
    ]

    results = []

    async with async_playwright() as p:
        log("\n[단계 1] 브라우저 시작...")

        browser = await p.chromium.launch(
            headless=False,
            args=['--disable-blink-features=AutomationControlled'],
        )
        context = await browser.new_context(
            viewport={'width': 1920, 'height': 1080},
            accept_downloads=True,
        )

        page = await context.new_page()

        # NotebookLM 접속
        log("\n[단계 2] NotebookLM 접속...")
        await page.goto("https://notebooklm.google.com/", wait_until='networkidle', timeout=60000)
        await asyncio.sleep(3)

        # 로그인 필요 확인
        if "accounts.google.com" in page.url:
            log("\n[단계 3] 앱 비밀번호로 로그인...")

            login_success = await google_login_with_app_password(page, GOOGLE_EMAIL, APP_PASSWORD)

            if not login_success:
                log("  ❌ 로그인 실패")
                await browser.close()
                return []
        else:
            log("  ✓ 이미 로그인됨")

        # 쿠키 저장
        log("\n[단계 4] 인증 정보 저장...")
        await extract_and_save_cookies(context)

        # 슬라이드 다운로드
        log("\n[단계 5] 슬라이드 다운로드...")
        for keyword, display_name in notebooks:
            log(f"\n--- {display_name} ---")
            result = await download_slides(page, keyword, display_name)
            results.append(result)
            await asyncio.sleep(2)

        await browser.close()

    # 결과 출력
    log("\n" + "=" * 60)
    log("완료 결과:")
    log("=" * 60)

    success_count = 0
    for r in results:
        status = "✓" if r.get('pptx') else "❌"
        log(f"  {status} {r['name']}")
        if r.get('pptx'):
            log(f"      PPTX: {Path(r['pptx']).name} ({r['slides']}장)")
            success_count += 1

    log(f"\n성공: {success_count}/{len(notebooks)}")
    log(f"저장 위치: {DOWNLOAD_DIR}")
    log("=" * 60)

    return results


if __name__ == "__main__":
    asyncio.run(main())
