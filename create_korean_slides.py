#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
NotebookLM 한글 슬라이드 생성 - 스튜디오 패널의 '슬라이드 자료' 버튼 사용
"""
import asyncio
import sys
import os
import io
import re
from pathlib import Path
from datetime import datetime

os.environ['PYTHONUNBUFFERED'] = '1'
if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8', line_buffering=True)

AUTH_DIR = Path.home() / ".notebooklm-mcp-cli"
PERSISTENT_PROFILE = AUTH_DIR / "persistent_browser"
DOWNLOAD_DIR = Path("G:/내 드라이브/notebooklm")

NOTEBOOKS = [
    ("족저근막염", "족저근막염"),
    ("족모지외반증", "족모지외반증"),
    ("발톱무좀", "발톱무좀"),
]


def log(msg):
    print(msg, flush=True)


def safe_filename(name):
    return re.sub(r'[<>:"/\\|?*]', '', name).strip()


def pdf_to_pptx(pdf_path, output_name=None):
    import fitz
    from pptx import Presentation
    from pptx.util import Inches

    pdf_path = Path(pdf_path)
    output_path = DOWNLOAD_DIR / f"{safe_filename(output_name or pdf_path.stem)}.pptx"

    doc = fitz.open(pdf_path)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for page_num in range(len(doc)):
        page = doc[page_num]
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img_data = pix.tobytes("png")
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(io.BytesIO(img_data), Inches(0), Inches(0),
                                  width=prs.slide_width, height=prs.slide_height)

    doc.close()
    prs.save(output_path)
    return output_path, len(prs.slides)


async def click_slides_button(page):
    """스튜디오 패널에서 '슬라이드 자료' 버튼 클릭"""
    log("  슬라이드 자료 버튼 찾기...")

    # 방법 1: 텍스트로 찾기
    try:
        slides_btn = await page.query_selector('text="슬라이드 자료"')
        if slides_btn:
            log("    '슬라이드 자료' 버튼 발견!")
            await slides_btn.click()
            await asyncio.sleep(3)
            # 클릭 후 스크린샷
            await page.screenshot(path=str(DOWNLOAD_DIR / "after_slides_click.png"))
            log("    스크린샷 저장: after_slides_click.png")
            return True
    except:
        pass

    # 방법 2: 영문 텍스트로 찾기
    try:
        slides_btn = await page.query_selector('text="Slides"')
        if not slides_btn:
            slides_btn = await page.query_selector('text="Presentation"')
        if slides_btn:
            log("    'Slides' 버튼 발견!")
            await slides_btn.click()
            await asyncio.sleep(3)
            return True
    except:
        pass

    # 방법 3: 스튜디오 패널 내 버튼들 검색
    try:
        buttons = await page.query_selector_all('button, [role="button"], div[class*="card"]')
        for btn in buttons:
            try:
                text = await btn.inner_text()
                if '슬라이드' in text or 'Slide' in text:
                    log(f"    슬라이드 버튼 발견: {text[:20]}")
                    await btn.click()
                    await asyncio.sleep(3)
                    return True
            except:
                continue
    except:
        pass

    # 방법 4: 좌표로 클릭 (스튜디오 패널의 슬라이드 자료 위치)
    # 스크린샷 기준 대략 (1240, 215) 부근
    try:
        log("    좌표로 슬라이드 버튼 클릭...")
        await page.mouse.click(1240, 215)
        await asyncio.sleep(3)
        return True
    except:
        pass

    return False


async def configure_and_generate(page):
    """슬라이드 설정 및 생성"""
    log("  슬라이드 설정...")

    # 다이얼로그가 열렸는지 확인
    await asyncio.sleep(2)

    # 언어 설정: 한국어
    try:
        # 드롭다운 찾기
        dropdowns = await page.query_selector_all('select, [role="combobox"], [role="listbox"], button[aria-haspopup="listbox"]')
        for dropdown in dropdowns:
            try:
                text = await dropdown.inner_text()
                aria = await dropdown.get_attribute('aria-label') or ''
                if 'language' in text.lower() or 'language' in aria.lower() or \
                   'english' in text.lower() or '영어' in text or '언어' in text:
                    log("    언어 드롭다운 발견")
                    await dropdown.click()
                    await asyncio.sleep(1)

                    # 한국어 선택
                    options = await page.query_selector_all('[role="option"], li, option, [role="menuitem"]')
                    for opt in options:
                        try:
                            opt_text = await opt.inner_text()
                            if '한국어' in opt_text or 'Korean' in opt_text:
                                log("    한국어 선택!")
                                await opt.click()
                                await asyncio.sleep(1)
                                break
                        except:
                            continue
                    break
            except:
                continue
    except Exception as e:
        log(f"    언어 설정 실패: {e}")

    # Generate/생성 버튼 클릭
    await asyncio.sleep(1)
    try:
        gen_buttons = await page.query_selector_all('button')
        for btn in gen_buttons:
            try:
                text = await btn.inner_text()
                if any(x in text for x in ['Generate', '생성', 'Create', '만들기']):
                    log(f"    생성 버튼 클릭: {text[:15]}")
                    await btn.click()
                    return True
            except:
                continue
    except:
        pass

    return False


async def wait_for_slides(page, timeout=180):
    """슬라이드 생성 완료 대기"""
    log(f"  생성 대기 중... (최대 {timeout}초)")

    start = asyncio.get_event_loop().time()

    while asyncio.get_event_loop().time() - start < timeout:
        await asyncio.sleep(10)
        elapsed = int(asyncio.get_event_loop().time() - start)

        if elapsed % 30 == 0:
            log(f"    {elapsed}초...")

        # 다운로드 가능 여부 확인
        try:
            # 스튜디오 패널에서 생성된 슬라이드 카드의 메뉴 확인
            menu_buttons = await page.query_selector_all('[aria-haspopup="menu"]')
            for btn in menu_buttons[-10:]:
                try:
                    await btn.click(force=True)
                    await asyncio.sleep(0.5)

                    dl_item = await page.query_selector('[role="menuitem"]:has-text("Download")')
                    await page.keyboard.press('Escape')

                    if dl_item:
                        log(f"    완료! ({elapsed}초)")
                        return True
                except:
                    await page.keyboard.press('Escape')
        except:
            pass

    log(f"    타임아웃 ({timeout}초)")
    return False


async def download_pdf(page, display_name):
    """PDF 다운로드"""
    log("  PDF 다운로드...")

    # 방법 1: 메뉴 버튼
    menu_buttons = await page.query_selector_all('[aria-haspopup="menu"]')
    log(f"    메뉴 버튼 {len(menu_buttons)}개")

    for btn in menu_buttons[-15:]:
        try:
            await btn.click(force=True)
            await asyncio.sleep(1.5)

            dl_item = await page.query_selector(
                '[role="menuitem"]:has-text("Download PDF"), '
                '[role="menuitem"]:has-text("Download"), '
                '[role="menuitem"]:has-text("다운로드")'
            )

            if dl_item:
                timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                filename = f"{safe_filename(display_name)}_{timestamp}.pdf"

                async with page.expect_download(timeout=60000) as download_info:
                    await dl_item.click()

                download = await download_info.value
                pdf_path = DOWNLOAD_DIR / filename
                await download.save_as(str(pdf_path))
                log(f"    PDF: {filename}")
                return pdf_path
        except:
            await page.keyboard.press('Escape')

    # 방법 2: 좌표 기반 (스튜디오 패널의 슬라이드 카드 메뉴)
    log("    좌표 시도...")
    for x, y in [(1400, 220), (1300, 220), (1350, 220), (1846, 365)]:
        try:
            await page.mouse.click(x, y)
            await asyncio.sleep(1.5)

            items = await page.query_selector_all('[role="menuitem"]')
            for item in items:
                try:
                    text = await item.inner_text()
                    if 'Download' in text or '다운로드' in text:
                        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                        filename = f"{safe_filename(display_name)}_{timestamp}.pdf"

                        async with page.expect_download(timeout=60000) as download_info:
                            await item.click()

                        download = await download_info.value
                        pdf_path = DOWNLOAD_DIR / filename
                        await download.save_as(str(pdf_path))
                        log(f"    PDF: {filename}")
                        return pdf_path
                except:
                    continue
        except:
            await page.keyboard.press('Escape')

    return None


async def process_notebook(page, keyword, display_name, idx, total):
    """단일 노트북 처리"""
    result = {'name': display_name, 'pdf': None, 'pptx': None, 'slides': 0}

    log(f"\n{'='*55}")
    log(f"[{idx}/{total}] {display_name}")
    log('='*55)

    try:
        # 홈으로
        await page.goto("https://notebooklm.google.com/", wait_until='networkidle', timeout=30000)
        await asyncio.sleep(3)

        # 노트북 찾기
        log(f"  노트북: '{keyword}'")
        notebook_link = None

        elements = await page.query_selector_all('a, div[role="button"], tr, [class*="notebook"]')
        for el in elements:
            try:
                text = await el.inner_text()
                if keyword in text:
                    notebook_link = el
                    break
            except:
                continue

        if not notebook_link:
            log("    찾기 실패!")
            return result

        log("    발견!")
        await notebook_link.click()
        await asyncio.sleep(5)

        # 슬라이드 자료 버튼 클릭
        clicked = await click_slides_button(page)

        if clicked:
            # 설정 및 생성
            await configure_and_generate(page)

            # 생성 대기
            await wait_for_slides(page, timeout=180)

        # 다운로드
        pdf_path = await download_pdf(page, display_name)

        if pdf_path and pdf_path.exists():
            result['pdf'] = str(pdf_path)
            log("  PPTX 변환...")
            pptx_path, slide_count = pdf_to_pptx(pdf_path, display_name)
            result['pptx'] = str(pptx_path)
            result['slides'] = slide_count
            log(f"    완료: {pptx_path.name} ({slide_count}장)")
        else:
            log("    다운로드 실패!")

    except Exception as e:
        log(f"  오류: {e}")

    return result


async def main():
    from playwright.async_api import async_playwright

    log("=" * 55)
    log("NotebookLM 15장 한글 슬라이드 생성")
    log("=" * 55)
    for _, name in NOTEBOOKS:
        log(f"  - {name}")
    log("=" * 55)

    DOWNLOAD_DIR.mkdir(parents=True, exist_ok=True)
    results = []

    async with async_playwright() as p:
        log("\n브라우저 시작...")

        context = await p.chromium.launch_persistent_context(
            user_data_dir=str(PERSISTENT_PROFILE),
            headless=False,
            downloads_path=str(DOWNLOAD_DIR),
            accept_downloads=True,
            args=['--disable-blink-features=AutomationControlled'],
            viewport={'width': 1920, 'height': 1080},
        )

        page = context.pages[0] if context.pages else await context.new_page()

        await page.goto("https://notebooklm.google.com/", wait_until='networkidle', timeout=30000)
        await asyncio.sleep(3)

        if "accounts.google.com" in page.url:
            log("로그인 대기 중... (브라우저에서 로그인)")
            for _ in range(60):
                await asyncio.sleep(5)
                if "notebooklm.google.com" in page.url and "accounts.google" not in page.url:
                    break
            else:
                await context.close()
                return []

        log("준비 완료!\n")

        for idx, (keyword, display_name) in enumerate(NOTEBOOKS, 1):
            result = await process_notebook(page, keyword, display_name, idx, len(NOTEBOOKS))
            results.append(result)
            await asyncio.sleep(3)

        await context.close()

    # 결과
    log("\n" + "=" * 55)
    log("결과:")
    log("=" * 55)

    success = sum(1 for r in results if r.get('pptx'))
    for r in results:
        status = "O" if r.get('pptx') else "X"
        log(f"  [{status}] {r['name']}")
        if r.get('pptx'):
            log(f"       {Path(r['pptx']).name} ({r['slides']}장)")

    log(f"\n성공: {success}/{len(NOTEBOOKS)}")
    log(f"저장: {DOWNLOAD_DIR}")

    return results


if __name__ == "__main__":
    asyncio.run(main())
