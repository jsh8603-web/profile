#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
NotebookLM 한글 슬라이드 생성 + 다운로드 + PPTX 변환
- 15장 한글 슬라이드 새로 생성
- PDF 다운로드
- PPTX 변환
"""
import asyncio
import sys
import os
import io
import re
from pathlib import Path
from datetime import datetime

os.environ['PYTHONUNBUFFERED'] = '1'
if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8', line_buffering=True)

# 설정
AUTH_DIR = Path.home() / ".notebooklm-mcp-cli"
PERSISTENT_PROFILE = AUTH_DIR / "persistent_browser"
DOWNLOAD_DIR = Path("G:/내 드라이브/notebooklm")

# 처리할 노트북
NOTEBOOKS = [
    ("족저근막염", "족저근막염"),
    ("족모지외반증", "족모지외반증"),
    ("발톱무좀", "발톱무좀"),
]


def log(msg):
    print(msg, flush=True)


def safe_filename(name):
    return re.sub(r'[<>:"/\\|?*]', '', name).strip()


def pdf_to_pptx(pdf_path, output_name=None):
    """PDF를 PPTX로 변환"""
    import fitz
    from pptx import Presentation
    from pptx.util import Inches

    pdf_path = Path(pdf_path)
    if output_name:
        output_path = DOWNLOAD_DIR / f"{safe_filename(output_name)}.pptx"
    else:
        output_path = pdf_path.with_suffix('.pptx')

    doc = fitz.open(pdf_path)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for page_num in range(len(doc)):
        page = doc[page_num]
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img_data = pix.tobytes("png")

        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            io.BytesIO(img_data),
            Inches(0), Inches(0),
            width=prs.slide_width,
            height=prs.slide_height
        )

    doc.close()
    prs.save(output_path)
    return output_path, len(prs.slides)


async def create_korean_slides(page, notebook_name):
    """한글 슬라이드 15장 생성"""
    log(f"  슬라이드 생성 시작...")

    # Studio 패널에서 Create/Add 버튼 찾기
    create_selectors = [
        'button:has-text("Create")',
        'button:has-text("만들기")',
        'button:has-text("Add")',
        'button:has-text("추가")',
        'button[aria-label*="Create"]',
        'button[aria-label*="Add"]',
        '[data-tooltip*="Create"]',
    ]

    create_btn = None
    for selector in create_selectors:
        create_btn = await page.query_selector(selector)
        if create_btn:
            break

    if create_btn:
        await create_btn.click()
        await asyncio.sleep(2)
        log(f"  Create 버튼 클릭")

    # Slides 옵션 선택
    slides_selectors = [
        '[role="menuitem"]:has-text("Slides")',
        '[role="menuitem"]:has-text("슬라이드")',
        'button:has-text("Slides")',
        'button:has-text("슬라이드")',
        'text="Slides"',
    ]

    slides_option = None
    for selector in slides_selectors:
        slides_option = await page.query_selector(selector)
        if slides_option:
            break

    if slides_option:
        await slides_option.click()
        await asyncio.sleep(3)
        log(f"  Slides 옵션 선택")

    # 설정 다이얼로그에서 언어 설정
    # 언어 드롭다운 찾기
    lang_selectors = [
        '[aria-label*="language"]',
        '[aria-label*="Language"]',
        'select',
        '[role="combobox"]',
        '[role="listbox"]',
        'button:has-text("English")',
        'button:has-text("영어")',
    ]

    for selector in lang_selectors:
        lang_elements = await page.query_selector_all(selector)
        for lang_el in lang_elements:
            try:
                await lang_el.click()
                await asyncio.sleep(1)

                # 한국어 옵션 찾기
                korean_options = [
                    '[role="option"]:has-text("한국어")',
                    '[role="option"]:has-text("Korean")',
                    'option:has-text("한국어")',
                    'li:has-text("한국어")',
                    'text="한국어"',
                    'text="Korean"',
                ]

                for ko_selector in korean_options:
                    korean_option = await page.query_selector(ko_selector)
                    if korean_option:
                        await korean_option.click()
                        log(f"  ✓ 한국어 선택")
                        await asyncio.sleep(1)
                        break
                break
            except:
                continue

    # Generate/생성 버튼 클릭
    generate_selectors = [
        'button:has-text("Generate")',
        'button:has-text("생성")',
        'button:has-text("Create")',
        'button[type="submit"]',
    ]

    for selector in generate_selectors:
        generate_btn = await page.query_selector(selector)
        if generate_btn:
            try:
                await generate_btn.click()
                log(f"  ✓ 생성 시작!")
                break
            except:
                continue

    return True


async def wait_for_slides(page, timeout=180):
    """슬라이드 생성 완료 대기"""
    log(f"  슬라이드 생성 대기 중... (최대 {timeout}초)")

    start = asyncio.get_event_loop().time()
    while asyncio.get_event_loop().time() - start < timeout:
        await asyncio.sleep(5)
        elapsed = int(asyncio.get_event_loop().time() - start)

        # 완료 확인: 다운로드 가능한 슬라이드 카드 존재
        try:
            # 메뉴 버튼이 있는 슬라이드 카드 확인
            menu_buttons = await page.query_selector_all('[aria-haspopup="menu"]')
            for menu_btn in menu_buttons[-5:]:
                try:
                    await menu_btn.click(force=True)
                    await asyncio.sleep(0.5)

                    download_item = await page.query_selector(
                        '[role="menuitem"]:has-text("Download")'
                    )
                    await page.keyboard.press('Escape')

                    if download_item:
                        log(f"  ✓ 슬라이드 생성 완료! ({elapsed}초)")
                        return True
                except:
                    await page.keyboard.press('Escape')

            # 로딩 상태 확인
            loading = await page.query_selector('[class*="loading"], [class*="progress"], [class*="spinner"]')
            if loading and elapsed % 15 == 0:
                log(f"  생성 중... {elapsed}초")

        except:
            pass

        if elapsed % 30 == 0:
            log(f"  대기 중... {elapsed}초")

    log(f"  ⚠️ 타임아웃 ({timeout}초) - 그래도 다운로드 시도")
    return False


async def download_pdf(page, display_name):
    """PDF 다운로드"""
    pdf_path = None

    # 방법 1: 메뉴 버튼
    menu_buttons = await page.query_selector_all(
        '[aria-haspopup="menu"], button[aria-label*="more"], button[aria-label*="More"]'
    )
    log(f"  메뉴 버튼 {len(menu_buttons)}개")

    for menu_btn in menu_buttons[-15:]:
        try:
            await menu_btn.click(force=True)
            await asyncio.sleep(1.5)

            download_item = await page.query_selector(
                '[role="menuitem"]:has-text("Download PDF"), '
                '[role="menuitem"]:has-text("Download"), '
                '[role="menuitem"]:has-text("다운로드")'
            )

            if download_item:
                timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                filename = f"{safe_filename(display_name)}_{timestamp}.pdf"

                async with page.expect_download(timeout=60000) as download_info:
                    await download_item.click()

                download = await download_info.value
                pdf_path = DOWNLOAD_DIR / filename
                await download.save_as(str(pdf_path))
                log(f"  ✓ PDF: {filename}")
                return pdf_path
        except:
            await page.keyboard.press('Escape')

    # 방법 2: 좌표 기반
    if not pdf_path:
        log(f"  좌표 기반 시도...")
        coordinates = [(1846, 365), (1800, 400), (1850, 350), (1750, 380)]

        for x, y in coordinates:
            try:
                await page.mouse.click(x, y)
                await asyncio.sleep(1.5)

                menu_items = await page.query_selector_all('[role="menuitem"]')
                for item in menu_items:
                    text = await item.inner_text()
                    if 'Download' in text or '다운로드' in text:
                        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                        filename = f"{safe_filename(display_name)}_{timestamp}.pdf"

                        async with page.expect_download(timeout=60000) as download_info:
                            await item.click()

                        download = await download_info.value
                        pdf_path = DOWNLOAD_DIR / filename
                        await download.save_as(str(pdf_path))
                        log(f"  ✓ PDF (좌표): {filename}")
                        return pdf_path
            except:
                await page.keyboard.press('Escape')

    return None


async def process_notebook(page, keyword, display_name):
    """단일 노트북 처리: 슬라이드 생성 → 다운로드 → 변환"""
    result = {'name': display_name, 'pdf': None, 'pptx': None, 'slides': 0}

    try:
        # 1. 홈으로 이동
        log(f"  홈으로 이동...")
        await page.goto("https://notebooklm.google.com/", wait_until='networkidle', timeout=30000)
        await asyncio.sleep(3)

        # 2. 노트북 찾기
        log(f"  노트북 찾기: '{keyword}'...")
        notebook_link = None

        all_elements = await page.query_selector_all('a, div[role="button"], [class*="notebook"], tr')
        for el in all_elements:
            try:
                text = await el.inner_text()
                if keyword in text:
                    notebook_link = el
                    log(f"  → 발견!")
                    break
            except:
                continue

        if not notebook_link:
            log(f"  ❌ 노트북 찾기 실패")
            return result

        # 3. 노트북 열기
        await notebook_link.click()
        await asyncio.sleep(5)
        log(f"  노트북 열림")

        # 4. 한글 슬라이드 생성
        await create_korean_slides(page, display_name)

        # 5. 생성 완료 대기
        await wait_for_slides(page, timeout=180)

        # 6. PDF 다운로드
        log(f"  PDF 다운로드...")
        pdf_path = await download_pdf(page, display_name)

        if pdf_path and pdf_path.exists():
            result['pdf'] = str(pdf_path)

            # 7. PPTX 변환
            log(f"  PPTX 변환...")
            pptx_path, slide_count = pdf_to_pptx(pdf_path, display_name)
            result['pptx'] = str(pptx_path)
            result['slides'] = slide_count
            log(f"  ✓ PPTX: {pptx_path.name} ({slide_count}장)")
        else:
            log(f"  ❌ 다운로드 실패")

    except Exception as e:
        log(f"  ❌ 오류: {e}")

    return result


async def main():
    from playwright.async_api import async_playwright

    log("=" * 60)
    log("NotebookLM 15장 한글 슬라이드 생성")
    log("=" * 60)
    for _, name in NOTEBOOKS:
        log(f"  - {name}")
    log("=" * 60)

    DOWNLOAD_DIR.mkdir(parents=True, exist_ok=True)
    PERSISTENT_PROFILE.mkdir(parents=True, exist_ok=True)

    results = []

    async with async_playwright() as p:
        log("\n[1] 브라우저 시작...")

        context = await p.chromium.launch_persistent_context(
            user_data_dir=str(PERSISTENT_PROFILE),
            headless=False,
            downloads_path=str(DOWNLOAD_DIR),
            accept_downloads=True,
            args=['--disable-blink-features=AutomationControlled'],
            viewport={'width': 1920, 'height': 1080},
        )

        page = context.pages[0] if context.pages else await context.new_page()

        # NotebookLM 접속
        log("\n[2] NotebookLM 접속...")
        await page.goto("https://notebooklm.google.com/", wait_until='networkidle', timeout=30000)
        await asyncio.sleep(3)

        # 로그인 확인
        if "accounts.google.com" in page.url:
            log("  로그인 대기 중...")
            for i in range(60):
                await asyncio.sleep(5)
                if "notebooklm.google.com" in page.url and "accounts.google" not in page.url:
                    log("  ✓ 로그인 완료!")
                    break
            else:
                log("  ❌ 로그인 타임아웃")
                await context.close()
                return []
        else:
            log("  ✓ 이미 로그인됨")

        # 각 노트북 처리
        log("\n[3] 슬라이드 생성...")
        for idx, (keyword, display_name) in enumerate(NOTEBOOKS, 1):
            log(f"\n{'='*50}")
            log(f"[{idx}/{len(NOTEBOOKS)}] {display_name}")
            log('='*50)
            result = await process_notebook(page, keyword, display_name)
            results.append(result)
            await asyncio.sleep(3)

        await context.close()

    # 결과 출력
    log("\n" + "=" * 60)
    log("완료 결과:")
    log("=" * 60)

    success_count = 0
    for r in results:
        status = "✓" if r.get('pptx') else "❌"
        log(f"  {status} {r['name']}")
        if r.get('pptx'):
            log(f"      PPTX: {Path(r['pptx']).name} ({r['slides']}장)")
            success_count += 1

    log(f"\n성공: {success_count}/{len(NOTEBOOKS)}")
    log(f"저장 위치: {DOWNLOAD_DIR}")
    log("=" * 60)

    return results


if __name__ == "__main__":
    asyncio.run(main())
