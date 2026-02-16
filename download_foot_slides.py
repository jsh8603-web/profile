#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
기존 노트북에서 슬라이드 다운로드 (슬라이드가 이미 생성되어 있는 경우)
"""
import asyncio
import sys
import os
import io
import re
from pathlib import Path
from datetime import datetime

os.environ['PYTHONUNBUFFERED'] = '1'
if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8', line_buffering=True)

DOWNLOAD_DIR = Path("G:/내 드라이브/notebooklm")
AUTH_DIR = Path.home() / ".notebooklm-mcp-cli"

# 찾을 노트북 키워드 (전체 이름 대신 키워드로 찾음)
NOTEBOOKS = [
    ("족저근막염", "족저근막염"),
    ("족모지외반증", "족모지외반증"),
    ("발톱무좀", "발톱무좀"),
]


def log(msg):
    print(msg, flush=True)


def safe_filename(name):
    """파일명에 사용할 수 없는 문자 제거"""
    return re.sub(r'[<>:"/\\|?*]', '', name).strip()


def pdf_to_pptx(pdf_path, output_name=None):
    """PDF를 PPTX로 변환"""
    import fitz
    from pptx import Presentation
    from pptx.util import Inches

    pdf_path = Path(pdf_path)
    if output_name:
        output_path = DOWNLOAD_DIR / f"{safe_filename(output_name)}.pptx"
    else:
        output_path = pdf_path.with_suffix('.pptx')

    doc = fitz.open(pdf_path)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for page_num in range(len(doc)):
        page = doc[page_num]
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img_data = pix.tobytes("png")

        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            io.BytesIO(img_data),
            Inches(0), Inches(0),
            width=prs.slide_width,
            height=prs.slide_height
        )

    doc.close()
    prs.save(output_path)
    return output_path, len(prs.slides)


async def download_from_notebook(page, keyword, display_name):
    """노트북에서 슬라이드 다운로드"""
    log(f"\n{'='*50}")
    log(f"[{display_name}] 처리 중...")
    log('='*50)

    result = {'name': display_name, 'pdf': None, 'pptx': None, 'slides': 0}

    try:
        # 홈으로 이동
        await page.goto("https://notebooklm.google.com/", wait_until='networkidle', timeout=30000)
        await asyncio.sleep(3)

        # 노트북 찾기 (키워드로)
        log(f"  노트북 찾기: '{keyword}'...")
        notebook_link = None

        all_elements = await page.query_selector_all('a, div[role="button"], [class*="notebook"]')
        for el in all_elements:
            try:
                text = await el.inner_text()
                if keyword in text:
                    notebook_link = el
                    log(f"  → 발견: {text[:40]}...")
                    break
            except:
                continue

        if not notebook_link:
            log(f"  ❌ 노트북 찾기 실패")
            return result

        # 노트북 클릭
        await notebook_link.click()
        await asyncio.sleep(5)

        # URL에서 ID 추출
        current_url = page.url
        match = re.search(r'/notebook/([a-zA-Z0-9_-]+)', current_url)
        notebook_id = match.group(1) if match else None
        log(f"  노트북 ID: {notebook_id[:12]}...")

        # 페이지 로드 대기
        log(f"  페이지 로드 대기...")
        await asyncio.sleep(5)

        # 스크린샷 저장 (디버깅용)
        screenshot_name = safe_filename(f"page_{display_name}_{datetime.now().strftime('%H%M%S')}")
        await page.screenshot(path=str(DOWNLOAD_DIR / f"{screenshot_name}.png"))
        log(f"  스크린샷: {screenshot_name}.png")

        # 다운로드 시도 - 여러 방법
        pdf_path = None

        # 방법 1: 메뉴 버튼 찾기
        log(f"  다운로드 버튼 찾기...")
        menu_buttons = await page.query_selector_all(
            '[aria-haspopup="menu"], button[aria-label*="more"], button[aria-label*="More"], '
            '[data-tooltip*="more"], [aria-label*="옵션"]'
        )
        log(f"  메뉴 버튼 {len(menu_buttons)}개 발견")

        for idx, menu_btn in enumerate(menu_buttons[-10:]):
            try:
                await menu_btn.click(force=True)
                await asyncio.sleep(1.5)

                # 다운로드 메뉴 항목 찾기
                download_item = await page.query_selector(
                    '[role="menuitem"]:has-text("Download PDF"), '
                    '[role="menuitem"]:has-text("Download"), '
                    '[role="menuitem"]:has-text("다운로드")'
                )

                if download_item:
                    log(f"  다운로드 항목 발견! (메뉴 {idx})")
                    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                    filename = f"{safe_filename(display_name)}_{timestamp}.pdf"

                    async with page.expect_download(timeout=60000) as download_info:
                        await download_item.click()

                    download = await download_info.value
                    pdf_path = DOWNLOAD_DIR / filename
                    await download.save_as(str(pdf_path))
                    log(f"  ✓ PDF 다운로드 성공: {filename}")
                    break

            except Exception as e:
                await page.keyboard.press('Escape')
                await asyncio.sleep(0.3)

        # 방법 2: 좌표 기반 클릭 (스튜디오 패널 슬라이드 카드)
        if not pdf_path:
            log(f"  좌표 기반 다운로드 시도...")
            coordinates = [(1846, 365), (1800, 400), (1850, 350), (1750, 400)]

            for x, y in coordinates:
                try:
                    await page.mouse.click(x, y)
                    await asyncio.sleep(1.5)

                    menu_items = await page.query_selector_all('[role="menuitem"]')
                    for item in menu_items:
                        text = await item.inner_text()
                        if 'Download' in text or '다운로드' in text:
                            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                            filename = f"{safe_filename(display_name)}_{timestamp}.pdf"

                            async with page.expect_download(timeout=60000) as download_info:
                                await item.click()

                            download = await download_info.value
                            pdf_path = DOWNLOAD_DIR / filename
                            await download.save_as(str(pdf_path))
                            log(f"  ✓ PDF 다운로드 성공 (좌표): {filename}")
                            break

                    if pdf_path:
                        break

                except:
                    await page.keyboard.press('Escape')

        if pdf_path and pdf_path.exists():
            result['pdf'] = str(pdf_path)

            # PPTX 변환
            log(f"  PPTX 변환 중...")
            pptx_path, slide_count = pdf_to_pptx(pdf_path, display_name)
            result['pptx'] = str(pptx_path)
            result['slides'] = slide_count
            log(f"  ✓ PPTX: {pptx_path.name} ({slide_count}장)")
        else:
            log(f"  ❌ 다운로드 실패")

    except Exception as e:
        log(f"  ❌ 오류: {e}")

    return result


async def main():
    from playwright.async_api import async_playwright

    log("=" * 60)
    log("족부 질환 슬라이드 다운로드")
    log("  - 족저근막염")
    log("  - 족모지외반증")
    log("  - 발톱무좀")
    log("=" * 60)

    DOWNLOAD_DIR.mkdir(parents=True, exist_ok=True)
    results = []

    async with async_playwright() as p:
        log("\n[메인] 브라우저 시작...")

        browser = await p.chromium.launch(
            channel='chrome',
            headless=False,
            args=['--disable-blink-features=AutomationControlled'],
        )
        context = await browser.new_context(
            viewport={'width': 1920, 'height': 1080},
            accept_downloads=True,
        )
        context._downloads_path = str(DOWNLOAD_DIR)

        page = context.pages[0] if context.pages else await context.new_page()

        log("[메인] NotebookLM 접속...")
        await page.goto("https://notebooklm.google.com/", wait_until='networkidle', timeout=60000)
        await asyncio.sleep(3)

        # 로그인 확인
        if "accounts.google.com" in page.url:
            log("  ⚠️ 로그인 필요 - 브라우저에서 로그인해주세요...")

            for i in range(60):
                await asyncio.sleep(5)
                log(f"  로그인 대기 중... {(i+1)*5}초")
                if "notebooklm.google.com" in page.url and "accounts.google" not in page.url:
                    log("  ✓ 로그인 완료!")
                    break
            else:
                log("  ❌ 로그인 타임아웃")
                await context.close()
                await browser.close()
                return []

        log("  ✓ NotebookLM 준비 완료")

        # 각 노트북 처리
        for keyword, display_name in NOTEBOOKS:
            result = await download_from_notebook(page, keyword, display_name)
            results.append(result)
            await asyncio.sleep(2)

        await context.close()
        await browser.close()

    # 결과 출력
    log("\n" + "=" * 60)
    log("완료 결과:")
    log("=" * 60)

    success_count = 0
    for r in results:
        status = "✓" if r.get('pptx') else "❌"
        log(f"  {status} {r['name']}")
        if r.get('pptx'):
            log(f"      PDF:  {Path(r['pdf']).name}")
            log(f"      PPTX: {Path(r['pptx']).name} ({r['slides']}장)")
            success_count += 1

    log(f"\n성공: {success_count}/{len(NOTEBOOKS)}")
    log(f"저장 위치: {DOWNLOAD_DIR}")
    log("=" * 60)

    return results


if __name__ == "__main__":
    asyncio.run(main())
