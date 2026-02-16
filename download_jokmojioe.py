#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""족모지외반증 노트북 슬라이드 다운로드 및 PPTX 변환"""
import asyncio
import sys
from pathlib import Path
from datetime import datetime
import time

if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8')

DOWNLOAD_DIR = Path('G:/내 드라이브/notebooklm')
AUTH_DIR = Path.home() / '.notebooklm-mcp-cli'
BROWSER_PROFILE = AUTH_DIR / 'browser_profile'
NOTEBOOK_NAME = '족모지외반증'


async def main():
    from playwright.async_api import async_playwright

    DOWNLOAD_DIR.mkdir(parents=True, exist_ok=True)

    print('NotebookLM 슬라이드 다운로드')
    print('=' * 50)

    async with async_playwright() as p:
        context = await p.chromium.launch_persistent_context(
            user_data_dir=str(BROWSER_PROFILE),
            headless=False,
            downloads_path=str(DOWNLOAD_DIR),
            accept_downloads=True,
            args=['--disable-blink-features=AutomationControlled', '--start-maximized'],
            viewport={'width': 1920, 'height': 1080},
        )

        page = context.pages[0] if context.pages else await context.new_page()

        # NotebookLM 접속
        print('[1/5] NotebookLM 접속...')
        await page.goto('https://notebooklm.google.com', wait_until='domcontentloaded', timeout=60000)
        await asyncio.sleep(5)

        # 노트북 열기
        print('[2/5] 노트북 열기...')
        notebook_link = await page.query_selector('text=족모지외반증')
        if notebook_link:
            await notebook_link.click()
            print('  ✓ 노트북 발견')
            await asyncio.sleep(8)
        else:
            print('  ❌ 노트북 없음')
            await context.close()
            return None

        await page.screenshot(path=str(DOWNLOAD_DIR / 'ready_to_download.png'))
        print('  스크린샷: ready_to_download.png')

        # 다운로드 시도
        print('[3/5] 슬라이드 다운로드 시도...')

        downloaded = False
        download_path = None
        start_time = time.time()

        # 모든 메뉴 버튼 시도
        menu_buttons = await page.query_selector_all('[aria-haspopup="menu"]')
        print(f'  메뉴 버튼: {len(menu_buttons)}개')

        for i, btn in enumerate(menu_buttons):
            try:
                box = await btn.bounding_box()
                if not box:
                    continue

                # 스튜디오 패널의 슬라이드 카드 영역 (x > 1050, 280 < y < 400)
                if box['x'] > 1050 and 280 < box['y'] < 400:
                    print(f'  [{i}] 메뉴 버튼 클릭 (x={int(box["x"])}, y={int(box["y"])})')
                    await btn.click()
                    await asyncio.sleep(1)

                    # Download 메뉴 찾기
                    download_item = await page.query_selector('[role="menuitem"]:has-text("Download")')
                    if download_item:
                        print('  ✓ Download 메뉴 발견!')

                        try:
                            async with page.expect_download(timeout=30000) as download_info:
                                await download_item.click()
                                await asyncio.sleep(1)

                                # PDF 선택
                                pdf_item = await page.query_selector('[role="menuitem"]:has-text("PDF")')
                                if pdf_item:
                                    print('  PDF 선택...')
                                    await pdf_item.click()

                            download = await download_info.value
                            filename = f"{NOTEBOOK_NAME}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
                            download_path = DOWNLOAD_DIR / filename
                            await download.save_as(str(download_path))
                            print(f'  ✓ 다운로드 완료: {download_path}')
                            downloaded = True
                            break
                        except Exception as e:
                            print(f'  다운로드 오류: {str(e)[:50]}')
                    else:
                        await page.keyboard.press('Escape')
                        await asyncio.sleep(0.5)
            except Exception as e:
                continue

        # 자동 다운로드 실패 시 수동 다운로드 안내
        if not downloaded:
            print()
            print('=' * 50)
            print('수동 다운로드가 필요합니다!')
            print('=' * 50)
            print()
            print('브라우저에서 다음 단계를 따라주세요:')
            print('1. 오른쪽 스튜디오 패널에서')
            print('2. "무지외반증 원인과 진단 및 치료 관리 가이드" 카드의')
            print('3. 점 3개 메뉴(⋮) 클릭')
            print('4. "Download" 선택')
            print('5. "PDF" 선택')
            print()
            print(f'저장 위치: {DOWNLOAD_DIR}')
            print()
            print('다운로드 파일 감지 대기 중... (120초)')

            # 새 파일 감지 대기
            while time.time() - start_time < 120:
                pdf_files = list(DOWNLOAD_DIR.glob('*.pdf'))
                new_files = [f for f in pdf_files if f.stat().st_mtime > start_time]
                if new_files:
                    download_path = sorted(new_files, key=lambda x: x.stat().st_mtime)[-1]
                    print(f'  ✓ 다운로드 감지: {download_path.name}')
                    downloaded = True
                    break
                await asyncio.sleep(2)
                elapsed = int(time.time() - start_time)
                print(f'\r  대기 중... {elapsed}초', end='', flush=True)

        await context.close()

        # PDF → PPTX 변환
        if download_path and download_path.exists():
            print()
            print('[4/5] PDF → PPTX 변환...')

            import fitz
            from pptx import Presentation
            from pptx.util import Inches
            import io

            pptx_path = download_path.with_suffix('.pptx')

            doc = fitz.open(download_path)
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            blank_layout = prs.slide_layouts[6]

            for page_num in range(len(doc)):
                pg = doc[page_num]
                mat = fitz.Matrix(2, 2)
                pix = pg.get_pixmap(matrix=mat)
                img_data = pix.tobytes('png')

                slide = prs.slides.add_slide(blank_layout)
                img_stream = io.BytesIO(img_data)
                slide.shapes.add_picture(
                    img_stream,
                    Inches(0), Inches(0),
                    width=prs.slide_width,
                    height=prs.slide_height
                )

            doc.close()
            prs.save(pptx_path)

            print(f'  ✓ 변환 완료: {pptx_path}')
            print(f'  슬라이드 수: {len(prs.slides)}')

            print()
            print('[5/5] 완료!')
            print('=' * 50)
            print(f'PDF: {download_path}')
            print(f'PPTX: {pptx_path}')

            return str(pptx_path)

        return None


if __name__ == "__main__":
    result = asyncio.run(main())
    if result:
        print()
        print('✓ 작업 완료!')
    else:
        print()
        print('❌ 작업 실패')
