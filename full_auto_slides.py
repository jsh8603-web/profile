#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
NotebookLM 완전 자동화 - 슬라이드 생성/다운로드/변환
"""
import asyncio
import sys
import os
import io
import re
from pathlib import Path
from datetime import datetime

os.environ['PYTHONUNBUFFERED'] = '1'
if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8', line_buffering=True)

AUTH_DIR = Path.home() / ".notebooklm-mcp-cli"
PERSISTENT_PROFILE = AUTH_DIR / "persistent_browser"
DOWNLOAD_DIR = Path("G:/내 드라이브/notebooklm")

NOTEBOOKS = [
    ("족저근막염", "족저근막염"),
    ("족모지외반증", "족모지외반증"),
    ("발톱무좀", "발톱무좀"),
]


def log(msg):
    print(msg, flush=True)


def safe_filename(name):
    return re.sub(r'[<>:"/\\|?*]', '', name).strip()


def pdf_to_pptx(pdf_path, output_name=None):
    import fitz
    from pptx import Presentation
    from pptx.util import Inches

    pdf_path = Path(pdf_path)
    output_path = DOWNLOAD_DIR / f"{safe_filename(output_name or pdf_path.stem)}.pptx"

    doc = fitz.open(pdf_path)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for page_num in range(len(doc)):
        page = doc[page_num]
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img_data = pix.tobytes("png")
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(io.BytesIO(img_data), Inches(0), Inches(0),
                                  width=prs.slide_width, height=prs.slide_height)

    doc.close()
    prs.save(output_path)
    return output_path, len(prs.slides)


async def click_create_slides(page):
    """슬라이드 생성 버튼 클릭 (여러 방법 시도)"""
    log("  슬라이드 생성 시작...")

    # 스크린샷 (디버깅용)
    await page.screenshot(path=str(DOWNLOAD_DIR / "before_create.png"))

    # 방법 1: Studio 패널에서 Create 버튼
    try:
        # + 버튼 또는 Create 버튼 찾기
        add_buttons = await page.query_selector_all('button')
        for btn in add_buttons:
            try:
                text = await btn.inner_text()
                aria = await btn.get_attribute('aria-label') or ''
                if any(x in text.lower() for x in ['create', 'add', '만들기', '추가']) or \
                   any(x in aria.lower() for x in ['create', 'add']):
                    log(f"    Create 버튼 발견: {text[:20]}")
                    await btn.click()
                    await asyncio.sleep(2)
                    break
            except:
                continue
    except Exception as e:
        log(f"    방법1 실패: {e}")

    # 방법 2: 메뉴에서 Slides 선택
    await asyncio.sleep(1)
    try:
        menu_items = await page.query_selector_all('[role="menuitem"], [role="option"]')
        for item in menu_items:
            try:
                text = await item.inner_text()
                if 'slide' in text.lower() or '슬라이드' in text:
                    log(f"    Slides 메뉴 발견: {text[:20]}")
                    await item.click()
                    await asyncio.sleep(2)
                    break
            except:
                continue
    except Exception as e:
        log(f"    방법2 실패: {e}")

    # 방법 3: 언어 설정 (한국어)
    await asyncio.sleep(1)
    try:
        # 드롭다운/select 찾기
        selects = await page.query_selector_all('select, [role="combobox"], [role="listbox"], button')
        for sel in selects:
            try:
                text = await sel.inner_text()
                if 'english' in text.lower() or 'language' in text.lower() or '영어' in text:
                    log(f"    언어 드롭다운 발견")
                    await sel.click()
                    await asyncio.sleep(1)

                    # 한국어 선택
                    options = await page.query_selector_all('[role="option"], li, option')
                    for opt in options:
                        opt_text = await opt.inner_text()
                        if '한국어' in opt_text or 'korean' in opt_text.lower():
                            log(f"    한국어 선택")
                            await opt.click()
                            await asyncio.sleep(1)
                            break
                    break
            except:
                continue
    except Exception as e:
        log(f"    방법3 실패: {e}")

    # 방법 4: Generate/생성 버튼 클릭
    await asyncio.sleep(1)
    try:
        buttons = await page.query_selector_all('button')
        for btn in buttons:
            try:
                text = await btn.inner_text()
                if any(x in text.lower() for x in ['generate', '생성', 'create']):
                    log(f"    Generate 버튼 발견: {text[:20]}")
                    await btn.click()
                    log("    생성 요청!")
                    return True
            except:
                continue
    except Exception as e:
        log(f"    방법4 실패: {e}")

    return False


async def wait_for_slides_ready(page, timeout=180):
    """슬라이드 생성 완료 대기"""
    log(f"  슬라이드 생성 대기... (최대 {timeout}초)")

    start = asyncio.get_event_loop().time()
    last_log = 0

    while asyncio.get_event_loop().time() - start < timeout:
        await asyncio.sleep(5)
        elapsed = int(asyncio.get_event_loop().time() - start)

        # 30초마다 로그
        if elapsed - last_log >= 30:
            log(f"    대기 중... {elapsed}초")
            last_log = elapsed

        # 다운로드 메뉴 확인
        try:
            menu_buttons = await page.query_selector_all('[aria-haspopup="menu"]')
            for btn in menu_buttons[-5:]:
                try:
                    await btn.click(force=True)
                    await asyncio.sleep(0.5)

                    dl_item = await page.query_selector('[role="menuitem"]:has-text("Download")')
                    await page.keyboard.press('Escape')

                    if dl_item:
                        log(f"    완료! ({elapsed}초)")
                        return True
                except:
                    await page.keyboard.press('Escape')
        except:
            pass

    log(f"    타임아웃 ({timeout}초)")
    await page.screenshot(path=str(DOWNLOAD_DIR / "timeout_state.png"))
    return False


async def download_pdf(page, display_name):
    """PDF 다운로드"""
    log("  PDF 다운로드...")
    pdf_path = None

    # 메뉴 버튼들
    menu_buttons = await page.query_selector_all(
        '[aria-haspopup="menu"], button[aria-label*="more"], button[aria-label*="More"]'
    )
    log(f"    메뉴 버튼 {len(menu_buttons)}개")

    for btn in menu_buttons[-15:]:
        try:
            await btn.click(force=True)
            await asyncio.sleep(1.5)

            dl_item = await page.query_selector(
                '[role="menuitem"]:has-text("Download PDF"), '
                '[role="menuitem"]:has-text("Download"), '
                '[role="menuitem"]:has-text("다운로드")'
            )

            if dl_item:
                timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                filename = f"{safe_filename(display_name)}_{timestamp}.pdf"

                async with page.expect_download(timeout=60000) as download_info:
                    await dl_item.click()

                download = await download_info.value
                pdf_path = DOWNLOAD_DIR / filename
                await download.save_as(str(pdf_path))
                log(f"    PDF: {filename}")
                return pdf_path
        except:
            await page.keyboard.press('Escape')

    # 좌표 기반
    if not pdf_path:
        log("    좌표 시도...")
        for x, y in [(1846, 365), (1800, 400), (1850, 350)]:
            try:
                await page.mouse.click(x, y)
                await asyncio.sleep(1.5)

                items = await page.query_selector_all('[role="menuitem"]')
                for item in items:
                    text = await item.inner_text()
                    if 'Download' in text or '다운로드' in text:
                        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                        filename = f"{safe_filename(display_name)}_{timestamp}.pdf"

                        async with page.expect_download(timeout=60000) as download_info:
                            await item.click()

                        download = await download_info.value
                        pdf_path = DOWNLOAD_DIR / filename
                        await download.save_as(str(pdf_path))
                        log(f"    PDF (좌표): {filename}")
                        return pdf_path
            except:
                await page.keyboard.press('Escape')

    return None


async def process_notebook(page, keyword, display_name, idx, total):
    """단일 노트북 처리"""
    result = {'name': display_name, 'pdf': None, 'pptx': None, 'slides': 0}

    log(f"\n{'='*50}")
    log(f"[{idx}/{total}] {display_name}")
    log('='*50)

    try:
        # 홈
        await page.goto("https://notebooklm.google.com/", wait_until='networkidle', timeout=30000)
        await asyncio.sleep(3)

        # 노트북 찾기
        log(f"  노트북 찾기: '{keyword}'...")
        notebook_link = None

        elements = await page.query_selector_all('a, div[role="button"], [class*="notebook"], tr')
        for el in elements:
            try:
                text = await el.inner_text()
                if keyword in text:
                    notebook_link = el
                    log("    발견!")
                    break
            except:
                continue

        if not notebook_link:
            log("    실패!")
            return result

        await notebook_link.click()
        await asyncio.sleep(5)

        # 슬라이드 생성
        created = await click_create_slides(page)

        if created:
            # 생성 대기
            await wait_for_slides_ready(page, timeout=180)

        # 다운로드
        pdf_path = await download_pdf(page, display_name)

        if pdf_path and pdf_path.exists():
            result['pdf'] = str(pdf_path)
            log("  PPTX 변환...")
            pptx_path, slide_count = pdf_to_pptx(pdf_path, display_name)
            result['pptx'] = str(pptx_path)
            result['slides'] = slide_count
            log(f"    PPTX: {pptx_path.name} ({slide_count}장)")
        else:
            log("    다운로드 실패")

    except Exception as e:
        log(f"  오류: {e}")

    return result


async def main():
    from playwright.async_api import async_playwright

    log("=" * 60)
    log("NotebookLM 15장 한글 슬라이드 완전 자동화")
    log("=" * 60)

    DOWNLOAD_DIR.mkdir(parents=True, exist_ok=True)
    results = []

    async with async_playwright() as p:
        log("\n브라우저 시작...")

        context = await p.chromium.launch_persistent_context(
            user_data_dir=str(PERSISTENT_PROFILE),
            headless=False,
            downloads_path=str(DOWNLOAD_DIR),
            accept_downloads=True,
            args=['--disable-blink-features=AutomationControlled'],
            viewport={'width': 1920, 'height': 1080},
        )

        page = context.pages[0] if context.pages else await context.new_page()

        await page.goto("https://notebooklm.google.com/", wait_until='networkidle', timeout=30000)
        await asyncio.sleep(3)

        if "accounts.google.com" in page.url:
            log("로그인 대기... (5분)")
            for i in range(60):
                await asyncio.sleep(5)
                if "notebooklm.google.com" in page.url and "accounts.google" not in page.url:
                    break
            else:
                await context.close()
                return []

        log("로그인 완료!\n")

        for idx, (keyword, display_name) in enumerate(NOTEBOOKS, 1):
            result = await process_notebook(page, keyword, display_name, idx, len(NOTEBOOKS))
            results.append(result)
            await asyncio.sleep(3)

        await context.close()

    # 결과
    log("\n" + "=" * 60)
    log("결과:")
    log("=" * 60)

    success = sum(1 for r in results if r.get('pptx'))
    for r in results:
        status = "O" if r.get('pptx') else "X"
        log(f"  [{status}] {r['name']}")
        if r.get('pptx'):
            log(f"       {Path(r['pptx']).name} ({r['slides']}장)")

    log(f"\n성공: {success}/{len(NOTEBOOKS)}")
    log(f"저장: {DOWNLOAD_DIR}")

    return results


if __name__ == "__main__":
    asyncio.run(main())
