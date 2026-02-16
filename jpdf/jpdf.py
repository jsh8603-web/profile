#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
JPDF v1.1 - PDF/이미지 → 편집 가능 PPTX 변환기

지원 파일: PDF, PNG, JPG, JPEG
"""
import io
import os
import sys
import base64
import requests
import cv2
import numpy as np
from pathlib import Path
from typing import List, Tuple, Optional
from dataclasses import dataclass

__version__ = "1.1.0"

if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8')

# .env.local 로드
try:
    from dotenv import load_dotenv
    env_path = Path(__file__).parent / '.env.local'
    if env_path.exists():
        load_dotenv(env_path)
    else:
        load_dotenv()
except ImportError:
    pass


# 지원 폰트 목록
SUPPORTED_FONTS = [
    'Arial',
    'Calibri',
    'Times New Roman',
    'Georgia',
    'Verdana',
    'Tahoma',
    'Trebuchet MS',
    'Malgun Gothic',  # 맑은 고딕
    'NanumGothic',    # 나눔고딕
    'NanumBarunGothic',
]


@dataclass
class TextBlock:
    """텍스트 블록"""
    text: str
    x: int
    y: int
    width: int
    height: int
    font_size: float = 24.0
    line_height: float = 1.2
    alignment: str = "left"


@dataclass
class PageData:
    """페이지 데이터"""
    page_num: int
    original_image: np.ndarray
    cleaned_image: np.ndarray
    text_blocks: List[TextBlock]
    width: int
    height: int


class JPDF:
    """JPDF - PDF/이미지 → 편집 가능 PPTX 변환기"""

    def __init__(self, api_key: str = None):
        self.api_key = api_key or os.getenv('GOOGLE_VISION_API_KEY')
        if not self.api_key:
            raise ValueError(
                "Google Vision API 키가 필요합니다.\n"
                ".env.local 파일에 GOOGLE_VISION_API_KEY를 설정하거나\n"
                "--api-key 옵션을 사용하세요."
            )
        self.api_url = "https://vision.googleapis.com/v1/images:annotate"

    def _ocr_image(self, image: np.ndarray) -> List[TextBlock]:
        """이미지에서 텍스트 블록 추출"""
        img_height, img_width = image.shape[:2]

        rgb = cv2.cvtColor(image, cv2.COLOR_BGR2RGB)
        _, buffer = cv2.imencode('.png', rgb)
        img_base64 = base64.b64encode(buffer).decode('utf-8')

        payload = {
            "requests": [{
                "image": {"content": img_base64},
                "features": [{"type": "DOCUMENT_TEXT_DETECTION"}]
            }]
        }

        response = requests.post(
            f"{self.api_url}?key={self.api_key}",
            json=payload,
            timeout=60
        )
        result = response.json()

        if 'error' in result:
            raise Exception(f"Vision API 오류: {result['error']}")

        blocks = []
        responses = result.get('responses', [{}])
        if not responses:
            return blocks

        full_annotation = responses[0].get('fullTextAnnotation', {})

        for page in full_annotation.get('pages', []):
            for block in page.get('blocks', []):
                block_text = ""
                word_heights = []
                line_xs = []

                for para in block.get('paragraphs', []):
                    para_text = ""
                    para_word_xs = []

                    for word in para.get('words', []):
                        word_text = "".join(
                            s.get('text', '') for s in word.get('symbols', [])
                        )
                        para_text += word_text + " "

                        word_vertices = word.get('boundingBox', {}).get('vertices', [])
                        if len(word_vertices) >= 4:
                            word_ys = [v.get('y', 0) for v in word_vertices]
                            word_xs = [v.get('x', 0) for v in word_vertices]
                            word_height = max(word_ys) - min(word_ys)
                            if word_height > 0:
                                word_heights.append(word_height)
                            para_word_xs.append(min(word_xs))

                    block_text += para_text.strip() + "\n"
                    if para_word_xs:
                        line_xs.append(min(para_word_xs))

                block_text = block_text.strip()
                if not block_text:
                    continue

                vertices = block.get('boundingBox', {}).get('vertices', [])
                if len(vertices) >= 4:
                    xs = [v.get('x', 0) for v in vertices]
                    ys = [v.get('y', 0) for v in vertices]
                    block_x = min(xs)
                    block_y = min(ys)
                    block_width = max(xs) - min(xs)
                    block_height = max(ys) - min(ys)

                    # 폰트 크기 추정
                    if word_heights:
                        median_height = sorted(word_heights)[len(word_heights) // 2]
                        font_size = (median_height / 2.0) * 0.5 + 1
                        font_size = max(8, min(48, font_size))
                    else:
                        font_size = 12.0

                    # 정렬 감지
                    alignment = "left"
                    if line_xs and block_width > 0:
                        avg_line_x = sum(line_xs) / len(line_xs)
                        relative_pos = (avg_line_x - block_x) / block_width
                        if relative_pos > 0.35:
                            alignment = "center"
                        elif relative_pos > 0.6:
                            alignment = "right"

                    # 줄 높이 비율
                    line_count = block_text.count('\n') + 1
                    if line_count > 1 and font_size > 0:
                        line_height = (block_height / line_count) / (font_size * 1.33)
                        line_height = max(1.0, min(2.0, line_height))
                    else:
                        line_height = 1.2

                    blocks.append(TextBlock(
                        text=block_text,
                        x=block_x,
                        y=block_y,
                        width=block_width,
                        height=block_height,
                        font_size=font_size,
                        line_height=line_height,
                        alignment=alignment
                    ))

        return self._merge_nearby_blocks(blocks)

    def _merge_nearby_blocks(
        self,
        blocks: List[TextBlock],
        y_threshold: int = 50,
        x_threshold: int = 100
    ) -> List[TextBlock]:
        """가까운 블록들을 하나로 병합"""
        if not blocks:
            return blocks

        sorted_blocks = sorted(blocks, key=lambda b: (b.y, b.x))
        merged = []
        current_group = [sorted_blocks[0]]

        for block in sorted_blocks[1:]:
            last = current_group[-1]

            y_close = (block.y - (last.y + last.height)) < y_threshold
            x_overlap = not (block.x > last.x + last.width + x_threshold or
                           block.x + block.width < last.x - x_threshold)

            if y_close and x_overlap:
                current_group.append(block)
            else:
                merged.append(self._merge_block_group(current_group))
                current_group = [block]

        if current_group:
            merged.append(self._merge_block_group(current_group))

        return merged

    def _merge_block_group(self, group: List[TextBlock]) -> TextBlock:
        """블록 그룹을 하나로 병합"""
        if len(group) == 1:
            return group[0]

        texts = [b.text for b in group]
        merged_text = "\n".join(texts)

        min_x = min(b.x for b in group)
        min_y = min(b.y for b in group)
        max_x = max(b.x + b.width for b in group)
        max_y = max(b.y + b.height for b in group)

        avg_font_size = sum(b.font_size for b in group) / len(group)
        avg_line_height = sum(b.line_height for b in group) / len(group)

        return TextBlock(
            text=merged_text,
            x=min_x,
            y=min_y,
            width=max_x - min_x,
            height=max_y - min_y,
            font_size=avg_font_size,
            line_height=avg_line_height,
            alignment=group[0].alignment
        )

    def _inpaint_image(
        self,
        image: np.ndarray,
        blocks: List[TextBlock],
        padding: int = 10,
        radius: int = 7
    ) -> np.ndarray:
        """텍스트 영역 제거 및 배경 복원"""
        if not blocks:
            return image.copy()

        mask = np.zeros(image.shape[:2], dtype=np.uint8)
        for b in blocks:
            x1 = max(0, b.x - padding)
            y1 = max(0, b.y - padding)
            x2 = min(image.shape[1], b.x + b.width + padding)
            y2 = min(image.shape[0], b.y + b.height + padding)
            cv2.rectangle(mask, (x1, y1), (x2, y2), 255, -1)

        return cv2.inpaint(image, mask, radius, cv2.INPAINT_TELEA)

    def load_image(self, path: str) -> np.ndarray:
        """이미지 파일 로드"""
        path = Path(path)
        if not path.exists():
            raise FileNotFoundError(f"파일을 찾을 수 없습니다: {path}")

        img = cv2.imread(str(path))
        if img is None:
            raise ValueError(f"이미지를 읽을 수 없습니다: {path}")
        return img

    def process_pdf(
        self,
        pdf_path: str,
        zoom: float = 2.0,
        inpaint: bool = True,
        padding: int = 10,
        inpaint_radius: int = 7
    ) -> List[PageData]:
        """PDF 처리"""
        import fitz

        pdf_path = Path(pdf_path)
        if not pdf_path.exists():
            raise FileNotFoundError(f"파일을 찾을 수 없습니다: {pdf_path}")

        doc = fitz.open(str(pdf_path))
        pages = []

        total = len(doc)
        print(f"📄 JPDF: {pdf_path.name} ({total}페이지)")

        for i in range(total):
            print(f"  [{i+1}/{total}] OCR + 처리 중...", end=" ", flush=True)

            page = doc[i]
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat)

            img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(
                pix.height, pix.width, pix.n
            )
            if pix.n == 4:
                img = cv2.cvtColor(img, cv2.COLOR_RGBA2BGR)
            else:
                img = cv2.cvtColor(img, cv2.COLOR_RGB2BGR)

            blocks = self._ocr_image(img)

            if inpaint and blocks:
                cleaned = self._inpaint_image(img, blocks, padding, inpaint_radius)
            else:
                cleaned = img.copy()

            pages.append(PageData(
                page_num=i + 1,
                original_image=img,
                cleaned_image=cleaned,
                text_blocks=blocks,
                width=pix.width,
                height=pix.height
            ))

            print(f"텍스트 {len(blocks)}개")

        doc.close()
        return pages

    def process_image(
        self,
        image_path: str,
        inpaint: bool = True,
        padding: int = 10,
        inpaint_radius: int = 7
    ) -> PageData:
        """이미지 파일 처리"""
        print(f"🖼️ JPDF: {Path(image_path).name}")
        print(f"  OCR + 처리 중...", end=" ", flush=True)

        img = self.load_image(image_path)
        blocks = self._ocr_image(img)

        if inpaint and blocks:
            cleaned = self._inpaint_image(img, blocks, padding, inpaint_radius)
        else:
            cleaned = img.copy()

        print(f"텍스트 {len(blocks)}개")

        return PageData(
            page_num=1,
            original_image=img,
            cleaned_image=cleaned,
            text_blocks=blocks,
            width=img.shape[1],
            height=img.shape[0]
        )

    def create_pptx(
        self,
        pages: List[PageData],
        output_path: str,
        font_size: int = None,
        font_family: str = 'Arial',
        font_color: Tuple[int, int, int] = (0x33, 0x33, 0x33)
    ) -> Path:
        """PPTX 생성"""
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]

        for page_idx, page in enumerate(pages):
            slide = prs.slides.add_slide(blank_layout)

            scale_x = prs.slide_width.emu / page.width
            scale_y = prs.slide_height.emu / page.height

            # 배경 이미지
            _, buffer = cv2.imencode('.png', page.cleaned_image)
            slide.shapes.add_picture(
                io.BytesIO(buffer.tobytes()),
                Inches(0), Inches(0),
                width=prs.slide_width,
                height=prs.slide_height
            )

            # 텍스트 박스
            for block_idx, block in enumerate(page.text_blocks):
                left = Emu(int(block.x * scale_x))
                top = Emu(int(block.y * scale_y))
                width = Emu(int(block.width * scale_x))
                height = Emu(int(block.height * scale_y))

                textbox = slide.shapes.add_textbox(left, top, width, height)
                tf = textbox.text_frame
                tf.word_wrap = True
                tf.auto_size = None

                alignment_map = {
                    'left': PP_ALIGN.LEFT,
                    'center': PP_ALIGN.CENTER,
                    'right': PP_ALIGN.RIGHT
                }

                lines = block.text.split('\n')
                for i, line in enumerate(lines):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()

                    p.text = line
                    p.alignment = alignment_map.get(block.alignment, PP_ALIGN.LEFT)

                    # 첫 슬라이드 제목: 24pt 볼드
                    is_first_slide_title = (page_idx == 0 and block_idx == 0)

                    if font_size:
                        size = font_size
                    elif is_first_slide_title:
                        size = 24
                    else:
                        size = int(block.font_size)

                    p.font.size = Pt(size)
                    p.font.name = font_family
                    p.font.color.rgb = RGBColor(*font_color)
                    p.font.bold = is_first_slide_title
                    p.line_spacing = block.line_height

        output_path = Path(output_path)
        prs.save(str(output_path))
        return output_path

    def convert(
        self,
        input_path: str,
        output_path: str = None,
        inpaint: bool = True,
        zoom: float = 2.0,
        padding: int = 10,
        inpaint_radius: int = 7,
        font_size: int = None,
        font_family: str = 'Arial'
    ) -> Tuple[Path, int]:
        """파일 → 편집 가능 PPTX 변환"""
        input_path = Path(input_path)
        ext = input_path.suffix.lower()

        if output_path is None:
            output_path = input_path.with_name(f"{input_path.stem}_편집가능.pptx")

        # 파일 타입별 처리
        if ext == '.pdf':
            pages = self.process_pdf(
                str(input_path), zoom=zoom, inpaint=inpaint,
                padding=padding, inpaint_radius=inpaint_radius
            )
        elif ext in ['.png', '.jpg', '.jpeg']:
            page = self.process_image(
                str(input_path), inpaint=inpaint,
                padding=padding, inpaint_radius=inpaint_radius
            )
            pages = [page]
        else:
            raise ValueError(f"지원하지 않는 파일 형식: {ext}")

        result_path = self.create_pptx(pages, output_path, font_size, font_family)

        print(f"\n✅ 변환 완료!")
        print(f"   파일: {result_path}")
        print(f"   슬라이드: {len(pages)}장")

        return result_path, len(pages)


def convert_images_to_pptx(
    image_paths: List[str],
    output_path: str,
    api_key: str = None,
    editable: bool = True,
    font_size: int = None,
    font_family: str = 'Arial'
) -> Tuple[Path, int]:
    """여러 이미지를 하나의 PPTX로 변환"""
    converter = JPDF(api_key)
    pages = []

    for path in image_paths:
        page = converter.process_image(path, inpaint=editable)
        pages.append(page)

    result_path = converter.create_pptx(pages, output_path, font_size, font_family)
    return result_path, len(pages)


def main():
    """CLI"""
    import argparse

    parser = argparse.ArgumentParser(
        prog="jpdf",
        description=f"JPDF v{__version__} - PDF/이미지 → 편집 가능 PPTX 변환기",
        epilog="예시: python jpdf.py slides.pdf -o output.pptx"
    )
    parser.add_argument("input_path", help="입력 파일 (PDF, PNG, JPG)")
    parser.add_argument("-o", "--output", help="출력 PPTX 경로")
    parser.add_argument("--api-key", help="Google Vision API 키")
    parser.add_argument("--no-inpaint", action="store_true", help="텍스트 제거 안함")
    parser.add_argument("--zoom", type=float, default=2.0, help="확대 비율 (기본: 2.0)")
    parser.add_argument("--padding", type=int, default=10, help="패딩 (기본: 10)")
    parser.add_argument("--inpaint-radius", type=int, default=7, help="Inpaint 반경 (기본: 7)")
    parser.add_argument("--font-size", type=int, help="폰트 크기 고정")
    parser.add_argument("--font-family", default='Arial', help=f"폰트 (기본: Arial)")
    parser.add_argument("-v", "--version", action="version", version=f"JPDF v{__version__}")

    args = parser.parse_args()

    try:
        converter = JPDF(args.api_key)
        converter.convert(
            args.input_path,
            args.output,
            inpaint=not args.no_inpaint,
            zoom=args.zoom,
            padding=args.padding,
            inpaint_radius=args.inpaint_radius,
            font_size=args.font_size,
            font_family=args.font_family
        )
    except Exception as e:
        print(f"❌ 오류: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
