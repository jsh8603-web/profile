#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
NotebookLM 한글 슬라이드 생성 v2
- 슬라이드 자료 카드의 연필 아이콘 클릭 → 언어 설정 → 재생성
- 또는 새로 생성
"""
import asyncio
import sys
import os
import io
import re
from pathlib import Path
from datetime import datetime

os.environ['PYTHONUNBUFFERED'] = '1'
if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8', line_buffering=True)

AUTH_DIR = Path.home() / ".notebooklm-mcp-cli"
PERSISTENT_PROFILE = AUTH_DIR / "persistent_browser"
DOWNLOAD_DIR = Path("G:/내 드라이브/notebooklm")

NOTEBOOKS = [
    ("족저근막염", "족저근막염"),
    ("족모지외반증", "족모지외반증"),
    ("발톱무좀", "발톱무좀"),
]


def log(msg):
    print(msg, flush=True)


def safe_filename(name):
    return re.sub(r'[<>:"/\\|?*]', '', name).strip()


def pdf_to_pptx(pdf_path, output_name=None):
    import fitz
    from pptx import Presentation
    from pptx.util import Inches

    pdf_path = Path(pdf_path)
    output_path = DOWNLOAD_DIR / f"{safe_filename(output_name or pdf_path.stem)}.pptx"

    doc = fitz.open(pdf_path)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for page_num in range(len(doc)):
        page = doc[page_num]
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img_data = pix.tobytes("png")
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(io.BytesIO(img_data), Inches(0), Inches(0),
                                  width=prs.slide_width, height=prs.slide_height)

    doc.close()
    prs.save(output_path)
    return output_path, len(prs.slides)


async def create_korean_slides(page):
    """한국어 슬라이드 생성 (연필 아이콘 또는 새로 생성)"""
    log("  한국어 슬라이드 생성...")

    # 스튜디오 패널에서 슬라이드 자료 카드 찾기
    # 연필 아이콘(편집) 클릭하여 설정 변경

    # 방법 1: 슬라이드 자료 카드의 연필 아이콘 클릭
    try:
        # 슬라이드 자료 영역 찾기
        slides_area = await page.query_selector('text="슬라이드 자료"')
        if slides_area:
            # 근처의 연필 아이콘 찾기
            parent = await slides_area.evaluate_handle('el => el.parentElement')
            edit_btn = await parent.query_selector('button, [role="button"]')
            if edit_btn:
                log("    슬라이드 설정 열기...")
                await edit_btn.click()
                await asyncio.sleep(2)
    except:
        pass

    # 방법 2: 좌표로 연필 아이콘 클릭 (슬라이드 자료 카드 옆)
    # 스크린샷 기준 대략 (1295, 215) 부근
    try:
        log("    연필 아이콘 클릭 시도...")
        await page.mouse.click(1295, 215)
        await asyncio.sleep(2)
        await page.screenshot(path=str(DOWNLOAD_DIR / "after_edit_click.png"))
    except:
        pass

    # 설정 다이얼로그에서 언어 변경
    await asyncio.sleep(1)

    # 언어 드롭다운 찾기
    try:
        # 모든 드롭다운/select 요소
        all_elements = await page.query_selector_all('button, select, [role="combobox"], [role="listbox"]')
        for el in all_elements:
            try:
                text = await el.inner_text()
                # English, 영어, Language 등 언어 관련 텍스트 찾기
                if any(x in text.lower() for x in ['english', 'language', '영어', '언어']):
                    log(f"    언어 드롭다운: {text[:20]}")
                    await el.click()
                    await asyncio.sleep(1)

                    # 한국어 옵션 찾기
                    options = await page.query_selector_all('[role="option"], [role="menuitem"], li, option')
                    for opt in options:
                        try:
                            opt_text = await opt.inner_text()
                            if '한국어' in opt_text or 'Korean' in opt_text:
                                log("    한국어 선택!")
                                await opt.click()
                                await asyncio.sleep(1)
                                break
                        except:
                            continue
                    break
            except:
                continue
    except:
        pass

    # 생성/재생성 버튼 클릭
    await asyncio.sleep(1)
    try:
        buttons = await page.query_selector_all('button')
        for btn in buttons:
            try:
                text = await btn.inner_text()
                # Generate, Regenerate, 생성, 재생성 등
                if any(x in text for x in ['Generate', 'Regenerate', '생성', '재생성', 'Create']):
                    # "add 노트북 만들기" 제외
                    if '노트북' not in text and 'notebook' not in text.lower():
                        log(f"    생성 버튼: {text[:20]}")
                        await btn.click()
                        await asyncio.sleep(2)
                        return True
            except:
                continue
    except:
        pass

    return False


async def wait_and_download(page, display_name, timeout=180):
    """슬라이드 생성 완료 대기 및 다운로드"""
    log(f"  생성 대기... (최대 {timeout}초)")

    start = asyncio.get_event_loop().time()

    while asyncio.get_event_loop().time() - start < timeout:
        await asyncio.sleep(10)
        elapsed = int(asyncio.get_event_loop().time() - start)

        if elapsed % 30 == 0:
            log(f"    {elapsed}초...")

        # 완료 확인: 슬라이드 카드에 다운로드 메뉴 있는지
        try:
            # 모든 메뉴 버튼 확인
            menu_buttons = await page.query_selector_all('[aria-haspopup="menu"], [aria-label*="more"], [aria-label*="옵션"]')

            for btn in menu_buttons:
                try:
                    await btn.click(force=True)
                    await asyncio.sleep(1)

                    # Download 메뉴 항목 찾기
                    dl_items = await page.query_selector_all('[role="menuitem"]')
                    for item in dl_items:
                        try:
                            text = await item.inner_text()
                            if 'Download' in text or '다운로드' in text:
                                log(f"    다운로드 가능! ({elapsed}초)")

                                # 다운로드 실행
                                timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                                filename = f"{safe_filename(display_name)}_{timestamp}.pdf"

                                async with page.expect_download(timeout=60000) as download_info:
                                    await item.click()

                                download = await download_info.value
                                pdf_path = DOWNLOAD_DIR / filename
                                await download.save_as(str(pdf_path))
                                log(f"    PDF: {filename}")
                                return pdf_path
                        except:
                            continue

                    await page.keyboard.press('Escape')
                except:
                    await page.keyboard.press('Escape')
        except:
            pass

    log(f"    타임아웃 ({timeout}초)")
    return None


async def process_notebook(page, keyword, display_name, idx, total):
    """단일 노트북 처리"""
    result = {'name': display_name, 'pdf': None, 'pptx': None, 'slides': 0}

    log(f"\n{'='*55}")
    log(f"[{idx}/{total}] {display_name}")
    log('='*55)

    try:
        # 홈으로
        await page.goto("https://notebooklm.google.com/", wait_until='domcontentloaded', timeout=60000)
        await asyncio.sleep(3)

        # 노트북 찾기
        log(f"  노트북: '{keyword}'")
        notebook_link = None

        elements = await page.query_selector_all('a, div[role="button"], tr')
        for el in elements:
            try:
                text = await el.inner_text()
                if keyword in text:
                    notebook_link = el
                    break
            except:
                continue

        if not notebook_link:
            log("    찾기 실패!")
            return result

        log("    발견!")
        await notebook_link.click()
        await asyncio.sleep(5)

        # 한국어 슬라이드 생성
        await create_korean_slides(page)

        # 대기 및 다운로드
        pdf_path = await wait_and_download(page, display_name, timeout=180)

        if pdf_path and pdf_path.exists():
            result['pdf'] = str(pdf_path)
            log("  PPTX 변환...")
            pptx_path, slide_count = pdf_to_pptx(pdf_path, display_name)
            result['pptx'] = str(pptx_path)
            result['slides'] = slide_count
            log(f"    완료: {pptx_path.name} ({slide_count}장)")
        else:
            log("    다운로드 실패!")
            await page.screenshot(path=str(DOWNLOAD_DIR / f"fail_{safe_filename(display_name)}.png"))

    except Exception as e:
        log(f"  오류: {e}")

    return result


async def main():
    from playwright.async_api import async_playwright

    log("=" * 55)
    log("NotebookLM 15장 한글 슬라이드 생성 v2")
    log("=" * 55)

    DOWNLOAD_DIR.mkdir(parents=True, exist_ok=True)
    results = []

    async with async_playwright() as p:
        log("\n브라우저 시작...")

        context = await p.chromium.launch_persistent_context(
            user_data_dir=str(PERSISTENT_PROFILE),
            headless=False,
            downloads_path=str(DOWNLOAD_DIR),
            accept_downloads=True,
            args=['--disable-blink-features=AutomationControlled'],
            viewport={'width': 1920, 'height': 1080},
        )

        page = context.pages[0] if context.pages else await context.new_page()

        await page.goto("https://notebooklm.google.com/", wait_until='domcontentloaded', timeout=60000)
        await asyncio.sleep(3)

        if "accounts.google.com" in page.url:
            log("로그인 대기 중...")
            for _ in range(60):
                await asyncio.sleep(5)
                if "notebooklm.google.com" in page.url:
                    break
            else:
                await context.close()
                return []

        log("준비 완료!\n")

        for idx, (keyword, display_name) in enumerate(NOTEBOOKS, 1):
            result = await process_notebook(page, keyword, display_name, idx, len(NOTEBOOKS))
            results.append(result)
            await asyncio.sleep(3)

        await context.close()

    # 결과
    log("\n" + "=" * 55)
    log("결과:")
    log("=" * 55)

    success = sum(1 for r in results if r.get('pptx'))
    for r in results:
        status = "O" if r.get('pptx') else "X"
        log(f"  [{status}] {r['name']}")
        if r.get('pptx'):
            log(f"       {Path(r['pptx']).name} ({r['slides']}장)")

    log(f"\n성공: {success}/{len(NOTEBOOKS)}")
    log(f"저장: {DOWNLOAD_DIR}")

    return results


if __name__ == "__main__":
    asyncio.run(main())
