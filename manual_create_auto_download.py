#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
수동 슬라이드 생성 + 자동 다운로드/변환
- 사용자가 브라우저에서 슬라이드 생성
- 자동으로 다운로드 및 PPTX 변환
"""
import asyncio
import sys
import os
import io
import re
from pathlib import Path
from datetime import datetime

os.environ['PYTHONUNBUFFERED'] = '1'
if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8', line_buffering=True)

AUTH_DIR = Path.home() / ".notebooklm-mcp-cli"
PERSISTENT_PROFILE = AUTH_DIR / "persistent_browser"
DOWNLOAD_DIR = Path("G:/내 드라이브/notebooklm")

NOTEBOOKS = [
    ("족저근막염", "족저근막염"),
    ("족모지외반증", "족모지외반증"),
    ("발톱무좀", "발톱무좀"),
]


def log(msg):
    print(msg, flush=True)


def safe_filename(name):
    return re.sub(r'[<>:"/\\|?*]', '', name).strip()


def pdf_to_pptx(pdf_path, output_name=None):
    """PDF를 PPTX로 변환"""
    import fitz
    from pptx import Presentation
    from pptx.util import Inches

    pdf_path = Path(pdf_path)
    if output_name:
        output_path = DOWNLOAD_DIR / f"{safe_filename(output_name)}.pptx"
    else:
        output_path = pdf_path.with_suffix('.pptx')

    doc = fitz.open(pdf_path)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for page_num in range(len(doc)):
        page = doc[page_num]
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img_data = pix.tobytes("png")

        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            io.BytesIO(img_data),
            Inches(0), Inches(0),
            width=prs.slide_width,
            height=prs.slide_height
        )

    doc.close()
    prs.save(output_path)
    return output_path, len(prs.slides)


async def wait_for_user_action(prompt_msg):
    """사용자 입력 대기"""
    log(f"\n  >>> {prompt_msg}")
    log(f"  >>> 완료되면 Enter를 누르세요...")
    await asyncio.get_event_loop().run_in_executor(None, input)


async def download_pdf(page, display_name):
    """PDF 다운로드"""
    pdf_path = None

    # 메뉴 버튼 찾기
    menu_buttons = await page.query_selector_all(
        '[aria-haspopup="menu"], button[aria-label*="more"], button[aria-label*="More"]'
    )
    log(f"  메뉴 버튼 {len(menu_buttons)}개 발견")

    for menu_btn in menu_buttons[-15:]:
        try:
            await menu_btn.click(force=True)
            await asyncio.sleep(1.5)

            download_item = await page.query_selector(
                '[role="menuitem"]:has-text("Download PDF"), '
                '[role="menuitem"]:has-text("Download"), '
                '[role="menuitem"]:has-text("다운로드")'
            )

            if download_item:
                timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                filename = f"{safe_filename(display_name)}_{timestamp}.pdf"

                async with page.expect_download(timeout=60000) as download_info:
                    await download_item.click()

                download = await download_info.value
                pdf_path = DOWNLOAD_DIR / filename
                await download.save_as(str(pdf_path))
                log(f"  ✓ PDF 저장: {filename}")
                return pdf_path
        except:
            await page.keyboard.press('Escape')

    # 좌표 기반 시도
    if not pdf_path:
        log(f"  좌표 기반 시도...")
        coordinates = [(1846, 365), (1800, 400), (1850, 350), (1750, 380), (1820, 380)]

        for x, y in coordinates:
            try:
                await page.mouse.click(x, y)
                await asyncio.sleep(1.5)

                menu_items = await page.query_selector_all('[role="menuitem"]')
                for item in menu_items:
                    text = await item.inner_text()
                    if 'Download' in text or '다운로드' in text:
                        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                        filename = f"{safe_filename(display_name)}_{timestamp}.pdf"

                        async with page.expect_download(timeout=60000) as download_info:
                            await item.click()

                        download = await download_info.value
                        pdf_path = DOWNLOAD_DIR / filename
                        await download.save_as(str(pdf_path))
                        log(f"  ✓ PDF 저장 (좌표): {filename}")
                        return pdf_path
            except:
                await page.keyboard.press('Escape')

    return None


async def process_notebook(page, keyword, display_name, idx, total):
    """단일 노트북 처리"""
    result = {'name': display_name, 'pdf': None, 'pptx': None, 'slides': 0}

    log(f"\n{'='*60}")
    log(f"[{idx}/{total}] {display_name}")
    log('='*60)

    try:
        # 1. 홈으로 이동
        await page.goto("https://notebooklm.google.com/", wait_until='networkidle', timeout=30000)
        await asyncio.sleep(2)

        # 2. 노트북 찾기
        log(f"  노트북 찾기: '{keyword}'...")
        notebook_link = None

        all_elements = await page.query_selector_all('a, div[role="button"], [class*="notebook"], tr')
        for el in all_elements:
            try:
                text = await el.inner_text()
                if keyword in text:
                    notebook_link = el
                    log(f"  → 발견!")
                    break
            except:
                continue

        if not notebook_link:
            log(f"  ❌ 노트북 찾기 실패")
            return result

        # 3. 노트북 열기
        await notebook_link.click()
        await asyncio.sleep(5)
        log(f"  노트북 열림")

        # 4. 사용자에게 슬라이드 생성 요청
        log(f"\n  ===== 수동 작업 필요 =====")
        log(f"  1. 오른쪽 Studio 패널에서 '+ Create' 클릭")
        log(f"  2. 'Slides' 선택")
        log(f"  3. Language: 한국어 (Korean) 선택")
        log(f"  4. 'Generate' 클릭")
        log(f"  5. 슬라이드 생성 완료까지 대기 (1-2분)")
        log(f"  ============================")

        await wait_for_user_action("슬라이드 생성이 완료되면 Enter를 누르세요")

        # 5. PDF 다운로드
        log(f"\n  PDF 다운로드 중...")
        pdf_path = await download_pdf(page, display_name)

        if pdf_path and pdf_path.exists():
            result['pdf'] = str(pdf_path)

            # 6. PPTX 변환
            log(f"  PPTX 변환 중...")
            pptx_path, slide_count = pdf_to_pptx(pdf_path, display_name)
            result['pptx'] = str(pptx_path)
            result['slides'] = slide_count
            log(f"  ✓ PPTX: {pptx_path.name} ({slide_count}장)")
        else:
            log(f"  ❌ 다운로드 실패")
            log(f"  수동으로 다운로드해주세요 (메뉴 → Download PDF)")
            await wait_for_user_action("수동 다운로드 후 Enter")

            # 수동 다운로드된 파일 찾기
            recent_pdfs = sorted(DOWNLOAD_DIR.glob("*.pdf"), key=lambda x: x.stat().st_mtime, reverse=True)
            if recent_pdfs:
                pdf_path = recent_pdfs[0]
                result['pdf'] = str(pdf_path)
                log(f"  발견: {pdf_path.name}")

                pptx_path, slide_count = pdf_to_pptx(pdf_path, display_name)
                result['pptx'] = str(pptx_path)
                result['slides'] = slide_count
                log(f"  ✓ PPTX: {pptx_path.name} ({slide_count}장)")

    except Exception as e:
        log(f"  ❌ 오류: {e}")

    return result


async def main():
    from playwright.async_api import async_playwright

    log("=" * 60)
    log("NotebookLM 슬라이드 생성 (수동) + 다운로드/변환 (자동)")
    log("=" * 60)
    for _, name in NOTEBOOKS:
        log(f"  - {name}")
    log("=" * 60)

    DOWNLOAD_DIR.mkdir(parents=True, exist_ok=True)
    PERSISTENT_PROFILE.mkdir(parents=True, exist_ok=True)

    results = []

    async with async_playwright() as p:
        log("\n[시작] 브라우저 열기...")

        context = await p.chromium.launch_persistent_context(
            user_data_dir=str(PERSISTENT_PROFILE),
            headless=False,
            downloads_path=str(DOWNLOAD_DIR),
            accept_downloads=True,
            args=['--disable-blink-features=AutomationControlled'],
            viewport={'width': 1920, 'height': 1080},
        )

        page = context.pages[0] if context.pages else await context.new_page()

        await page.goto("https://notebooklm.google.com/", wait_until='networkidle', timeout=30000)
        await asyncio.sleep(2)

        if "accounts.google.com" in page.url:
            log("  로그인 필요...")
            for i in range(60):
                await asyncio.sleep(5)
                if "notebooklm.google.com" in page.url:
                    break
            else:
                await context.close()
                return []

        log("  ✓ 준비 완료")

        # 각 노트북 처리
        for idx, (keyword, display_name) in enumerate(NOTEBOOKS, 1):
            result = await process_notebook(page, keyword, display_name, idx, len(NOTEBOOKS))
            results.append(result)

        await context.close()

    # 결과
    log("\n" + "=" * 60)
    log("완료 결과:")
    log("=" * 60)

    success_count = 0
    for r in results:
        status = "✓" if r.get('pptx') else "❌"
        log(f"  {status} {r['name']}")
        if r.get('pptx'):
            log(f"      PPTX: {Path(r['pptx']).name} ({r['slides']}장)")
            success_count += 1

    log(f"\n성공: {success_count}/{len(NOTEBOOKS)}")
    log(f"저장 위치: {DOWNLOAD_DIR}")
    log("=" * 60)

    return results


if __name__ == "__main__":
    asyncio.run(main())
