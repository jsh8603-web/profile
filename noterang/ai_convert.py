#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
AI 기반 PDF → PPTX 변환 모듈
Grok Vision API를 사용하여 PDF 슬라이드 구조를 분석하고
수정 가능한 PPTX로 변환
"""
import io
import os
import sys
import json
import base64
from pathlib import Path
from typing import List, Dict, Optional, Tuple
from dataclasses import dataclass, field

if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8')

# 환경 변수 로드
from dotenv import load_dotenv
load_dotenv(Path(__file__).parent.parent / '.env.local')


@dataclass
class SlideElement:
    """슬라이드 요소"""
    type: str  # "title", "subtitle", "body", "bullet", "image", "chart", "table"
    content: str
    x: float = 0.0  # 0-1 비율
    y: float = 0.0
    width: float = 1.0
    height: float = 0.2
    font_size: int = 24
    bold: bool = False
    color: str = "#333333"
    alignment: str = "left"  # "left", "center", "right"


@dataclass
class SlideStructure:
    """슬라이드 구조"""
    page_number: int
    elements: List[SlideElement] = field(default_factory=list)
    background_color: str = "#FFFFFF"
    layout_type: str = "title_and_content"  # "title", "title_and_content", "two_column", "blank"


class GrokVisionAnalyzer:
    """Grok Vision API를 사용한 슬라이드 분석"""

    def __init__(self, api_key: str = None):
        self.api_key = api_key or os.getenv('GROK_API_KEY')
        if not self.api_key:
            raise ValueError("GROK_API_KEY가 필요합니다")

        self.base_url = "https://api.x.ai/v1"
        self.model = "grok-2-vision-latest"

    def analyze_slide_image(self, image_data: bytes) -> SlideStructure:
        """
        이미지를 분석하여 슬라이드 구조 추출

        Args:
            image_data: PNG/JPEG 이미지 바이트

        Returns:
            SlideStructure 객체
        """
        import requests

        # 이미지를 base64로 인코딩
        image_base64 = base64.b64encode(image_data).decode('utf-8')

        prompt = """이 슬라이드 이미지를 분석하여 다음 JSON 형식으로 구조를 추출해주세요.
모든 텍스트는 원본 그대로 추출하고, 위치는 0-1 비율로 표시합니다.

JSON 형식:
{
    "layout_type": "title_and_content",
    "background_color": "#FFFFFF",
    "elements": [
        {
            "type": "title",
            "content": "제목 텍스트",
            "x": 0.05,
            "y": 0.05,
            "width": 0.9,
            "height": 0.15,
            "font_size": 44,
            "bold": true,
            "color": "#333333",
            "alignment": "left"
        },
        {
            "type": "bullet",
            "content": "• 첫 번째 항목\\n• 두 번째 항목",
            "x": 0.05,
            "y": 0.25,
            "width": 0.9,
            "height": 0.6,
            "font_size": 24,
            "bold": false,
            "color": "#666666",
            "alignment": "left"
        }
    ]
}

type 종류: title, subtitle, body, bullet, image_placeholder, chart_placeholder, table_placeholder

중요:
- 모든 텍스트를 정확히 추출
- 불릿 포인트는 bullet 타입으로, 줄바꿈으로 구분
- 이미지/차트가 있으면 placeholder로 표시하고 content에 설명
- 색상은 hex 코드로
- 위치(x, y, width, height)는 0-1 비율로

JSON만 반환하세요:"""

        headers = {
            "Authorization": f"Bearer {self.api_key}",
            "Content-Type": "application/json"
        }

        payload = {
            "model": self.model,
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/png;base64,{image_base64}"
                            }
                        },
                        {
                            "type": "text",
                            "text": prompt
                        }
                    ]
                }
            ],
            "max_tokens": 8192,
            "temperature": 0.1
        }

        try:
            response = requests.post(
                f"{self.base_url}/chat/completions",
                headers=headers,
                json=payload,
                timeout=60
            )
            response.raise_for_status()

            result = response.json()
            text = result['choices'][0]['message']['content'].strip()

            # ```json ... ``` 형식 처리
            if '```json' in text:
                text = text.split('```json')[1].split('```')[0].strip()
            elif '```' in text:
                text = text.split('```')[1].split('```')[0].strip()

            data = json.loads(text)

            # SlideStructure로 변환
            elements = []
            for elem in data.get('elements', []):
                elements.append(SlideElement(
                    type=elem.get('type', 'body'),
                    content=elem.get('content', ''),
                    x=elem.get('x', 0.05),
                    y=elem.get('y', 0.1),
                    width=elem.get('width', 0.9),
                    height=elem.get('height', 0.2),
                    font_size=elem.get('font_size', 24),
                    bold=elem.get('bold', False),
                    color=elem.get('color', '#333333'),
                    alignment=elem.get('alignment', 'left')
                ))

            return SlideStructure(
                page_number=0,
                elements=elements,
                background_color=data.get('background_color', '#FFFFFF'),
                layout_type=data.get('layout_type', 'title_and_content')
            )

        except json.JSONDecodeError as e:
            print(f"  ⚠️ JSON 파싱 실패: {e}")
            print(f"  응답: {text[:500] if 'text' in dir() else 'N/A'}")
            return SlideStructure(page_number=0, elements=[])
        except requests.exceptions.RequestException as e:
            print(f"  ⚠️ Grok API 오류: {e}")
            return SlideStructure(page_number=0, elements=[])
        except Exception as e:
            print(f"  ⚠️ 오류: {e}")
            return SlideStructure(page_number=0, elements=[])


class AIConverter:
    """AI 기반 PDF → PPTX 변환기"""

    def __init__(self, api_key: str = None, provider: str = "grok"):
        """
        Args:
            api_key: API 키 (None이면 환경 변수 사용)
            provider: "grok" 또는 "gemini"
        """
        self.provider = provider
        if provider == "grok":
            self.analyzer = GrokVisionAnalyzer(api_key)
        else:
            raise ValueError(f"지원하지 않는 provider: {provider}")

    def pdf_to_images(self, pdf_path: Path, zoom: float = 2.0) -> List[bytes]:
        """PDF를 이미지 리스트로 변환"""
        import fitz

        doc = fitz.open(pdf_path)
        images = []

        for page_num in range(len(doc)):
            page = doc[page_num]
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat)
            img_data = pix.tobytes("png")
            images.append(img_data)

        doc.close()
        return images

    def analyze_pdf(self, pdf_path: Path, images: List[bytes] = None, progress_callback=None) -> List[SlideStructure]:
        """PDF 전체 분석"""
        print(f"📄 PDF 분석 중: {pdf_path.name}")

        if images is None:
            images = self.pdf_to_images(pdf_path)

        structures = []

        for i, img_data in enumerate(images):
            if progress_callback:
                progress_callback(i + 1, len(images))
            print(f"  🔍 페이지 {i+1}/{len(images)} 분석 중...")

            structure = self.analyzer.analyze_slide_image(img_data)
            structure.page_number = i + 1
            structures.append(structure)

        return structures

    def create_pptx(
        self,
        structures: List[SlideStructure],
        output_path: Path,
        images: List[bytes] = None,
        slide_width_inches: float = 13.333,
        slide_height_inches: float = 7.5,
        overlay_mode: bool = True
    ) -> Path:
        """
        분석된 구조로 PPTX 생성

        Args:
            structures: 슬라이드 구조 리스트
            output_path: 출력 경로
            images: PDF 페이지 이미지 리스트 (overlay_mode=True일 때 사용)
            slide_width_inches: 슬라이드 너비
            slide_height_inches: 슬라이드 높이
            overlay_mode: True면 이미지 배경 + 텍스트 오버레이
        """
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width = Inches(slide_width_inches)
        prs.slide_height = Inches(slide_height_inches)

        blank_layout = prs.slide_layouts[6]

        for i, struct in enumerate(structures):
            slide = prs.slides.add_slide(blank_layout)

            # 오버레이 모드: PDF 이미지를 배경으로
            if overlay_mode and images and i < len(images):
                slide.shapes.add_picture(
                    io.BytesIO(images[i]),
                    Inches(0), Inches(0),
                    width=prs.slide_width,
                    height=prs.slide_height
                )
            else:
                # 배경색 설정
                try:
                    bg_color = struct.background_color.lstrip('#')
                    r, g, b = int(bg_color[0:2], 16), int(bg_color[2:4], 16), int(bg_color[4:6], 16)
                    background = slide.background
                    fill = background.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(r, g, b)
                except:
                    pass

            # 텍스트 요소만 추가 (이미지/차트는 이미 배경에 있음)
            for elem in struct.elements:
                if elem.type not in ['image_placeholder', 'chart_placeholder', 'table_placeholder']:
                    self._add_textbox(slide, elem, prs, transparent=(overlay_mode and images))

        prs.save(output_path)
        print(f"  ✅ PPTX 저장: {output_path}")
        return output_path

    def _add_textbox(self, slide, elem: SlideElement, prs, transparent: bool = False):
        """텍스트 박스 추가"""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        # 위치 계산 (비율 → 인치)
        left = Inches(elem.x * 13.333)
        top = Inches(elem.y * 7.5)
        width = Inches(elem.width * 13.333)
        height = Inches(elem.height * 7.5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True

        # 투명 모드에서는 배경 투명하게 (기본적으로 텍스트박스는 투명)
        # 텍스트 박스 자체는 투명하므로 별도 처리 불필요

        # 텍스트 설정
        p = tf.paragraphs[0]
        p.text = elem.content
        p.font.size = Pt(elem.font_size)
        p.font.bold = elem.bold

        # 색상 - 투명 모드에서는 투명 텍스트 (선택 가능하지만 안 보임)
        if transparent:
            # 텍스트를 거의 투명하게 (배경 이미지가 보이도록)
            # 하지만 선택하면 편집 가능
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # 흰색
            # 투명도 설정은 python-pptx에서 직접 지원 안함
            # 대신 텍스트 색상을 배경과 비슷하게 설정
        else:
            try:
                color = elem.color.lstrip('#')
                r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
                p.font.color.rgb = RGBColor(r, g, b)
            except:
                pass

        # 정렬
        alignment_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT
        }
        p.alignment = alignment_map.get(elem.alignment, PP_ALIGN.LEFT)

    def _add_placeholder(self, slide, elem: SlideElement, prs):
        """플레이스홀더 추가 (이미지/차트 영역)"""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE

        left = Inches(elem.x * 13.333)
        top = Inches(elem.y * 7.5)
        width = Inches(elem.width * 13.333)
        height = Inches(elem.height * 7.5)

        # 점선 박스로 표시
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
        shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        shape.line.dash_style = 2  # 점선

        # 설명 텍스트 추가
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"[{elem.type.replace('_', ' ').title()}]\n{elem.content}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    def convert(
        self,
        pdf_path: Path,
        output_path: Path = None,
        progress_callback=None,
        overlay_mode: bool = True
    ) -> Tuple[Path, int]:
        """
        PDF를 수정 가능한 PPTX로 변환

        Args:
            pdf_path: PDF 파일 경로
            output_path: 출력 경로 (None이면 자동 생성)
            progress_callback: 진행 콜백 함수 (current, total)
            overlay_mode: True면 이미지 배경 + 텍스트 오버레이

        Returns:
            (출력 파일 경로, 슬라이드 수)
        """
        pdf_path = Path(pdf_path)

        if output_path is None:
            output_path = pdf_path.with_suffix('.pptx')
        else:
            output_path = Path(output_path)

        # PDF를 이미지로 변환
        print(f"📄 PDF 이미지 변환 중...")
        images = self.pdf_to_images(pdf_path)

        # PDF 분석 (이미지 재사용)
        structures = self.analyze_pdf(pdf_path, images=images, progress_callback=progress_callback)

        # PPTX 생성 (이미지 + 텍스트 오버레이)
        self.create_pptx(structures, output_path, images=images, overlay_mode=overlay_mode)

        return output_path, len(structures)


def convert_pdf_to_editable_pptx(
    pdf_path: str,
    output_path: str = None,
    api_key: str = None,
    provider: str = "grok"
) -> Tuple[Path, int]:
    """
    PDF를 수정 가능한 PPTX로 변환 (간편 함수)

    Args:
        pdf_path: PDF 파일 경로
        output_path: 출력 경로 (선택)
        api_key: API 키 (선택, 환경 변수 사용 가능)
        provider: "grok" (기본)

    Returns:
        (출력 파일 경로, 슬라이드 수)

    Example:
        >>> pptx_path, count = convert_pdf_to_editable_pptx("slides.pdf")
        >>> print(f"변환 완료: {pptx_path} ({count}장)")
    """
    converter = AIConverter(api_key, provider)
    return converter.convert(Path(pdf_path), Path(output_path) if output_path else None)


# CLI 지원
if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="AI 기반 PDF → PPTX 변환 (Grok Vision)")
    parser.add_argument("pdf_path", help="PDF 파일 경로")
    parser.add_argument("-o", "--output", help="출력 경로")
    parser.add_argument("--api-key", help="Grok API 키")
    parser.add_argument("--no-overlay", action="store_true", help="배경 이미지 없이 텍스트만")

    args = parser.parse_args()

    try:
        converter = AIConverter(args.api_key)
        pptx_path, count = converter.convert(
            Path(args.pdf_path),
            Path(args.output) if args.output else None,
            overlay_mode=not args.no_overlay
        )
        print(f"\n✅ 변환 완료!")
        print(f"   파일: {pptx_path}")
        print(f"   슬라이드: {count}장")
    except Exception as e:
        print(f"❌ 오류: {e}")
        sys.exit(1)
