#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
노트랑 PDF → PPTX 변환 모듈

PyMuPDF로 PDF 페이지를 이미지로 추출하고
python-pptx로 PPTX 슬라이드에 삽입합니다.
"""
import sys
from pathlib import Path
from typing import Optional

if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8')


def pdf_to_pptx(
    pdf_path: str,
    pptx_path: str = None,
    dpi: int = 200,
) -> bool:
    """
    PDF를 PPTX로 변환

    Args:
        pdf_path: 입력 PDF 파일 경로
        pptx_path: 출력 PPTX 파일 경로 (None이면 PDF와 같은 이름)
        dpi: 이미지 해상도 (기본 200)

    Returns:
        성공 여부
    """
    try:
        import fitz  # PyMuPDF
        from pptx import Presentation
        from pptx.util import Inches
        from io import BytesIO
    except ImportError as e:
        print(f"필요한 패키지가 설치되지 않았습니다: {e}")
        print("설치: pip install pymupdf python-pptx")
        return False

    pdf_path = Path(pdf_path)
    if not pdf_path.exists():
        print(f"PDF 파일을 찾을 수 없습니다: {pdf_path}")
        return False

    if pptx_path is None:
        pptx_path = pdf_path.with_suffix('.pptx')
    else:
        pptx_path = Path(pptx_path)

    try:
        # PDF 열기
        pdf_doc = fitz.open(str(pdf_path))
        page_count = len(pdf_doc)

        print(f"PDF 변환 중: {pdf_path.name} ({page_count}페이지)")

        # PPTX 생성 (16:9 비율)
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]  # 빈 슬라이드

        for page_num in range(page_count):
            page = pdf_doc[page_num]

            # 페이지를 이미지로 렌더링
            mat = fitz.Matrix(dpi / 72, dpi / 72)
            pix = page.get_pixmap(matrix=mat)

            # 이미지를 BytesIO로 변환
            img_bytes = BytesIO(pix.tobytes("png"))

            # 슬라이드 추가
            slide = prs.slides.add_slide(blank_layout)

            # 이미지 삽입 (전체 슬라이드 크기)
            slide.shapes.add_picture(
                img_bytes,
                Inches(0),
                Inches(0),
                width=prs.slide_width,
                height=prs.slide_height,
            )

            print(f"  페이지 {page_num + 1}/{page_count} 변환 완료")

        pdf_doc.close()

        # PPTX 저장
        prs.save(str(pptx_path))
        print(f"✓ PPTX 저장 완료: {pptx_path}")

        return True

    except Exception as e:
        print(f"변환 실패: {e}")
        return False


def batch_convert(
    pdf_dir: str,
    output_dir: str = None,
    dpi: int = 200,
) -> int:
    """
    디렉토리 내 모든 PDF를 PPTX로 변환

    Args:
        pdf_dir: PDF 파일이 있는 디렉토리
        output_dir: 출력 디렉토리 (None이면 같은 디렉토리)
        dpi: 이미지 해상도

    Returns:
        변환 성공한 파일 수
    """
    pdf_dir = Path(pdf_dir)
    if not pdf_dir.exists():
        print(f"디렉토리를 찾을 수 없습니다: {pdf_dir}")
        return 0

    if output_dir:
        output_dir = Path(output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)
    else:
        output_dir = pdf_dir

    pdf_files = list(pdf_dir.glob("*.pdf"))
    if not pdf_files:
        print(f"PDF 파일이 없습니다: {pdf_dir}")
        return 0

    print(f"\n{len(pdf_files)}개 PDF 파일 변환 시작\n")

    success_count = 0
    for pdf_file in pdf_files:
        pptx_path = output_dir / pdf_file.with_suffix('.pptx').name

        if pdf_to_pptx(str(pdf_file), str(pptx_path), dpi):
            success_count += 1
        print()

    print(f"\n변환 완료: {success_count}/{len(pdf_files)} 파일")
    return success_count


# CLI 실행
if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="PDF → PPTX 변환")
    parser.add_argument("pdf_path", help="PDF 파일 또는 디렉토리 경로")
    parser.add_argument("-o", "--output", help="출력 경로")
    parser.add_argument("--dpi", type=int, default=200, help="이미지 해상도 (기본 200)")
    parser.add_argument("--batch", action="store_true", help="디렉토리 일괄 변환")

    args = parser.parse_args()

    if args.batch:
        batch_convert(args.pdf_path, args.output, args.dpi)
    else:
        success = pdf_to_pptx(args.pdf_path, args.output, args.dpi)
        sys.exit(0 if success else 1)
