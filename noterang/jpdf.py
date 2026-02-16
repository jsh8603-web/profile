#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
JPDF - PDF → 편집 가능 PPTX 변환기

Google Cloud Vision OCR + OpenCV Inpainting을 사용하여
PDF 슬라이드에서 텍스트를 추출하고, 배경을 복원한 후
편집 가능한 텍스트 박스로 PPTX를 생성합니다.

Usage:
    python -m noterang.jpdf input.pdf -o output.pptx
    python -m noterang.jpdf input.pdf --no-inpaint

    # 또는 직접 실행
    jpdf input.pdf -o output.pptx
"""
import io
import os
import sys
import base64
import requests
import cv2
import numpy as np
from pathlib import Path
from typing import List, Tuple, Optional
from dataclasses import dataclass, field

if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8')

from dotenv import load_dotenv
load_dotenv(Path(__file__).parent.parent / '.env.local')


@dataclass
class TextBlock:
    """텍스트 블록"""
    text: str
    x: int
    y: int
    width: int
    height: int
    font_size: float = 24.0  # 추정 폰트 크기 (pt)
    line_height: float = 1.2  # 줄 높이 비율
    alignment: str = "left"  # left, center, right


@dataclass
class PageData:
    """페이지 데이터"""
    page_num: int
    original_image: np.ndarray
    cleaned_image: np.ndarray
    text_blocks: List[TextBlock]
    width: int
    height: int


class JPDF:
    """JPDF - PDF → 편집 가능 PPTX 변환기"""

    def __init__(self, api_key: str = None):
        """
        Args:
            api_key: Google Cloud Vision API 키
        """
        self.api_key = api_key or os.getenv('GOOGLE_VISION_API_KEY')
        if not self.api_key:
            raise ValueError(
                "GOOGLE_VISION_API_KEY가 필요합니다.\n"
                ".env.local 파일에 설정하거나 --api-key 옵션을 사용하세요."
            )
        self.api_url = "https://vision.googleapis.com/v1/images:annotate"

    def _ocr_image(self, image: np.ndarray) -> List[TextBlock]:
        """이미지에서 텍스트 블록 추출 (폰트 크기, 정렬 포함)"""
        img_height, img_width = image.shape[:2]

        # BGR → RGB → PNG
        rgb = cv2.cvtColor(image, cv2.COLOR_BGR2RGB)
        _, buffer = cv2.imencode('.png', rgb)
        img_base64 = base64.b64encode(buffer).decode('utf-8')

        payload = {
            "requests": [{
                "image": {"content": img_base64},
                "features": [{"type": "DOCUMENT_TEXT_DETECTION"}]
            }]
        }

        response = requests.post(
            f"{self.api_url}?key={self.api_key}",
            json=payload,
            timeout=60
        )
        result = response.json()

        if 'error' in result:
            raise Exception(f"Vision API 오류: {result['error']}")

        blocks = []
        responses = result.get('responses', [{}])
        if not responses:
            return blocks

        full_annotation = responses[0].get('fullTextAnnotation', {})

        for page in full_annotation.get('pages', []):
            for block in page.get('blocks', []):
                # 블록 텍스트 및 폰트 정보 추출
                block_text = ""
                word_heights = []  # 단어 높이들 (폰트 크기 추정용)
                line_xs = []  # 각 줄의 x 좌표들 (정렬 감지용)

                for para in block.get('paragraphs', []):
                    para_text = ""
                    para_word_xs = []

                    for word in para.get('words', []):
                        word_text = "".join(
                            s.get('text', '') for s in word.get('symbols', [])
                        )
                        para_text += word_text + " "

                        # 단어 높이 추출 (폰트 크기 추정)
                        word_vertices = word.get('boundingBox', {}).get('vertices', [])
                        if len(word_vertices) >= 4:
                            word_ys = [v.get('y', 0) for v in word_vertices]
                            word_xs = [v.get('x', 0) for v in word_vertices]
                            word_height = max(word_ys) - min(word_ys)
                            if word_height > 0:
                                word_heights.append(word_height)
                            para_word_xs.append(min(word_xs))

                    block_text += para_text.strip() + "\n"
                    if para_word_xs:
                        line_xs.append(min(para_word_xs))

                block_text = block_text.strip()
                if not block_text:
                    continue

                # 바운딩 박스
                vertices = block.get('boundingBox', {}).get('vertices', [])
                if len(vertices) >= 4:
                    xs = [v.get('x', 0) for v in vertices]
                    ys = [v.get('y', 0) for v in vertices]
                    block_x = min(xs)
                    block_y = min(ys)
                    block_width = max(xs) - min(xs)
                    block_height = max(ys) - min(ys)

                    # 폰트 크기 추정 (단어 높이의 중간값 사용)
                    if word_heights:
                        median_height = sorted(word_heights)[len(word_heights) // 2]
                        # 픽셀 → 포인트 변환 (zoom 2.0 보정 + 1pt 추가)
                        font_size = (median_height / 2.0) * 0.5 + 1  # 원래 크기 + 1pt
                        font_size = max(8, min(48, font_size))  # 8~48pt 범위
                    else:
                        font_size = 12.0

                    # 정렬 감지 (블록 내 텍스트 위치 기반)
                    alignment = "left"
                    if line_xs and block_width > 0:
                        avg_line_x = sum(line_xs) / len(line_xs)
                        relative_pos = (avg_line_x - block_x) / block_width
                        if relative_pos > 0.35:
                            alignment = "center"
                        elif relative_pos > 0.6:
                            alignment = "right"

                    # 줄 높이 비율 계산
                    line_count = block_text.count('\n') + 1
                    if line_count > 1 and font_size > 0:
                        line_height = (block_height / line_count) / (font_size * 1.33)
                        line_height = max(1.0, min(2.0, line_height))
                    else:
                        line_height = 1.2

                    blocks.append(TextBlock(
                        text=block_text,
                        x=block_x,
                        y=block_y,
                        width=block_width,
                        height=block_height,
                        font_size=font_size,
                        line_height=line_height,
                        alignment=alignment
                    ))

        # 가까운 블록들을 하나로 병합
        merged_blocks = self._merge_nearby_blocks(blocks)
        return merged_blocks

    def _merge_nearby_blocks(
        self,
        blocks: List[TextBlock],
        y_threshold: int = 50,  # Y 방향 거리 임계값 (픽셀)
        x_threshold: int = 100  # X 방향 거리 임계값 (픽셀)
    ) -> List[TextBlock]:
        """가까운 블록들을 하나로 병합"""
        if not blocks:
            return blocks

        # Y 좌표로 정렬
        sorted_blocks = sorted(blocks, key=lambda b: (b.y, b.x))
        merged = []
        current_group = [sorted_blocks[0]]

        for block in sorted_blocks[1:]:
            last = current_group[-1]

            # 이전 블록과 Y 방향으로 가깝고, X 범위가 겹치면 같은 그룹
            y_close = (block.y - (last.y + last.height)) < y_threshold
            x_overlap = not (block.x > last.x + last.width + x_threshold or
                           block.x + block.width < last.x - x_threshold)

            if y_close and x_overlap:
                current_group.append(block)
            else:
                # 그룹 병합 후 새 그룹 시작
                merged.append(self._merge_block_group(current_group))
                current_group = [block]

        # 마지막 그룹 병합
        if current_group:
            merged.append(self._merge_block_group(current_group))

        return merged

    def _merge_block_group(self, group: List[TextBlock]) -> TextBlock:
        """블록 그룹을 하나의 블록으로 병합"""
        if len(group) == 1:
            return group[0]

        # 모든 블록의 텍스트 합치기
        texts = [b.text for b in group]
        merged_text = "\n".join(texts)

        # 바운딩 박스 계산 (모든 블록을 포함)
        min_x = min(b.x for b in group)
        min_y = min(b.y for b in group)
        max_x = max(b.x + b.width for b in group)
        max_y = max(b.y + b.height for b in group)

        # 폰트 크기는 평균값
        avg_font_size = sum(b.font_size for b in group) / len(group)

        # 정렬은 첫 블록 기준
        alignment = group[0].alignment

        # 줄 높이는 평균값
        avg_line_height = sum(b.line_height for b in group) / len(group)

        return TextBlock(
            text=merged_text,
            x=min_x,
            y=min_y,
            width=max_x - min_x,
            height=max_y - min_y,
            font_size=avg_font_size,
            line_height=avg_line_height,
            alignment=alignment
        )

    def _inpaint_image(
        self,
        image: np.ndarray,
        blocks: List[TextBlock],
        padding: int = 10,
        radius: int = 7
    ) -> np.ndarray:
        """텍스트 영역 제거 및 배경 복원"""
        if not blocks:
            return image.copy()

        # 마스크 생성
        mask = np.zeros(image.shape[:2], dtype=np.uint8)
        for b in blocks:
            x1 = max(0, b.x - padding)
            y1 = max(0, b.y - padding)
            x2 = min(image.shape[1], b.x + b.width + padding)
            y2 = min(image.shape[0], b.y + b.height + padding)
            cv2.rectangle(mask, (x1, y1), (x2, y2), 255, -1)

        # Inpainting
        cleaned = cv2.inpaint(image, mask, radius, cv2.INPAINT_TELEA)
        return cleaned

    def process_pdf(
        self,
        pdf_path: str,
        zoom: float = 2.0,
        inpaint: bool = True,
        padding: int = 10,
        inpaint_radius: int = 7
    ) -> List[PageData]:
        """
        PDF 처리

        Args:
            pdf_path: PDF 파일 경로
            zoom: 이미지 확대 비율
            inpaint: 텍스트 제거 여부
            padding: 텍스트 영역 패딩
            inpaint_radius: Inpainting 반경

        Returns:
            페이지 데이터 리스트
        """
        import fitz

        pdf_path = Path(pdf_path)
        if not pdf_path.exists():
            raise FileNotFoundError(f"파일을 찾을 수 없습니다: {pdf_path}")

        doc = fitz.open(str(pdf_path))
        pages = []

        total = len(doc)
        print(f"📄 PDF: {pdf_path.name} ({total}페이지)")

        for i in range(total):
            print(f"  [{i+1}/{total}] OCR + 처리 중...", end=" ", flush=True)

            page = doc[i]
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat)

            # numpy array로 변환
            img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(
                pix.height, pix.width, pix.n
            )
            if pix.n == 4:
                img = cv2.cvtColor(img, cv2.COLOR_RGBA2BGR)
            else:
                img = cv2.cvtColor(img, cv2.COLOR_RGB2BGR)

            # OCR
            blocks = self._ocr_image(img)

            # Inpainting
            if inpaint and blocks:
                cleaned = self._inpaint_image(img, blocks, padding, inpaint_radius)
            else:
                cleaned = img.copy()

            pages.append(PageData(
                page_num=i + 1,
                original_image=img,
                cleaned_image=cleaned,
                text_blocks=blocks,
                width=pix.width,
                height=pix.height
            ))

            print(f"텍스트 {len(blocks)}개")

        doc.close()
        return pages

    def create_pptx(
        self,
        pages: List[PageData],
        output_path: str,
        font_size: int = None,
        font_color: Tuple[int, int, int] = (0x33, 0x33, 0x33)
    ) -> Path:
        """
        PPTX 생성

        Args:
            pages: 페이지 데이터 리스트
            output_path: 출력 경로
            font_size: 폰트 크기 (None이면 자동)
            font_color: 폰트 색상 (R, G, B)

        Returns:
            출력 파일 경로
        """
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]

        for page_idx, page in enumerate(pages):
            slide = prs.slides.add_slide(blank_layout)

            # 좌표 변환 비율
            scale_x = prs.slide_width.emu / page.width
            scale_y = prs.slide_height.emu / page.height

            # 1. 배경 이미지 (텍스트 제거됨)
            _, buffer = cv2.imencode('.png', page.cleaned_image)
            slide.shapes.add_picture(
                io.BytesIO(buffer.tobytes()),
                Inches(0), Inches(0),
                width=prs.slide_width,
                height=prs.slide_height
            )

            # 2. 텍스트 박스 추가
            for block_idx, block in enumerate(page.text_blocks):
                left = Emu(int(block.x * scale_x))
                top = Emu(int(block.y * scale_y))
                width = Emu(int(block.width * scale_x))
                height = Emu(int(block.height * scale_y))

                textbox = slide.shapes.add_textbox(left, top, width, height)
                tf = textbox.text_frame
                tf.word_wrap = True

                # 텍스트 프레임 설정
                from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
                tf.auto_size = None  # 자동 크기 조정 비활성화

                # 정렬 설정
                alignment_map = {
                    'left': PP_ALIGN.LEFT,
                    'center': PP_ALIGN.CENTER,
                    'right': PP_ALIGN.RIGHT
                }

                # 줄바꿈된 텍스트 처리
                lines = block.text.split('\n')
                for i, line in enumerate(lines):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()

                    p.text = line
                    p.alignment = alignment_map.get(block.alignment, PP_ALIGN.LEFT)

                    # 폰트 크기 결정
                    # 첫 슬라이드의 첫 블록(제목)은 24pt 볼드
                    is_first_slide_title = (page_idx == 0 and block_idx == 0)

                    if font_size:
                        size = font_size
                    elif is_first_slide_title:
                        size = 24  # 첫 슬라이드 제목은 24pt 고정
                    else:
                        size = int(block.font_size)

                    p.font.size = Pt(size)
                    p.font.color.rgb = RGBColor(*font_color)
                    p.font.bold = is_first_slide_title  # 첫 슬라이드 제목만 볼드

                    # 줄 간격 설정
                    p.line_spacing = block.line_height

        output_path = Path(output_path)
        prs.save(str(output_path))
        return output_path

    def convert(
        self,
        pdf_path: str,
        output_path: str = None,
        inpaint: bool = True,
        zoom: float = 2.0,
        padding: int = 10,
        inpaint_radius: int = 7,
        font_size: int = None
    ) -> Tuple[Path, int]:
        """
        PDF → 편집 가능 PPTX 변환

        Args:
            pdf_path: PDF 파일 경로
            output_path: 출력 경로 (None이면 자동)
            inpaint: 텍스트 제거 여부
            zoom: 이미지 확대 비율
            padding: 텍스트 영역 패딩
            inpaint_radius: Inpainting 반경
            font_size: 폰트 크기 (None이면 자동)

        Returns:
            (출력 파일 경로, 페이지 수)
        """
        pdf_path = Path(pdf_path)

        if output_path is None:
            output_path = pdf_path.with_name(f"{pdf_path.stem}_편집가능.pptx")

        # PDF 처리
        pages = self.process_pdf(
            pdf_path,
            zoom=zoom,
            inpaint=inpaint,
            padding=padding,
            inpaint_radius=inpaint_radius
        )

        # PPTX 생성
        result_path = self.create_pptx(pages, output_path, font_size)

        print(f"\n✅ 변환 완료!")
        print(f"   파일: {result_path}")
        print(f"   슬라이드: {len(pages)}장")

        return result_path, len(pages)


def convert(
    pdf_path: str,
    output_path: str = None,
    api_key: str = None,
    **kwargs
) -> Tuple[Path, int]:
    """
    PDF → 편집 가능 PPTX 변환 (간편 함수)

    Args:
        pdf_path: PDF 파일 경로
        output_path: 출력 경로
        api_key: Google Vision API 키
        **kwargs: 추가 옵션

    Returns:
        (출력 파일 경로, 페이지 수)

    Example:
        >>> from noterang.pdf2pptx import convert
        >>> pptx, count = convert("slides.pdf")
        >>> print(f"변환 완료: {pptx}")
    """
    converter = JPDF(api_key)
    return converter.convert(pdf_path, output_path, **kwargs)


def main():
    """CLI 엔트리포인트"""
    import argparse

    parser = argparse.ArgumentParser(
        prog="jpdf",
        description="JPDF - PDF → 편집 가능 PPTX 변환기",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
예시:
  jpdf slides.pdf
  jpdf slides.pdf -o output.pptx
  jpdf slides.pdf --no-inpaint
  jpdf slides.pdf --font-size 20

  # 또는 모듈로 실행
  python -m noterang.jpdf slides.pdf
        """
    )
    parser.add_argument("pdf_path", help="PDF 파일 경로")
    parser.add_argument("-o", "--output", help="출력 PPTX 경로")
    parser.add_argument("--api-key", help="Google Vision API 키")
    parser.add_argument(
        "--no-inpaint",
        action="store_true",
        help="텍스트 제거 없이 원본 배경 사용"
    )
    parser.add_argument(
        "--zoom",
        type=float,
        default=2.0,
        help="이미지 확대 비율 (기본: 2.0)"
    )
    parser.add_argument(
        "--padding",
        type=int,
        default=10,
        help="텍스트 영역 패딩 (기본: 10)"
    )
    parser.add_argument(
        "--inpaint-radius",
        type=int,
        default=7,
        help="Inpainting 반경 (기본: 7)"
    )
    parser.add_argument(
        "--font-size",
        type=int,
        help="폰트 크기 (기본: 자동)"
    )

    args = parser.parse_args()

    try:
        converter = JPDF(args.api_key)
        converter.convert(
            args.pdf_path,
            args.output,
            inpaint=not args.no_inpaint,
            zoom=args.zoom,
            padding=args.padding,
            inpaint_radius=args.inpaint_radius,
            font_size=args.font_size
        )
    except FileNotFoundError as e:
        print(f"❌ 오류: {e}")
        sys.exit(1)
    except Exception as e:
        print(f"❌ 오류: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
