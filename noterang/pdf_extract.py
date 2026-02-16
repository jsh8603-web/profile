#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PDF 구성 요소 분리 추출 → 편집 가능한 PPTX 생성
- 이미지/도형: 그대로 추출하여 PPTX에 삽입
- 텍스트: 편집 가능한 텍스트 박스로 분리
"""
import io
import os
import sys
from pathlib import Path
from typing import List, Dict, Tuple, Optional
from dataclasses import dataclass, field

if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8')


@dataclass
class TextBlock:
    """텍스트 블록"""
    text: str
    x: float  # 좌표 (pt)
    y: float
    width: float
    height: float
    font_size: float
    font_name: str
    color: Tuple[float, float, float]  # RGB 0-1
    bold: bool = False
    italic: bool = False


@dataclass
class ImageBlock:
    """이미지 블록"""
    data: bytes
    x: float
    y: float
    width: float
    height: float
    ext: str = "png"


@dataclass
class PageContent:
    """페이지 컨텐츠"""
    page_num: int
    width: float  # pt
    height: float
    texts: List[TextBlock] = field(default_factory=list)
    images: List[ImageBlock] = field(default_factory=list)
    background_image: bytes = None  # 도형/그래픽용 배경


class PDFExtractor:
    """PDF 구성 요소 추출기"""

    def __init__(self):
        try:
            import fitz
            self.fitz = fitz
        except ImportError:
            raise ImportError("pip install pymupdf 필요")

    def extract_page(self, page, page_num: int) -> PageContent:
        """한 페이지에서 텍스트와 이미지 추출"""
        content = PageContent(
            page_num=page_num,
            width=page.rect.width,
            height=page.rect.height
        )

        # 1. 텍스트 추출 (위치 정보 포함)
        blocks = page.get_text("dict", flags=self.fitz.TEXT_PRESERVE_WHITESPACE)

        for block in blocks.get("blocks", []):
            if block.get("type") == 0:  # 텍스트 블록
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        text = span.get("text", "").strip()
                        if not text:
                            continue

                        bbox = span.get("bbox", [0, 0, 0, 0])
                        font_size = span.get("size", 12)
                        font_name = span.get("font", "")
                        color = span.get("color", 0)

                        # 색상 변환 (int → RGB)
                        if isinstance(color, int):
                            r = ((color >> 16) & 0xFF) / 255
                            g = ((color >> 8) & 0xFF) / 255
                            b = (color & 0xFF) / 255
                        else:
                            r, g, b = 0, 0, 0

                        # 볼드/이탤릭 감지
                        flags = span.get("flags", 0)
                        bold = bool(flags & 2**4)  # bit 4
                        italic = bool(flags & 2**1)  # bit 1

                        content.texts.append(TextBlock(
                            text=text,
                            x=bbox[0],
                            y=bbox[1],
                            width=bbox[2] - bbox[0],
                            height=bbox[3] - bbox[1],
                            font_size=font_size,
                            font_name=font_name,
                            color=(r, g, b),
                            bold=bold,
                            italic=italic
                        ))

        # 2. 이미지 추출
        image_list = page.get_images(full=True)

        for img_index, img_info in enumerate(image_list):
            xref = img_info[0]
            try:
                base_image = page.parent.extract_image(xref)
                if base_image:
                    image_data = base_image.get("image")
                    ext = base_image.get("ext", "png")

                    # 이미지 위치 찾기
                    img_rect = page.get_image_rects(xref)
                    if img_rect:
                        rect = img_rect[0]
                        content.images.append(ImageBlock(
                            data=image_data,
                            x=rect.x0,
                            y=rect.y0,
                            width=rect.width,
                            height=rect.height,
                            ext=ext
                        ))
            except Exception as e:
                print(f"  ⚠️ 이미지 {img_index} 추출 실패: {e}")

        # 3. 도형/그래픽이 있는 경우 배경 이미지로 렌더링
        # (벡터 도형은 직접 변환이 어려우므로 이미지로 처리)
        drawings = page.get_drawings()
        if drawings:
            # 텍스트 없이 도형만 렌더링
            # 텍스트를 가리고 렌더링
            mat = self.fitz.Matrix(2, 2)  # 2배 확대

            # 원본 페이지를 이미지로 렌더링 (텍스트 포함)
            pix = page.get_pixmap(matrix=mat)
            content.background_image = pix.tobytes("png")

        return content

    def extract_pdf(self, pdf_path: Path) -> List[PageContent]:
        """PDF 전체 추출"""
        print(f"📄 PDF 추출 중: {pdf_path.name}")

        doc = self.fitz.open(pdf_path)
        pages = []

        for page_num in range(len(doc)):
            print(f"  📃 페이지 {page_num + 1}/{len(doc)}...")
            page = doc[page_num]
            content = self.extract_page(page, page_num + 1)
            pages.append(content)

        doc.close()
        return pages


class PPTXBuilder:
    """PPTX 생성기"""

    def __init__(self):
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            self.Presentation = Presentation
            self.Inches = Inches
            self.Pt = Pt
            self.RGBColor = RGBColor
            self.PP_ALIGN = PP_ALIGN
        except ImportError:
            raise ImportError("pip install python-pptx 필요")

    def pt_to_inches(self, pt: float) -> float:
        """포인트를 인치로 변환"""
        return pt / 72.0

    def build(
        self,
        pages: List[PageContent],
        output_path: Path,
        use_background: bool = True
    ) -> Path:
        """PageContent 리스트로 PPTX 생성"""
        prs = self.Presentation()

        if pages:
            # 슬라이드 크기를 PDF 페이지 크기에 맞춤
            first_page = pages[0]
            prs.slide_width = self.Inches(self.pt_to_inches(first_page.width))
            prs.slide_height = self.Inches(self.pt_to_inches(first_page.height))

        blank_layout = prs.slide_layouts[6]

        for page in pages:
            slide = prs.slides.add_slide(blank_layout)

            # 1. 배경 이미지 (도형/그래픽 포함)
            if use_background and page.background_image:
                slide.shapes.add_picture(
                    io.BytesIO(page.background_image),
                    self.Inches(0),
                    self.Inches(0),
                    width=prs.slide_width,
                    height=prs.slide_height
                )

            # 2. 개별 이미지 추가 (배경 사용 안 할 때)
            if not use_background:
                for img in page.images:
                    try:
                        slide.shapes.add_picture(
                            io.BytesIO(img.data),
                            self.Inches(self.pt_to_inches(img.x)),
                            self.Inches(self.pt_to_inches(img.y)),
                            width=self.Inches(self.pt_to_inches(img.width)),
                            height=self.Inches(self.pt_to_inches(img.height))
                        )
                    except Exception as e:
                        print(f"  ⚠️ 이미지 삽입 실패: {e}")

            # 3. 텍스트 박스 추가 (항상 맨 위에)
            for text_block in page.texts:
                self._add_text(slide, text_block, page)

        prs.save(output_path)
        print(f"  ✅ PPTX 저장: {output_path}")
        return output_path

    def _add_text(self, slide, tb: TextBlock, page: PageContent):
        """텍스트 박스 추가"""
        left = self.Inches(self.pt_to_inches(tb.x))
        top = self.Inches(self.pt_to_inches(tb.y))
        width = self.Inches(self.pt_to_inches(tb.width + 10))  # 약간 여유
        height = self.Inches(self.pt_to_inches(tb.height + 5))

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = False  # 원본 줄바꿈 유지

        p = tf.paragraphs[0]
        p.text = tb.text
        p.font.size = self.Pt(tb.font_size)
        p.font.bold = tb.bold
        p.font.italic = tb.italic

        # 색상
        r = int(tb.color[0] * 255)
        g = int(tb.color[1] * 255)
        b = int(tb.color[2] * 255)
        p.font.color.rgb = self.RGBColor(r, g, b)


def convert_pdf_to_editable_pptx(
    pdf_path: str,
    output_path: str = None,
    use_background: bool = True
) -> Tuple[Path, int]:
    """
    PDF를 편집 가능한 PPTX로 변환

    Args:
        pdf_path: PDF 파일 경로
        output_path: 출력 경로 (선택)
        use_background: 배경 이미지 사용 여부

    Returns:
        (출력 파일 경로, 슬라이드 수)
    """
    pdf_path = Path(pdf_path)
    if output_path is None:
        output_path = pdf_path.with_name(pdf_path.stem + "_editable.pptx")
    else:
        output_path = Path(output_path)

    # PDF 추출
    extractor = PDFExtractor()
    pages = extractor.extract_pdf(pdf_path)

    # PPTX 생성
    builder = PPTXBuilder()
    builder.build(pages, output_path, use_background=use_background)

    return output_path, len(pages)


# CLI
if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="PDF → 편집 가능 PPTX 변환")
    parser.add_argument("pdf_path", help="PDF 파일 경로")
    parser.add_argument("-o", "--output", help="출력 경로")
    parser.add_argument("--no-background", action="store_true", help="배경 이미지 제외")

    args = parser.parse_args()

    try:
        pptx_path, count = convert_pdf_to_editable_pptx(
            args.pdf_path,
            args.output,
            use_background=not args.no_background
        )
        print(f"\n✅ 변환 완료!")
        print(f"   파일: {pptx_path}")
        print(f"   슬라이드: {count}장")
    except Exception as e:
        print(f"❌ 오류: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
