#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PDF 텍스트 제거 + Inpainting → 편집 가능 PPTX

1. Google Vision OCR로 텍스트 위치 감지
2. 텍스트 영역 마스크 생성
3. OpenCV inpaint로 배경 복원
4. 복원된 이미지 + 편집 가능한 텍스트 박스 → PPTX
"""
import io
import os
import sys
import base64
import requests
import cv2
import numpy as np
from pathlib import Path
from typing import List, Tuple, Optional
from dataclasses import dataclass, field

if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8')

from dotenv import load_dotenv
load_dotenv(Path(__file__).parent.parent / '.env.local')


@dataclass
class TextRegion:
    """텍스트 영역"""
    text: str
    x: int
    y: int
    width: int
    height: int
    vertices: List[Tuple[int, int]] = field(default_factory=list)


@dataclass
class ProcessedPage:
    """처리된 페이지"""
    page_num: int
    original_image: np.ndarray
    cleaned_image: np.ndarray  # 텍스트 제거된 이미지
    text_regions: List[TextRegion]
    width: int
    height: int


class TextRemover:
    """텍스트 제거 및 Inpainting"""

    def __init__(self, api_key: str = None):
        self.api_key = api_key or os.getenv('GOOGLE_VISION_API_KEY')
        if not self.api_key:
            raise ValueError("GOOGLE_VISION_API_KEY가 필요합니다")
        self.api_url = "https://vision.googleapis.com/v1/images:annotate"

    def detect_text_regions(self, image: np.ndarray) -> List[TextRegion]:
        """이미지에서 텍스트 영역 감지"""
        # BGR to RGB, then encode
        rgb = cv2.cvtColor(image, cv2.COLOR_BGR2RGB)
        _, buffer = cv2.imencode('.png', rgb)
        img_base64 = base64.b64encode(buffer).decode('utf-8')

        payload = {
            "requests": [{
                "image": {"content": img_base64},
                "features": [{"type": "DOCUMENT_TEXT_DETECTION"}]
            }]
        }

        response = requests.post(
            f"{self.api_url}?key={self.api_key}",
            json=payload,
            timeout=60
        )
        result = response.json()

        if 'error' in result:
            raise Exception(f"Vision API 오류: {result['error']}")

        regions = []
        responses = result.get('responses', [{}])
        if not responses:
            return regions

        resp = responses[0]
        full_annotation = resp.get('fullTextAnnotation', {})

        for page in full_annotation.get('pages', []):
            for block in page.get('blocks', []):
                # 블록 단위로 텍스트 추출
                block_text = ""
                for paragraph in block.get('paragraphs', []):
                    para_text = ""
                    for word in paragraph.get('words', []):
                        word_text = "".join(
                            symbol.get('text', '')
                            for symbol in word.get('symbols', [])
                        )
                        para_text += word_text + " "
                    block_text += para_text.strip() + "\n"

                block_text = block_text.strip()
                if not block_text:
                    continue

                # 바운딩 박스
                vertices = block.get('boundingBox', {}).get('vertices', [])
                if len(vertices) >= 4:
                    points = [(v.get('x', 0), v.get('y', 0)) for v in vertices]
                    xs = [p[0] for p in points]
                    ys = [p[1] for p in points]
                    x, y = min(xs), min(ys)
                    w, h = max(xs) - x, max(ys) - y

                    regions.append(TextRegion(
                        text=block_text,
                        x=x, y=y,
                        width=w, height=h,
                        vertices=points
                    ))

        return regions

    def create_text_mask(
        self,
        image_shape: Tuple[int, int],
        regions: List[TextRegion],
        padding: int = 5
    ) -> np.ndarray:
        """텍스트 영역 마스크 생성"""
        mask = np.zeros(image_shape[:2], dtype=np.uint8)

        for region in regions:
            # 패딩 적용
            x1 = max(0, region.x - padding)
            y1 = max(0, region.y - padding)
            x2 = min(image_shape[1], region.x + region.width + padding)
            y2 = min(image_shape[0], region.y + region.height + padding)

            # 폴리곤 또는 사각형으로 마스크
            if region.vertices and len(region.vertices) >= 4:
                pts = np.array(region.vertices, dtype=np.int32)
                # 패딩 적용된 폴리곤
                center = pts.mean(axis=0)
                pts_padded = pts + (pts - center) * (padding / 50)
                pts_padded = pts_padded.astype(np.int32)
                cv2.fillPoly(mask, [pts_padded], 255)
            else:
                cv2.rectangle(mask, (x1, y1), (x2, y2), 255, -1)

        return mask

    def inpaint_image(
        self,
        image: np.ndarray,
        mask: np.ndarray,
        method: str = "telea",
        radius: int = 5
    ) -> np.ndarray:
        """마스크 영역 inpainting"""
        if method == "telea":
            flags = cv2.INPAINT_TELEA
        else:
            flags = cv2.INPAINT_NS

        # 마스크 확장 (더 넓은 영역 커버)
        kernel = np.ones((3, 3), np.uint8)
        mask_dilated = cv2.dilate(mask, kernel, iterations=2)

        result = cv2.inpaint(image, mask_dilated, radius, flags)
        return result

    def process_image(
        self,
        image: np.ndarray,
        padding: int = 8,
        inpaint_radius: int = 5
    ) -> Tuple[np.ndarray, List[TextRegion]]:
        """
        이미지에서 텍스트 제거

        Returns:
            (텍스트 제거된 이미지, 텍스트 영역 리스트)
        """
        # 텍스트 영역 감지
        regions = self.detect_text_regions(image)

        if not regions:
            return image.copy(), []

        # 마스크 생성
        mask = self.create_text_mask(image.shape, regions, padding)

        # Inpainting
        cleaned = self.inpaint_image(image, mask, radius=inpaint_radius)

        return cleaned, regions


class PDFToEditablePPTX:
    """PDF → 편집 가능 PPTX 변환 (텍스트 제거 + 재배치)"""

    def __init__(self, api_key: str = None):
        self.remover = TextRemover(api_key)

    def process_pdf(
        self,
        pdf_path: Path,
        zoom: float = 2.0
    ) -> List[ProcessedPage]:
        """PDF 처리"""
        import fitz

        pdf_path = Path(pdf_path)
        doc = fitz.open(pdf_path)
        pages = []

        print(f"📄 PDF 처리 중: {pdf_path.name} ({len(doc)}페이지)")

        for page_num in range(len(doc)):
            print(f"  🔍 페이지 {page_num + 1}/{len(doc)}: OCR + 텍스트 제거 중...")

            page = doc[page_num]
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat)

            # numpy array로 변환
            img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(
                pix.height, pix.width, pix.n
            )
            if pix.n == 4:
                img = cv2.cvtColor(img, cv2.COLOR_RGBA2BGR)
            else:
                img = cv2.cvtColor(img, cv2.COLOR_RGB2BGR)

            # 텍스트 제거
            cleaned, regions = self.remover.process_image(img)

            pages.append(ProcessedPage(
                page_num=page_num + 1,
                original_image=img,
                cleaned_image=cleaned,
                text_regions=regions,
                width=pix.width,
                height=pix.height
            ))

        doc.close()
        return pages

    def create_pptx(
        self,
        pages: List[ProcessedPage],
        output_path: Path
    ) -> Path:
        """PPTX 생성"""
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]

        for page in pages:
            slide = prs.slides.add_slide(blank_layout)

            # 좌표 변환 비율
            scale_x = prs.slide_width.emu / page.width
            scale_y = prs.slide_height.emu / page.height

            # 1. 클린 이미지를 배경으로
            _, buffer = cv2.imencode('.png', page.cleaned_image)
            slide.shapes.add_picture(
                io.BytesIO(buffer.tobytes()),
                Inches(0), Inches(0),
                width=prs.slide_width,
                height=prs.slide_height
            )

            # 2. 텍스트 박스 추가
            for region in page.text_regions:
                left = Emu(int(region.x * scale_x))
                top = Emu(int(region.y * scale_y))
                width = Emu(int(region.width * scale_x))
                height = Emu(int(region.height * scale_y))

                textbox = slide.shapes.add_textbox(left, top, width, height)
                tf = textbox.text_frame
                tf.word_wrap = True

                p = tf.paragraphs[0]
                p.text = region.text

                # 폰트 크기 추정
                est_font = max(8, min(44, int(region.height * scale_y / 914400 * 0.6)))
                p.font.size = Pt(est_font)
                p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        prs.save(output_path)
        print(f"  ✅ PPTX 저장: {output_path}")
        return output_path

    def convert(
        self,
        pdf_path: str,
        output_path: str = None
    ) -> Tuple[Path, int]:
        """
        PDF → 편집 가능 PPTX 변환

        Args:
            pdf_path: PDF 파일 경로
            output_path: 출력 경로

        Returns:
            (출력 파일 경로, 페이지 수)
        """
        pdf_path = Path(pdf_path)
        if output_path is None:
            output_path = pdf_path.with_name(f"{pdf_path.stem}_clean_편집가능.pptx")
        else:
            output_path = Path(output_path)

        # PDF 처리
        pages = self.process_pdf(pdf_path)

        # PPTX 생성
        self.create_pptx(pages, output_path)

        return output_path, len(pages)


def convert_pdf_to_clean_pptx(
    pdf_path: str,
    output_path: str = None,
    api_key: str = None
) -> Tuple[Path, int]:
    """
    PDF를 텍스트 제거된 깨끗한 PPTX로 변환

    Args:
        pdf_path: PDF 파일 경로
        output_path: 출력 경로
        api_key: Google Vision API 키

    Returns:
        (출력 파일 경로, 페이지 수)
    """
    converter = PDFToEditablePPTX(api_key)
    return converter.convert(pdf_path, output_path)


# CLI
if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(
        description="PDF 텍스트 제거 + Inpainting → 편집 가능 PPTX"
    )
    parser.add_argument("pdf_path", help="PDF 파일 경로")
    parser.add_argument("-o", "--output", help="출력 경로")
    parser.add_argument("--api-key", help="Google Vision API 키")

    args = parser.parse_args()

    try:
        pptx_path, count = convert_pdf_to_clean_pptx(
            args.pdf_path,
            args.output,
            args.api_key
        )
        print(f"\n✅ 변환 완료!")
        print(f"   파일: {pptx_path}")
        print(f"   슬라이드: {count}장")
    except Exception as e:
        print(f"❌ 오류: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
