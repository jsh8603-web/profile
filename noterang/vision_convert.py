#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Google Cloud Vision OCR 기반 PDF → 편집 가능 PPTX 변환

- PDF 페이지를 이미지로 렌더링
- Google Cloud Vision API로 텍스트 + 위치 추출
- 원본 이미지 배경 + 편집 가능한 텍스트 박스 오버레이
"""
import io
import os
import sys
import base64
import requests
from pathlib import Path
from typing import List, Tuple, Optional
from dataclasses import dataclass, field

if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8')

from dotenv import load_dotenv
load_dotenv(Path(__file__).parent.parent / '.env.local')


@dataclass
class TextBox:
    """텍스트 박스 정보"""
    text: str
    x: float      # 픽셀 좌표
    y: float
    width: float
    height: float
    confidence: float = 1.0


@dataclass
class PageOCR:
    """페이지 OCR 결과"""
    page_num: int
    width: float   # 이미지 크기 (픽셀)
    height: float
    image_data: bytes
    text_boxes: List[TextBox] = field(default_factory=list)
    full_text: str = ""


class GoogleVisionOCR:
    """Google Cloud Vision OCR"""

    def __init__(self, api_key: str = None):
        self.api_key = api_key or os.getenv('GOOGLE_VISION_API_KEY')
        if not self.api_key:
            raise ValueError("GOOGLE_VISION_API_KEY가 필요합니다")
        self.api_url = "https://vision.googleapis.com/v1/images:annotate"

    def ocr_image(self, image_data: bytes) -> Tuple[str, List[TextBox]]:
        """
        이미지에서 텍스트와 위치 추출

        Returns:
            (전체 텍스트, 텍스트 박스 리스트)
        """
        payload = {
            "requests": [{
                "image": {"content": base64.b64encode(image_data).decode('utf-8')},
                "features": [{"type": "DOCUMENT_TEXT_DETECTION"}]
            }]
        }

        response = requests.post(
            f"{self.api_url}?key={self.api_key}",
            json=payload,
            timeout=60
        )
        result = response.json()

        if 'error' in result:
            raise Exception(f"Vision API 오류: {result['error']}")

        responses = result.get('responses', [{}])
        if not responses:
            return "", []

        resp = responses[0]

        # 전체 텍스트
        full_text = ""
        annotations = resp.get('textAnnotations', [])
        if annotations:
            full_text = annotations[0].get('description', '')

        # 개별 텍스트 박스 (단어/블록 단위)
        text_boxes = []
        full_annotation = resp.get('fullTextAnnotation', {})

        for page in full_annotation.get('pages', []):
            for block in page.get('blocks', []):
                # 블록 단위로 텍스트 추출
                block_text = ""
                for paragraph in block.get('paragraphs', []):
                    para_text = ""
                    for word in paragraph.get('words', []):
                        word_text = "".join(
                            symbol.get('text', '')
                            for symbol in word.get('symbols', [])
                        )
                        para_text += word_text + " "
                    block_text += para_text.strip() + "\n"

                block_text = block_text.strip()
                if not block_text:
                    continue

                # 바운딩 박스
                vertices = block.get('boundingBox', {}).get('vertices', [])
                if len(vertices) >= 4:
                    xs = [v.get('x', 0) for v in vertices]
                    ys = [v.get('y', 0) for v in vertices]
                    x = min(xs)
                    y = min(ys)
                    w = max(xs) - x
                    h = max(ys) - y

                    confidence = block.get('confidence', 1.0)

                    text_boxes.append(TextBox(
                        text=block_text,
                        x=x, y=y,
                        width=w, height=h,
                        confidence=confidence
                    ))

        return full_text, text_boxes


class VisionConverter:
    """Vision OCR 기반 PDF → PPTX 변환기"""

    def __init__(self, api_key: str = None):
        self.ocr = GoogleVisionOCR(api_key)

    def pdf_to_pages(self, pdf_path: Path, zoom: float = 2.0) -> List[PageOCR]:
        """PDF를 페이지별 OCR 결과로 변환"""
        import fitz

        pdf_path = Path(pdf_path)
        doc = fitz.open(pdf_path)
        pages = []

        print(f"📄 PDF 처리 중: {pdf_path.name} ({len(doc)}페이지)")

        for page_num in range(len(doc)):
            print(f"  🔍 페이지 {page_num + 1}/{len(doc)} OCR 중...")

            page = doc[page_num]
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat)
            img_data = pix.tobytes("png")

            # OCR 실행
            full_text, text_boxes = self.ocr.ocr_image(img_data)

            pages.append(PageOCR(
                page_num=page_num + 1,
                width=pix.width,
                height=pix.height,
                image_data=img_data,
                text_boxes=text_boxes,
                full_text=full_text
            ))

        doc.close()
        return pages

    def create_pptx(
        self,
        pages: List[PageOCR],
        output_path: Path,
        include_background: bool = True
    ) -> Path:
        """OCR 결과로 PPTX 생성"""
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

        prs = Presentation()

        # 슬라이드 크기 설정 (16:9 와이드스크린)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]

        for page in pages:
            slide = prs.slides.add_slide(blank_layout)

            # 이미지 → 슬라이드 좌표 변환 비율
            scale_x = prs.slide_width.emu / page.width
            scale_y = prs.slide_height.emu / page.height

            # 1. 배경 이미지 삽입
            if include_background:
                slide.shapes.add_picture(
                    io.BytesIO(page.image_data),
                    Inches(0), Inches(0),
                    width=prs.slide_width,
                    height=prs.slide_height
                )

            # 2. 텍스트 박스 오버레이 (편집 가능)
            for tb in page.text_boxes:
                left = Emu(int(tb.x * scale_x))
                top = Emu(int(tb.y * scale_y))
                width = Emu(int(tb.width * scale_x))
                height = Emu(int(tb.height * scale_y))

                textbox = slide.shapes.add_textbox(left, top, width, height)
                tf = textbox.text_frame
                tf.word_wrap = True
                tf.auto_size = None

                # 텍스트 설정
                p = tf.paragraphs[0]
                p.text = tb.text

                # 폰트 크기 추정 (박스 높이 기반)
                estimated_font_size = max(8, min(48, int(tb.height * scale_y / 914400 * 0.7)))
                p.font.size = Pt(estimated_font_size)

                # 배경 포함시 텍스트 투명하게 (선택하면 보임)
                if include_background:
                    # 투명 텍스트 (RGBA 지원 안해서 배경색과 동일하게)
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    # 실제로는 선택하면 편집 가능
                else:
                    p.font.color.rgb = RGBColor(0, 0, 0)

        prs.save(output_path)
        print(f"  ✅ PPTX 저장: {output_path}")
        return output_path

    def convert(
        self,
        pdf_path: str,
        output_path: str = None,
        include_background: bool = True
    ) -> Tuple[Path, int]:
        """
        PDF → 편집 가능 PPTX 변환

        Args:
            pdf_path: PDF 파일 경로
            output_path: 출력 경로 (None이면 자동)
            include_background: 배경 이미지 포함 여부

        Returns:
            (출력 파일 경로, 페이지 수)
        """
        pdf_path = Path(pdf_path)
        if output_path is None:
            output_path = pdf_path.with_name(f"{pdf_path.stem}_편집가능.pptx")
        else:
            output_path = Path(output_path)

        # OCR 실행
        pages = self.pdf_to_pages(pdf_path)

        # PPTX 생성
        self.create_pptx(pages, output_path, include_background)

        return output_path, len(pages)


def convert_pdf_to_editable_pptx(
    pdf_path: str,
    output_path: str = None,
    include_background: bool = True,
    api_key: str = None
) -> Tuple[Path, int]:
    """
    PDF를 편집 가능한 PPTX로 변환 (간편 함수)

    Args:
        pdf_path: PDF 파일 경로
        output_path: 출력 경로
        include_background: 배경 이미지 포함 여부
        api_key: Google Vision API 키

    Returns:
        (출력 파일 경로, 페이지 수)

    Example:
        >>> pptx, count = convert_pdf_to_editable_pptx("slides.pdf")
        >>> print(f"변환 완료: {pptx} ({count}장)")
    """
    converter = VisionConverter(api_key)
    return converter.convert(pdf_path, output_path, include_background)


# CLI
if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(
        description="Google Cloud Vision OCR 기반 PDF → 편집 가능 PPTX 변환"
    )
    parser.add_argument("pdf_path", help="PDF 파일 경로")
    parser.add_argument("-o", "--output", help="출력 경로")
    parser.add_argument(
        "--no-background",
        action="store_true",
        help="배경 이미지 제외 (텍스트만)"
    )
    parser.add_argument("--api-key", help="Google Vision API 키")

    args = parser.parse_args()

    try:
        pptx_path, count = convert_pdf_to_editable_pptx(
            args.pdf_path,
            args.output,
            include_background=not args.no_background,
            api_key=args.api_key
        )
        print(f"\n✅ 변환 완료!")
        print(f"   파일: {pptx_path}")
        print(f"   슬라이드: {count}장")
    except Exception as e:
        print(f"❌ 오류: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
