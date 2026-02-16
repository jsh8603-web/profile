#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
노트랑 완전 자동화 v2 - 멀티 에이전트 시스템 통합
- 연구 → 슬라이드 생성 → 다운로드 → PPTX 변환
- 타임아웃 시 헬퍼 에이전트 투입
- 버그 발생 시 복구 에이전트 투입
"""
import json
import subprocess
import sys
import asyncio
import os
from pathlib import Path
from datetime import datetime
import time

if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8')

# 경로 설정
SCRIPT_DIR = Path(__file__).parent
sys.path.insert(0, str(SCRIPT_DIR))

from noterang.agent_manager import get_noterang_agent, AgentTask, AgentStatus

# 설정
NLM_EXE = Path.home() / "AppData/Roaming/Python/Python313/Scripts/nlm.exe"
DOWNLOAD_DIR = Path("G:/내 드라이브/notebooklm")
AUTH_DIR = Path.home() / ".notebooklm-mcp-cli"


def sync_auth():
    """인증 동기화"""
    root_auth = AUTH_DIR / "auth.json"
    profile_dir = AUTH_DIR / "profiles" / "default"

    if not root_auth.exists():
        return False

    with open(root_auth) as f:
        root_data = json.load(f)

    cookies_dict = root_data.get('cookies', {})
    cookies_list = [
        {"name": n, "value": v, "domain": ".google.com", "path": "/",
         "expires": -1, "httpOnly": False, "secure": True, "sameSite": "Lax"}
        for n, v in cookies_dict.items()
    ]

    profile_dir.mkdir(parents=True, exist_ok=True)

    with open(profile_dir / "cookies.json", "w") as f:
        json.dump(cookies_list, f)

    with open(profile_dir / "metadata.json", "w") as f:
        json.dump({
            "csrf_token": root_data.get("csrf_token", ""),
            "session_id": root_data.get("session_id", ""),
            "email": root_data.get("email", ""),
            "last_validated": datetime.now().isoformat()
        }, f)

    with open(profile_dir / "auth.json", "w") as f:
        json.dump(root_data, f)

    return True


def run_nlm(args, timeout=120):
    """nlm CLI 실행"""
    cmd = [str(NLM_EXE)] + args
    env = os.environ.copy()
    env['PYTHONIOENCODING'] = 'utf-8'

    try:
        result = subprocess.run(cmd, capture_output=True, timeout=timeout, env=env)
        stdout = result.stdout.decode('utf-8', errors='replace') if result.stdout else ''
        stderr = result.stderr.decode('utf-8', errors='replace') if result.stderr else ''
        return result.returncode == 0, stdout, stderr
    except Exception as e:
        return False, '', str(e)


def check_auth():
    """인증 확인"""
    sync_auth()
    success, stdout, _ = run_nlm(["login", "--check"])
    return success and stdout and "valid" in stdout.lower()


async def ensure_auth():
    """인증 확인 및 필요시 자동 로그인"""
    # 먼저 기존 인증 확인
    if check_auth():
        return True

    # 실패하면 자동 로그인 시도
    print("  인증 만료 - 자동 로그인 시도...")
    from auto_login import auto_login

    # headless로 먼저 시도
    if await auto_login(headless=True, timeout=30):
        sync_auth()
        return check_auth()

    # 실패하면 브라우저로 시도
    print("  백그라운드 실패 - 브라우저로 재시도...")
    if await auto_login(headless=False, timeout=120):
        sync_auth()
        return check_auth()

    return False


async def download_via_browser(notebook_id, output_dir):
    """브라우저를 통한 다운로드 (CLI 버그 우회)"""
    from playwright.async_api import async_playwright

    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    user_data_dir = AUTH_DIR / "browser_profile"

    async with async_playwright() as p:
        context = await p.chromium.launch_persistent_context(
            user_data_dir=str(user_data_dir),
            headless=False,
            downloads_path=str(output_dir),
            accept_downloads=True,
            args=['--disable-blink-features=AutomationControlled'],
            viewport={'width': 1920, 'height': 1080},
        )

        page = context.pages[0] if context.pages else await context.new_page()
        notebook_url = f"https://notebooklm.google.com/notebook/{notebook_id}"

        try:
            await page.goto(notebook_url, wait_until='domcontentloaded', timeout=30000)
        except:
            pass

        await asyncio.sleep(8)

        downloaded_path = None
        menu_btns = await page.query_selector_all('[aria-haspopup="menu"], button[aria-label*="more"]')

        for menu_btn in menu_btns[-10:]:
            try:
                await menu_btn.click(force=True)
                await asyncio.sleep(1)

                dl_item = await page.query_selector('[role="menuitem"]:has-text("다운로드"), [role="menuitem"]:has-text("Download")')
                if dl_item:
                    async with page.expect_download(timeout=30000) as download_info:
                        await dl_item.click()

                    download = await download_info.value
                    filename = f"slides_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
                    downloaded_path = output_dir / filename
                    await download.save_as(str(downloaded_path))
                    break
            except:
                await page.keyboard.press('Escape')

        await asyncio.sleep(2)
        await context.close()
        return downloaded_path


def pdf_to_pptx(pdf_path):
    """PDF를 PPTX로 변환"""
    import fitz
    from pptx import Presentation
    from pptx.util import Inches
    import io

    pdf_path = Path(pdf_path)
    output_path = pdf_path.with_suffix('.pptx')

    doc = fitz.open(pdf_path)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for page_num in range(len(doc)):
        page = doc[page_num]
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img_data = pix.tobytes("png")

        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            io.BytesIO(img_data),
            Inches(0), Inches(0),
            width=prs.slide_width,
            height=prs.slide_height
        )

    doc.close()
    prs.save(output_path)
    return output_path, len(prs.slides)


async def run_full_automation_async(title, research_queries, focus=None, language="ko"):
    """전체 자동화 실행 (멀티 에이전트 버전)"""
    agent = get_noterang_agent()

    print("=" * 60)
    print(f"노트랑 멀티 에이전트 자동화: {title}")
    print("=" * 60)

    # 1. 인증 확인 (자동 로그인 포함)
    print("\n[1/6] 인증 확인...")
    if not await ensure_auth():
        print("  ❌ 인증 실패 - 수동 로그인 필요")
        return None
    print("  ✓ 인증 유효")

    # 2. 노트북 확인/생성
    print(f"\n[2/6] 노트북 확인...")
    success, stdout, _ = run_nlm(["list", "notebooks"])

    notebook_id = None
    if success:
        try:
            notebooks = json.loads(stdout)
            for nb in notebooks:
                if nb.get('title') == title:
                    notebook_id = nb.get('id')
                    print(f"  기존 노트북: {notebook_id[:8]}...")
                    break
        except:
            pass

    if not notebook_id:
        success, stdout, _ = run_nlm(["notebook", "create", title])
        if success:
            try:
                data = json.loads(stdout)
                notebook_id = data.get('id')
                print(f"  새 노트북 생성: {notebook_id[:8]}...")
            except:
                pass

    if not notebook_id:
        print("  ❌ 노트북 생성 실패")
        return None

    # 3. 연구 (멀티 에이전트 모니터링)
    print(f"\n[3/6] 연구 자료 수집 (멀티 에이전트)...")
    total_sources = 0
    for query in research_queries:
        success, count = await agent.run_research_with_monitoring(notebook_id, query)
        total_sources += count
    print(f"\n  총 {total_sources}개 소스 추가")

    # 4. 슬라이드 생성 (멀티 에이전트 모니터링)
    print(f"\n[4/6] 슬라이드 생성 (멀티 에이전트)...")
    artifact_id = await agent.create_slides_with_monitoring(
        notebook_id,
        language=language,
        focus=focus
    )

    if not artifact_id:
        print("  ⚠️ 슬라이드 생성 실패")

    # 5. 다운로드
    print(f"\n[5/6] 다운로드...")
    pdf_path = await download_via_browser(notebook_id, DOWNLOAD_DIR)

    if not pdf_path or not pdf_path.exists():
        print("  ❌ 다운로드 실패")
        return None
    print(f"  ✓ PDF: {pdf_path.name}")

    # 6. PPTX 변환
    print(f"\n[6/6] PPTX 변환...")
    pptx_path, slide_count = pdf_to_pptx(pdf_path)
    print(f"  ✓ PPTX: {pptx_path.name} ({slide_count}슬라이드)")

    # 메모리 통계 출력
    stats = agent.get_memory_stats()
    print("\n" + "-" * 40)
    print("에이전트 통계:")
    print(f"  생성된 에이전트: {stats['agents_created']}개")
    print(f"  총 작업: {stats['performance']['total_tasks']}")
    print(f"  성공률: {stats['performance']['successful_tasks']}/{stats['performance']['total_tasks']}")

    print("\n" + "=" * 60)
    print("완료!")
    print(f"  PDF:  {pdf_path}")
    print(f"  PPTX: {pptx_path}")
    print("=" * 60)

    return {
        'notebook_id': notebook_id,
        'pdf': str(pdf_path),
        'pptx': str(pptx_path),
        'slides': slide_count,
        'agent_stats': stats
    }


def run_full_automation(title, research_queries, focus=None, language="ko"):
    """동기 버전 래퍼"""
    return asyncio.run(run_full_automation_async(title, research_queries, focus, language))


if __name__ == "__main__":
    result = run_full_automation(
        title="견관절회전근개 파열",
        research_queries=[
            "회전근개 파열 원인 병인",
            "회전근개 파열 수술 치료",
            "회전근개 파열 재활 운동",
        ],
        focus="병인, 치료방법, 재활법",
        language="ko"
    )

    if result:
        print(f"\n결과: {json.dumps(result, ensure_ascii=False, indent=2)}")
