#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""PDF를 PPTX로 변환하고 합치기"""
import sys
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from pathlib import Path
import io

if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8')

def pdf_to_pptx(pdf_path, output_path=None):
    """PDF를 PPTX로 변환"""
    pdf_path = Path(pdf_path)
    if output_path is None:
        output_path = pdf_path.with_suffix('.pptx')
    
    print(f"변환: {pdf_path.name}", end=" ")
    
    doc = fitz.open(pdf_path)
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    for page_num in range(len(doc)):
        page = doc[page_num]
        mat = fitz.Matrix(2, 2)
        pix = page.get_pixmap(matrix=mat)
        img_data = pix.tobytes("png")
        
        slide = prs.slides.add_slide(blank_layout)
        img_stream = io.BytesIO(img_data)
        slide.shapes.add_picture(
            img_stream, 
            Inches(0), Inches(0),
            width=prs.slide_width,
            height=prs.slide_height
        )
    
    doc.close()
    prs.save(output_path)
    print(f"→ {len(prs.slides)}슬라이드")
    
    return output_path

def merge_pptx(pptx_files, output_path):
    """여러 PPTX 합치기"""
    print(f"\n합치는 중: {len(pptx_files)}개 파일...")
    
    merged = Presentation()
    merged.slide_width = Inches(13.333)
    merged.slide_height = Inches(7.5)
    blank_layout = merged.slide_layouts[6]
    
    total_slides = 0
    
    for pptx_file in pptx_files:
        src = Presentation(pptx_file)
        
        for slide in src.slides:
            new_slide = merged.slides.add_slide(blank_layout)
            
            for shape in slide.shapes:
                if hasattr(shape, 'image'):
                    img = shape.image
                    img_stream = io.BytesIO(img.blob)
                    new_slide.shapes.add_picture(
                        img_stream,
                        shape.left, shape.top,
                        shape.width, shape.height
                    )
            total_slides += 1
    
    merged.save(output_path)
    print(f"→ 총 {total_slides}슬라이드")
    
    return output_path

def main():
    downloads_dir = Path("downloads")
    
    pdf_files = sorted(downloads_dir.glob("*.pdf"))
    print(f"PDF 파일 {len(pdf_files)}개 발견\n")
    
    pptx_files = []
    for pdf_file in pdf_files:
        pptx_path = pdf_to_pptx(pdf_file)
        pptx_files.append(pptx_path)
    
    if pptx_files:
        merged_path = downloads_dir / "n8n_마스터_가이드_통합.pptx"
        merge_pptx(pptx_files, merged_path)
        
        print(f"\n✓ 완료: {merged_path}")

if __name__ == "__main__":
    main()
