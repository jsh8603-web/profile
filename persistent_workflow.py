#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
NotebookLM 영구 인증 워크플로우
- 첫 실행: 수동 로그인 → 세션 영구 저장
- 이후 실행: 저장된 세션으로 자동 인증
"""
import asyncio
import json
import sys
import os
import io
import re
import time
from pathlib import Path
from datetime import datetime

os.environ['PYTHONUNBUFFERED'] = '1'
if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8', line_buffering=True)

# 설정
AUTH_DIR = Path.home() / ".notebooklm-mcp-cli"
PERSISTENT_PROFILE = AUTH_DIR / "persistent_browser"  # 영구 브라우저 프로필
DOWNLOAD_DIR = Path("G:/내 드라이브/notebooklm")


def log(msg):
    print(msg, flush=True)


def safe_filename(name):
    return re.sub(r'[<>:"/\\|?*]', '', name).strip()


def pdf_to_pptx(pdf_path, output_name=None):
    """PDF를 PPTX로 변환"""
    import fitz
    from pptx import Presentation
    from pptx.util import Inches

    pdf_path = Path(pdf_path)
    if output_name:
        output_path = DOWNLOAD_DIR / f"{safe_filename(output_name)}.pptx"
    else:
        output_path = pdf_path.with_suffix('.pptx')

    doc = fitz.open(pdf_path)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for page_num in range(len(doc)):
        page = doc[page_num]
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img_data = pix.tobytes("png")

        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            io.BytesIO(img_data),
            Inches(0), Inches(0),
            width=prs.slide_width,
            height=prs.slide_height
        )

    doc.close()
    prs.save(output_path)
    return output_path, len(prs.slides)


async def wait_for_login(page, timeout=300):
    """로그인 완료 대기 (최대 5분)"""
    log("  로그인 대기 중... (브라우저에서 로그인해주세요)")

    start = time.time()
    while time.time() - start < timeout:
        await asyncio.sleep(3)
        elapsed = int(time.time() - start)

        url = page.url
        if "notebooklm.google.com" in url and "accounts.google" not in url:
            log(f"  ✓ 로그인 완료! ({elapsed}초)")
            return True

        if elapsed % 15 == 0:
            log(f"  대기 중... {elapsed}초")

    log("  ❌ 로그인 타임아웃")
    return False


async def download_slides(page, keyword: str, display_name: str) -> dict:
    """노트북에서 슬라이드 다운로드"""
    result = {'name': display_name, 'pdf': None, 'pptx': None, 'slides': 0}

    try:
        # 홈으로 이동
        log(f"  홈으로 이동...")
        await page.goto("https://notebooklm.google.com/", wait_until='networkidle', timeout=30000)
        await asyncio.sleep(3)

        # 노트북 찾기
        log(f"  노트북 찾기: '{keyword}'...")
        notebook_link = None

        all_elements = await page.query_selector_all('a, div[role="button"], [class*="notebook"], tr')
        for el in all_elements:
            try:
                text = await el.inner_text()
                if keyword in text:
                    notebook_link = el
                    log(f"  → 발견: {text[:50]}...")
                    break
            except:
                continue

        if not notebook_link:
            # 텍스트로 직접 찾기
            notebook_link = await page.query_selector(f'text="{keyword}"')

        if not notebook_link:
            log(f"  ❌ 노트북 찾기 실패")
            # 스크린샷
            await page.screenshot(path=str(DOWNLOAD_DIR / f"debug_{safe_filename(display_name)}.png"))
            return result

        await notebook_link.click()
        await asyncio.sleep(5)

        log(f"  노트북 열림, 다운로드 시도...")

        # 다운로드 시도
        pdf_path = None

        # 방법 1: 메뉴 버튼
        menu_buttons = await page.query_selector_all(
            '[aria-haspopup="menu"], button[aria-label*="more"], button[aria-label*="More"], '
            'button[aria-label*="옵션"], [data-tooltip*="more"]'
        )
        log(f"  메뉴 버튼 {len(menu_buttons)}개 발견")

        for idx, menu_btn in enumerate(menu_buttons[-15:]):
            try:
                await menu_btn.click(force=True)
                await asyncio.sleep(1.5)

                download_item = await page.query_selector(
                    '[role="menuitem"]:has-text("Download PDF"), '
                    '[role="menuitem"]:has-text("Download"), '
                    '[role="menuitem"]:has-text("다운로드")'
                )

                if download_item:
                    log(f"  다운로드 메뉴 발견! (버튼 {idx})")
                    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                    filename = f"{safe_filename(display_name)}_{timestamp}.pdf"

                    async with page.expect_download(timeout=60000) as download_info:
                        await download_item.click()

                    download = await download_info.value
                    pdf_path = DOWNLOAD_DIR / filename
                    await download.save_as(str(pdf_path))
                    log(f"  ✓ PDF: {filename}")
                    break
            except Exception as e:
                await page.keyboard.press('Escape')
                await asyncio.sleep(0.3)

        # 방법 2: 좌표 기반
        if not pdf_path:
            log(f"  좌표 기반 시도...")
            coordinates = [(1846, 365), (1800, 400), (1850, 350), (1750, 380)]

            for x, y in coordinates:
                try:
                    await page.mouse.click(x, y)
                    await asyncio.sleep(1.5)

                    menu_items = await page.query_selector_all('[role="menuitem"]')
                    for item in menu_items:
                        text = await item.inner_text()
                        if 'Download' in text or '다운로드' in text:
                            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                            filename = f"{safe_filename(display_name)}_{timestamp}.pdf"

                            async with page.expect_download(timeout=60000) as download_info:
                                await item.click()

                            download = await download_info.value
                            pdf_path = DOWNLOAD_DIR / filename
                            await download.save_as(str(pdf_path))
                            log(f"  ✓ PDF (좌표): {filename}")
                            break
                    if pdf_path:
                        break
                except:
                    await page.keyboard.press('Escape')

        if pdf_path and pdf_path.exists():
            result['pdf'] = str(pdf_path)
            log(f"  PPTX 변환 중...")
            pptx_path, slide_count = pdf_to_pptx(pdf_path, display_name)
            result['pptx'] = str(pptx_path)
            result['slides'] = slide_count
            log(f"  ✓ PPTX: {pptx_path.name} ({slide_count}장)")
        else:
            log(f"  ❌ 다운로드 실패")
            await page.screenshot(path=str(DOWNLOAD_DIR / f"fail_{safe_filename(display_name)}.png"))

    except Exception as e:
        log(f"  ❌ 오류: {e}")

    return result


async def main():
    from playwright.async_api import async_playwright

    log("=" * 60)
    log("NotebookLM 영구 인증 워크플로우")
    log("=" * 60)
    log("  - 첫 실행: 브라우저에서 수동 로그인")
    log("  - 이후 실행: 저장된 세션 자동 사용")
    log("=" * 60)

    DOWNLOAD_DIR.mkdir(parents=True, exist_ok=True)
    PERSISTENT_PROFILE.mkdir(parents=True, exist_ok=True)

    # 처리할 노트북
    notebooks = [
        ("족저근막염", "족저근막염"),
        ("족모지외반증", "족모지외반증"),
        ("발톱무좀", "발톱무좀"),
    ]

    results = []

    async with async_playwright() as p:
        log("\n[1] 영구 브라우저 프로필 시작...")
        log(f"    프로필: {PERSISTENT_PROFILE}")

        # 영구 프로필로 브라우저 시작 (로그인 상태 유지됨)
        context = await p.chromium.launch_persistent_context(
            user_data_dir=str(PERSISTENT_PROFILE),
            headless=False,
            downloads_path=str(DOWNLOAD_DIR),
            accept_downloads=True,
            args=[
                '--disable-blink-features=AutomationControlled',
                '--disable-infobars',
            ],
            viewport={'width': 1920, 'height': 1080},
        )

        page = context.pages[0] if context.pages else await context.new_page()

        # NotebookLM 접속
        log("\n[2] NotebookLM 접속...")
        try:
            await page.goto("https://notebooklm.google.com/", wait_until='networkidle', timeout=30000)
        except:
            pass
        await asyncio.sleep(3)

        # 로그인 상태 확인
        if "accounts.google.com" in page.url:
            log("\n[3] 로그인 필요...")
            logged_in = await wait_for_login(page, timeout=300)
            if not logged_in:
                await context.close()
                return []
        else:
            log("\n[3] ✓ 이미 로그인됨 (저장된 세션)")

        # 슬라이드 다운로드
        log("\n[4] 슬라이드 다운로드...")
        for keyword, display_name in notebooks:
            log(f"\n{'='*40}")
            log(f"[{display_name}]")
            log('='*40)
            result = await download_slides(page, keyword, display_name)
            results.append(result)
            await asyncio.sleep(2)

        await context.close()

    # 결과 출력
    log("\n" + "=" * 60)
    log("완료 결과:")
    log("=" * 60)

    success_count = 0
    for r in results:
        status = "✓" if r.get('pptx') else "❌"
        log(f"  {status} {r['name']}")
        if r.get('pptx'):
            log(f"      PPTX: {Path(r['pptx']).name} ({r['slides']}장)")
            success_count += 1

    log(f"\n성공: {success_count}/{len(notebooks)}")
    log(f"저장 위치: {DOWNLOAD_DIR}")
    log("=" * 60)

    # 다음 실행 안내
    if success_count > 0:
        log("\n※ 다음 실행 시 자동 로그인됩니다 (세션 저장됨)")

    return results


if __name__ == "__main__":
    asyncio.run(main())
