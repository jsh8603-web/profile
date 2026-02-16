#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
기존 노트북에서 15장 한글 슬라이드 재생성
- 멀티 에이전트: 타임아웃/에러 시 헬퍼 에이전트 투입
- 슬라이드 생성 체크 기능 포함
- 한글 슬라이드 강제 설정
"""
import asyncio
import sys
import os
from pathlib import Path
from datetime import datetime
import re
import io

# 버퍼링 비활성화
os.environ['PYTHONUNBUFFERED'] = '1'
if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8', line_buffering=True)

def log(msg):
    """즉시 출력"""
    print(msg, flush=True)

DOWNLOAD_DIR = Path("G:/내 드라이브/notebooklm")
AUTH_DIR = Path.home() / ".notebooklm-mcp-cli"
# Chrome 기본 프로필 사용 (이미 Google 로그인됨)
CHROME_USER_DATA = Path.home() / "AppData/Local/Google/Chrome/User Data"
BROWSER_PROFILE = CHROME_USER_DATA if CHROME_USER_DATA.exists() else AUTH_DIR / "browser_profile"


def pdf_to_pptx(pdf_path, output_name=None):
    """PDF를 PPTX로 변환하여 교체"""
    import fitz
    from pptx import Presentation
    from pptx.util import Inches

    pdf_path = Path(pdf_path)
    if output_name:
        output_path = DOWNLOAD_DIR / f"{output_name}.pptx"
    else:
        output_path = pdf_path.with_suffix('.pptx')

    doc = fitz.open(pdf_path)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for page_num in range(len(doc)):
        page = doc[page_num]
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img_data = pix.tobytes("png")

        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            io.BytesIO(img_data),
            Inches(0), Inches(0),
            width=prs.slide_width,
            height=prs.slide_height
        )

    doc.close()
    prs.save(output_path)
    return output_path, len(prs.slides)


async def check_slides_created(page, timeout=180):
    """슬라이드 생성 완료 체크 (헬퍼 에이전트 역할)"""
    log("  [체크] 슬라이드 생성 상태 모니터링...")
    start_time = asyncio.get_event_loop().time()

    while asyncio.get_event_loop().time() - start_time < timeout:
        await asyncio.sleep(5)
        elapsed = int(asyncio.get_event_loop().time() - start_time)

        # 완료 확인 방법들
        try:
            # 1. 슬라이드 카드 확인
            slide_card = await page.query_selector('[class*="slide"], [class*="artifact-card"]')
            if slide_card:
                # 다운로드 버튼 확인
                download_available = await page.query_selector(
                    'button:has-text("Download"), button:has-text("다운로드"), '
                    '[aria-label*="download"], [aria-label*="Download"]'
                )
                if download_available:
                    log(f"  [체크] ✓ 슬라이드 생성 완료! ({elapsed}초)")
                    return True

            # 2. more_vert 메뉴 확인 (슬라이드 카드에)
            more_menu = await page.query_selector('[aria-label*="more"], [aria-haspopup="menu"]')
            if more_menu:
                # 클릭해서 다운로드 옵션 확인
                await more_menu.click()
                await asyncio.sleep(0.5)
                download_item = await page.query_selector('[role="menuitem"]:has-text("Download")')
                await page.keyboard.press('Escape')
                if download_item:
                    log(f"  [체크] ✓ 다운로드 가능 확인! ({elapsed}초)")
                    return True

            # 3. 로딩/진행 상태 확인
            loading = await page.query_selector('[class*="loading"], [class*="progress"], [class*="spinner"]')
            if loading:
                log(f"  [체크] 생성 중... ({elapsed}초)", )

        except Exception as e:
            log(f"  [체크] 체크 중 오류: {e}")

        # 타임아웃 경고
        if elapsed > 120:
            log(f"  [경고] 생성 지연됨 ({elapsed}초) - 헬퍼 대기 중...")

    log(f"  [체크] ⚠️ 타임아웃 ({timeout}초)")
    return False


async def helper_agent_download(page, notebook_name):
    """헬퍼 에이전트: 다운로드 전담"""
    log("  [헬퍼] 다운로드 에이전트 투입...")

    downloaded = False
    pdf_path = None

    # 여러 방법으로 다운로드 시도
    methods = [
        ("메뉴 버튼", try_menu_download),
        ("좌표 클릭", try_coordinate_download),
        ("키보드 단축키", try_keyboard_download),
    ]

    for method_name, method_func in methods:
        log(f"  [헬퍼] 시도: {method_name}...")
        try:
            pdf_path = await method_func(page, notebook_name)
            if pdf_path and pdf_path.exists():
                downloaded = True
                log(f"  [헬퍼] ✓ {method_name} 성공!")
                break
        except Exception as e:
            log(f"  [헬퍼] {method_name} 실패: {e}")

    return pdf_path if downloaded else None


async def try_menu_download(page, notebook_name):
    """메뉴 버튼으로 다운로드"""
    menu_buttons = await page.query_selector_all(
        '[aria-haspopup="menu"], button[aria-label*="more"], button[aria-label*="More"]'
    )

    for menu_btn in menu_buttons[-5:]:
        try:
            await menu_btn.click(force=True)
            await asyncio.sleep(1.5)

            download_item = await page.query_selector(
                '[role="menuitem"]:has-text("Download PDF"), '
                '[role="menuitem"]:has-text("Download"), '
                '[role="menuitem"]:has-text("다운로드")'
            )

            if download_item:
                timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                filename = f"{notebook_name}_{timestamp}.pdf"

                async with page.expect_download(timeout=60000) as download_info:
                    await download_item.click()

                download = await download_info.value
                pdf_path = DOWNLOAD_DIR / filename
                await download.save_as(str(pdf_path))
                return pdf_path
        except:
            await page.keyboard.press('Escape')

    return None


async def try_coordinate_download(page, notebook_name):
    """좌표 기반 다운로드 (이전 성공 위치)"""
    # 스튜디오 패널의 슬라이드 카드 more 버튼
    coordinates = [(1846, 365), (1800, 400), (1850, 350)]

    for x, y in coordinates:
        try:
            await page.mouse.click(x, y)
            await asyncio.sleep(1.5)

            menu_items = await page.query_selector_all('[role="menuitem"]')
            for item in menu_items:
                text = await item.inner_text()
                if 'Download' in text or '다운로드' in text:
                    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                    filename = f"{notebook_name}_{timestamp}.pdf"

                    async with page.expect_download(timeout=60000) as download_info:
                        await item.click()

                    download = await download_info.value
                    pdf_path = DOWNLOAD_DIR / filename
                    await download.save_as(str(pdf_path))
                    return pdf_path
        except:
            await page.keyboard.press('Escape')

    return None


async def try_keyboard_download(page, notebook_name):
    """키보드 단축키로 다운로드"""
    try:
        # Ctrl+Shift+S 또는 다른 단축키 시도
        await page.keyboard.press('Escape')
        await asyncio.sleep(0.5)

        # Tab으로 다운로드 버튼으로 이동 시도
        for _ in range(10):
            await page.keyboard.press('Tab')
            await asyncio.sleep(0.2)

            focused = await page.evaluate('document.activeElement.innerText')
            if focused and ('Download' in focused or '다운로드' in focused):
                timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                filename = f"{notebook_name}_{timestamp}.pdf"

                async with page.expect_download(timeout=60000) as download_info:
                    await page.keyboard.press('Enter')

                download = await download_info.value
                pdf_path = DOWNLOAD_DIR / filename
                await download.save_as(str(pdf_path))
                return pdf_path
    except:
        pass

    return None


async def debug_agent_screenshot(page, name):
    """디버그 에이전트: 스크린샷 저장"""
    try:
        screenshot_path = DOWNLOAD_DIR / f"debug_{name}_{datetime.now().strftime('%H%M%S')}.png"
        await page.screenshot(path=str(screenshot_path))
        log(f"  [디버그] 스크린샷 저장: {screenshot_path.name}")
        return screenshot_path
    except Exception as e:
        log(f"  [디버그] 스크린샷 실패: {e}")
        return None


async def create_korean_slides(page, notebook_name):
    """한글 슬라이드 생성 (강제)"""
    log(f"  슬라이드 생성 (한글, 15장)...")

    # 1. Create 버튼 찾기
    create_btn = await page.query_selector(
        'button:has-text("Create"), button:has-text("만들기"), '
        'button[aria-label*="Create"], button[aria-label*="Add"]'
    )
    if create_btn:
        await create_btn.click()
        await asyncio.sleep(1)

    # 2. Slides/슬라이드 옵션 선택
    slides_option = await page.query_selector(
        '[role="menuitem"]:has-text("Slides"), '
        '[role="menuitem"]:has-text("슬라이드"), '
        '[role="menuitem"]:has-text("Presentation"), '
        'button:has-text("Slides")'
    )
    if slides_option:
        await slides_option.click()
        await asyncio.sleep(2)

    # 3. 언어 설정 - 한국어로 강제
    log("  언어 설정: 한국어...")

    # 언어 드롭다운 찾기
    lang_elements = await page.query_selector_all(
        '[aria-label*="language"], [aria-label*="Language"], '
        'select, [role="combobox"], [role="listbox"]'
    )

    for lang_el in lang_elements:
        try:
            await lang_el.click()
            await asyncio.sleep(0.5)

            # 한국어 옵션 찾기
            korean_option = await page.query_selector(
                '[role="option"]:has-text("한국어"), '
                '[role="option"]:has-text("Korean"), '
                'option:has-text("한국어")'
            )
            if korean_option:
                await korean_option.click()
                log("  ✓ 한국어 선택됨")
                await asyncio.sleep(0.5)
                break
        except:
            continue

    # 4. 생성 버튼 클릭
    generate_btn = await page.query_selector(
        'button:has-text("Generate"), button:has-text("생성"), '
        'button:has-text("Create"), button[type="submit"]'
    )
    if generate_btn:
        await generate_btn.click()
        log("  생성 시작...")

    return True


async def process_single_notebook(page, notebook_name, idx, total):
    """단일 노트북 처리 (에이전트 협력)"""
    log(f"\n{'='*50}")
    log(f"[{idx}/{total}] {notebook_name} 처리 중...")
    log('='*50)

    result = {
        'name': notebook_name,
        'notebook_id': None,
        'pdf': None,
        'pptx': None,
        'slides': 0
    }

    try:
        # 홈으로 이동
        await page.goto("https://notebooklm.google.com/", wait_until='networkidle', timeout=30000)
        await asyncio.sleep(3)

        # 노트북 찾기
        log(f"  노트북 찾는 중...")
        notebook_link = None

        # 모든 요소에서 검색
        all_elements = await page.query_selector_all('a, div[role="button"], div[class*="notebook"]')
        for el in all_elements:
            try:
                text = await el.inner_text()
                if notebook_name in text:
                    notebook_link = el
                    break
            except:
                continue

        if not notebook_link:
            # 텍스트로 직접 찾기
            notebook_link = await page.query_selector(f'text="{notebook_name}"')

        if not notebook_link:
            log(f"  ❌ '{notebook_name}' 노트북을 찾을 수 없음")
            await debug_agent_screenshot(page, f"notfound_{notebook_name}")
            return result

        # 노트북 클릭
        await notebook_link.click()
        await asyncio.sleep(5)

        # URL에서 ID 추출
        current_url = page.url
        match = re.search(r'/notebook/([a-zA-Z0-9_-]+)', current_url)
        result['notebook_id'] = match.group(1) if match else None
        log(f"  노트북 ID: {result['notebook_id'][:12] if result['notebook_id'] else 'N/A'}...")

        # 한글 슬라이드 생성
        await create_korean_slides(page, notebook_name)

        # 생성 완료 체크 (모니터링 에이전트)
        slides_ready = await check_slides_created(page, timeout=180)

        if not slides_ready:
            log("  ⚠️ 슬라이드 생성 타임아웃 - 그래도 다운로드 시도...")
            await debug_agent_screenshot(page, f"timeout_{notebook_name}")

        # PDF 다운로드 (헬퍼 에이전트)
        log(f"  PDF 다운로드 중...")
        pdf_path = await helper_agent_download(page, notebook_name)

        if pdf_path and pdf_path.exists():
            result['pdf'] = str(pdf_path)
            log(f"  ✓ PDF: {pdf_path.name}")

            # PPTX 변환하여 교체
            log(f"  PPTX 변환 중...")
            pptx_path, slide_count = pdf_to_pptx(pdf_path)
            result['pptx'] = str(pptx_path)
            result['slides'] = slide_count
            log(f"  ✓ PPTX: {pptx_path.name} ({slide_count}장)")
        else:
            log(f"  ❌ 다운로드 실패")
            await debug_agent_screenshot(page, f"download_fail_{notebook_name}")

    except Exception as e:
        log(f"  ❌ 오류 발생: {e}")
        await debug_agent_screenshot(page, f"error_{notebook_name}")

    return result


async def process_notebooks():
    """족부 질환 노트북 3개 처리 (멀티 에이전트)"""
    from playwright.async_api import async_playwright

    DOWNLOAD_DIR.mkdir(parents=True, exist_ok=True)
    BROWSER_PROFILE.mkdir(parents=True, exist_ok=True)

    notebooks_to_process = [
        "족저근막염: 원인, 증상, 치료 및 재활 가이드",
        "족모지외반증 (무지외반증): 원인, 증상, 치료 완벽 가이드",
        "발톱무좀 (조갑진균증): 원인, 증상, 치료 완벽 가이드"
    ]
    results = []

    async with async_playwright() as p:
        log("\n[메인] 브라우저 시작 (Chrome 채널 사용)...")

        # Chrome이 실행 중이면 새 프로필 사용
        try:
            browser = await p.chromium.launch(
                channel='chrome',
                headless=False,
                args=['--disable-blink-features=AutomationControlled'],
            )
            context = await browser.new_context(
                viewport={'width': 1920, 'height': 1080},
                accept_downloads=True,
            )
            context._downloads_path = str(DOWNLOAD_DIR)
            log("  Chrome 브라우저 사용")
        except:
            # 실패하면 persistent context 사용
            log("  Chrome 실패, Chromium 프로필 사용...")
            BROWSER_PROFILE.mkdir(parents=True, exist_ok=True)
            context = await p.chromium.launch_persistent_context(
                user_data_dir=str(AUTH_DIR / "browser_profile"),
                headless=False,
                downloads_path=str(DOWNLOAD_DIR),
                accept_downloads=True,
                args=['--disable-blink-features=AutomationControlled'],
                viewport={'width': 1920, 'height': 1080},
            )

        page = context.pages[0] if context.pages else await context.new_page()

        log("[메인] NotebookLM 접속...")
        await page.goto("https://notebooklm.google.com/", wait_until='networkidle', timeout=60000)
        await asyncio.sleep(3)

        # 로그인 확인
        if "accounts.google.com" in page.url:
            log("  ⚠️ 로그인 필요 - 브라우저에서 로그인해주세요...")

            for i in range(60):  # 5분 대기
                await asyncio.sleep(5)
                elapsed = (i + 1) * 5
                log(f"  로그인 대기 중... {elapsed}초")
                if "notebooklm.google.com" in page.url and "accounts.google" not in page.url:
                    log("  ✓ 로그인 감지됨!")
                    break
            else:
                log("  ❌ 로그인 타임아웃 (5분)")
                await context.close()
                return []

        log("  ✓ NotebookLM 접속 완료")

        # 각 노트북 처리
        for idx, notebook_name in enumerate(notebooks_to_process, 1):
            result = await process_single_notebook(page, notebook_name, idx, len(notebooks_to_process))
            results.append(result)
            await asyncio.sleep(2)

        await context.close()

    return results


async def main():
    log("=" * 60)
    log("노트랑 - 족부 질환 슬라이드 생성")
    log("  - 족저근막염")
    log("  - 족모지외반증 (무지외반증)")
    log("  - 발톱무좀 (조갑진균증)")
    log("  → 각 한글 슬라이드 생성")
    log("  → G:\\내 드라이브\\notebooklm 저장")
    log("  → PPTX 변환")
    log("=" * 60)

    results = await process_notebooks()

    log("\n" + "=" * 60)
    log("완료 결과:")
    log("=" * 60)

    success_count = 0
    for r in results:
        status = "✓" if r.get('pptx') else "❌"
        log(f"  {status} {r['name']}")
        if r.get('pptx'):
            log(f"      PDF:  {Path(r['pdf']).name}")
            log(f"      PPTX: {Path(r['pptx']).name} ({r['slides']}장)")
            success_count += 1

    log(f"\n성공: {success_count}/3")
    log(f"저장 위치: {DOWNLOAD_DIR}")
    log("=" * 60)

    return results


if __name__ == "__main__":
    asyncio.run(main())
