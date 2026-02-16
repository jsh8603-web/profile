#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
병렬 노트북 생성 - 여러 주제를 동시에 처리
"""
import asyncio
import json
import subprocess
import sys
import os
from pathlib import Path
from datetime import datetime
import time

if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8')

# 설정
NLM_EXE = Path.home() / "AppData/Roaming/Python/Python313/Scripts/nlm.exe"
DOWNLOAD_DIR = Path("G:/내 드라이브/notebooklm")
AUTH_DIR = Path.home() / ".notebooklm-mcp-cli"


def sync_auth():
    """인증 동기화"""
    root_auth = AUTH_DIR / "auth.json"
    profile_dir = AUTH_DIR / "profiles" / "default"

    if not root_auth.exists():
        return False

    with open(root_auth) as f:
        root_data = json.load(f)

    cookies_dict = root_data.get('cookies', {})
    cookies_list = [
        {"name": n, "value": v, "domain": ".google.com", "path": "/",
         "expires": -1, "httpOnly": False, "secure": True, "sameSite": "Lax"}
        for n, v in cookies_dict.items()
    ]

    profile_dir.mkdir(parents=True, exist_ok=True)

    with open(profile_dir / "cookies.json", "w") as f:
        json.dump(cookies_list, f)

    with open(profile_dir / "metadata.json", "w") as f:
        json.dump({
            "csrf_token": root_data.get("csrf_token", ""),
            "session_id": root_data.get("session_id", ""),
            "last_validated": datetime.now().isoformat()
        }, f)

    with open(profile_dir / "auth.json", "w") as f:
        json.dump(root_data, f)

    return True


def run_nlm(args, timeout=120):
    """nlm CLI 실행"""
    cmd = [str(NLM_EXE)] + args
    env = os.environ.copy()
    env['PYTHONIOENCODING'] = 'utf-8'

    try:
        result = subprocess.run(cmd, capture_output=True, timeout=timeout, env=env)
        stdout = result.stdout.decode('utf-8', errors='replace') if result.stdout else ''
        stderr = result.stderr.decode('utf-8', errors='replace') if result.stderr else ''
        return result.returncode == 0, stdout, stderr
    except Exception as e:
        return False, '', str(e)


async def create_notebook_workflow(title: str, research_queries: list, task_id: int):
    """단일 노트북 워크플로우"""
    prefix = f"[{task_id}] {title[:10]}"
    print(f"{prefix}: 시작")

    try:
        # 1. 노트북 생성
        print(f"{prefix}: 노트북 생성...")
        success, stdout, _ = run_nlm(["notebook", "create", title])

        notebook_id = None
        if success:
            try:
                # JSON 파싱 시도
                data = json.loads(stdout)
                notebook_id = data.get('id')
            except:
                # 텍스트에서 ID 추출 시도
                for line in stdout.split('\n'):
                    if 'id' in line.lower():
                        parts = line.split(':')
                        if len(parts) > 1:
                            notebook_id = parts[1].strip().strip('"').strip(',')
                            break

        if not notebook_id:
            print(f"{prefix}: ❌ 노트북 생성 실패")
            return None

        print(f"{prefix}: 노트북 ID: {notebook_id[:8]}...")

        # 2. 연구 및 소스 추가
        total_sources = 0
        for query in research_queries:
            print(f"{prefix}: 연구 '{query[:15]}...'")

            success, stdout, _ = run_nlm([
                "research", "start", query,
                "--notebook-id", notebook_id,
                "--mode", "fast"
            ], timeout=60)

            if not success:
                continue

            # Task ID 추출
            task_id_str = None
            for line in stdout.split('\n'):
                if 'Task ID:' in line:
                    task_id_str = line.split('Task ID:')[1].strip()
                    break

            if task_id_str:
                # 완료 대기
                for _ in range(12):  # 최대 1분
                    await asyncio.sleep(5)
                    success, stdout, _ = run_nlm(["research", "status", notebook_id])
                    if "completed" in stdout.lower():
                        break

                # 가져오기
                success, stdout, _ = run_nlm(["research", "import", notebook_id, task_id_str])
                if "Imported" in stdout:
                    try:
                        count = int(stdout.split("Imported")[1].split("source")[0].strip())
                        total_sources += count
                    except:
                        pass

        print(f"{prefix}: 소스 {total_sources}개 추가됨")

        # 3. 슬라이드 생성 (15장)
        print(f"{prefix}: 슬라이드 생성...")
        success, stdout, _ = run_nlm([
            "slides", "create", notebook_id,
            "--language", "ko",
            "--focus", f"{title}의 병인, 진단, 치료, 재활",
            "--length", "default",
            "--confirm"
        ], timeout=60)

        # 완료 대기
        print(f"{prefix}: 슬라이드 생성 대기...")
        for i in range(36):  # 최대 3분
            await asyncio.sleep(5)
            success, stdout, _ = run_nlm(["studio", "status", notebook_id])
            if '"status": "completed"' in stdout:
                print(f"{prefix}: 슬라이드 생성 완료!")
                break

        # 4. 다운로드
        print(f"{prefix}: 다운로드...")
        pdf_path = await download_single(notebook_id, title)

        if pdf_path:
            # 5. PPTX 변환
            print(f"{prefix}: PPTX 변환...")
            pptx_path = pdf_to_pptx(pdf_path)
            print(f"{prefix}: ✓ 완료: {pptx_path.name}")
            return {
                'title': title,
                'notebook_id': notebook_id,
                'pdf': str(pdf_path),
                'pptx': str(pptx_path)
            }
        else:
            print(f"{prefix}: ⚠️ 다운로드 실패")
            return {'title': title, 'notebook_id': notebook_id, 'pdf': None, 'pptx': None}

    except Exception as e:
        print(f"{prefix}: ❌ 오류: {e}")
        return None


async def download_single(notebook_id: str, title: str):
    """단일 노트북 다운로드"""
    from playwright.async_api import async_playwright

    DOWNLOAD_DIR.mkdir(parents=True, exist_ok=True)
    user_data_dir = AUTH_DIR / "browser_profile"

    # 파일명 생성
    safe_title = "".join(c for c in title if c.isalnum() or c in (' ', '_', '-')).strip()
    filename = f"{safe_title}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"

    async with async_playwright() as p:
        context = await p.chromium.launch_persistent_context(
            user_data_dir=str(user_data_dir),
            headless=True,
            downloads_path=str(DOWNLOAD_DIR),
            accept_downloads=True,
            args=['--disable-blink-features=AutomationControlled'],
            viewport={'width': 1920, 'height': 1080},
        )

        page = context.pages[0] if context.pages else await context.new_page()

        try:
            await page.goto(
                f"https://notebooklm.google.com/notebook/{notebook_id}",
                wait_until='domcontentloaded',
                timeout=30000
            )
        except:
            pass

        await asyncio.sleep(8)

        downloaded_path = None
        menu_btns = await page.query_selector_all('[aria-haspopup="menu"], button[aria-label*="more"]')

        for menu_btn in menu_btns[-10:]:
            try:
                await menu_btn.click(force=True)
                await asyncio.sleep(1)

                dl_item = await page.query_selector(
                    '[role="menuitem"]:has-text("다운로드"), [role="menuitem"]:has-text("Download")'
                )
                if dl_item:
                    async with page.expect_download(timeout=30000) as download_info:
                        await dl_item.click()

                    download = await download_info.value
                    downloaded_path = DOWNLOAD_DIR / filename
                    await download.save_as(str(downloaded_path))
                    break
            except:
                await page.keyboard.press('Escape')

        await context.close()
        return downloaded_path


def pdf_to_pptx(pdf_path):
    """PDF를 PPTX로 변환"""
    import fitz
    from pptx import Presentation
    from pptx.util import Inches
    import io

    pdf_path = Path(pdf_path)
    output_path = pdf_path.with_suffix('.pptx')

    doc = fitz.open(pdf_path)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for page_num in range(len(doc)):
        page = doc[page_num]
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img_data = pix.tobytes("png")

        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            io.BytesIO(img_data),
            Inches(0), Inches(0),
            width=prs.slide_width,
            height=prs.slide_height
        )

    doc.close()
    prs.save(output_path)
    return output_path


async def run_parallel(topics: list):
    """여러 주제를 병렬로 처리"""
    print("=" * 60)
    print(f"병렬 노트북 생성: {len(topics)}개")
    print("=" * 60)

    # 인증 확인
    print("\n[준비] 인증 동기화...")
    sync_auth()

    success, stdout, _ = run_nlm(["login", "--check"])
    if not (success and "valid" in stdout.lower()):
        print("❌ 인증 실패")
        return []

    print("✓ 인증 유효\n")

    # 병렬 실행
    tasks = []
    for i, topic in enumerate(topics):
        title = topic['title']
        queries = topic['queries']
        task = create_notebook_workflow(title, queries, i + 1)
        tasks.append(task)

    print(f"[실행] {len(tasks)}개 작업 동시 시작...\n")
    results = await asyncio.gather(*tasks, return_exceptions=True)

    # 결과 정리
    print("\n" + "=" * 60)
    print("결과:")
    successful = []
    for r in results:
        if isinstance(r, dict) and r:
            print(f"  ✓ {r['title']}")
            if r.get('pptx'):
                print(f"    → {r['pptx']}")
            successful.append(r)
        elif isinstance(r, Exception):
            print(f"  ❌ 오류: {r}")
        else:
            print(f"  ⚠️ 실패: {r}")

    print("=" * 60)
    print(f"완료: {len(successful)}/{len(topics)}")

    return successful


if __name__ == "__main__":
    # 족관절 질환 3개
    topics = [
        {
            "title": "족관절 염좌",
            "queries": [
                "족관절 염좌 원인 병인",
                "족관절 염좌 진단 치료",
                "족관절 염좌 재활 운동"
            ]
        },
        {
            "title": "족관절 골관절염",
            "queries": [
                "족관절 골관절염 원인 병인",
                "족관절 골관절염 진단 치료",
                "족관절 골관절염 재활"
            ]
        },
        {
            "title": "족관절 골절",
            "queries": [
                "족관절 골절 원인 분류",
                "족관절 골절 수술 치료",
                "족관절 골절 재활 운동"
            ]
        }
    ]

    results = asyncio.run(run_parallel(topics))

    if results:
        print(f"\n저장 위치: {DOWNLOAD_DIR}")
