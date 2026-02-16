#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""슬라이드 한글화 - PDF 텍스트 추출 → 번역 → 노트 추가"""
import sys
import fitz
from pptx import Presentation
from pathlib import Path
from deep_translator import GoogleTranslator
import time

if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8')

def extract_text_from_pdf(pdf_path):
    """PDF에서 페이지별 텍스트 추출"""
    doc = fitz.open(pdf_path)
    pages_text = []
    
    for page in doc:
        text = page.get_text()
        pages_text.append(text.strip())
    
    doc.close()
    return pages_text

def translate_text(text, max_chunk=4500):
    """텍스트를 한글로 번역"""
    if not text or len(text) < 10:
        return ""
    
    translator = GoogleTranslator(source='en', target='ko')
    
    # 긴 텍스트는 분할
    if len(text) > max_chunk:
        chunks = [text[i:i+max_chunk] for i in range(0, len(text), max_chunk)]
        translated = []
        for chunk in chunks:
            try:
                result = translator.translate(chunk)
                translated.append(result)
                time.sleep(0.5)  # API 제한 방지
            except Exception as e:
                translated.append(f"[번역 오류: {e}]")
        return " ".join(translated)
    else:
        try:
            return translator.translate(text)
        except Exception as e:
            return f"[번역 오류: {e}]"

def add_notes_to_pptx(pptx_path, pdf_files, output_path):
    """PPTX 슬라이드에 번역된 노트 추가"""
    prs = Presentation(pptx_path)
    
    # 모든 PDF에서 텍스트 추출
    all_texts = []
    for pdf_file in pdf_files:
        print(f"텍스트 추출: {pdf_file.name}")
        texts = extract_text_from_pdf(pdf_file)
        all_texts.extend(texts)
    
    print(f"\n총 {len(all_texts)}페이지, {len(prs.slides)}슬라이드")
    
    # 번역 및 노트 추가
    for i, slide in enumerate(prs.slides):
        if i < len(all_texts) and all_texts[i]:
            print(f"\r번역 중: {i+1}/{len(prs.slides)}", end="", flush=True)
            
            original = all_texts[i][:2000]  # 최대 2000자
            translated = translate_text(original)
            
            # 노트 추가
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = f"[한글 번역]\n{translated}\n\n[원문]\n{original[:500]}..."
    
    print("\n\n저장 중...")
    prs.save(output_path)
    print(f"완료: {output_path}")

def main():
    downloads_dir = Path("downloads")
    
    # PDF 파일 (순서대로)
    pdf_files = sorted(downloads_dir.glob("*.pdf"))
    
    # 통합 PPTX
    merged_pptx = downloads_dir / "n8n_마스터_가이드_통합.pptx"
    output_pptx = downloads_dir / "n8n_마스터_가이드_한글.pptx"
    
    if not merged_pptx.exists():
        print("통합 PPTX 파일이 없습니다")
        return
    
    add_notes_to_pptx(merged_pptx, pdf_files, output_pptx)

if __name__ == "__main__":
    main()
