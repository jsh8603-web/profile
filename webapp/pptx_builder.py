#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
편집 가능한 PPTX 빌더
- 배경 이미지 + 텍스트박스로 구성
- 텍스트는 PowerPoint에서 편집 가능
"""
import io
import sys
from pathlib import Path
from typing import List, Optional, Tuple, Union

from PIL import Image

if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8')

from .ocr_engine import TextBlock


class EditablePPTXBuilder:
    """편집 가능한 PPTX 생성기"""

    def __init__(
        self,
        slide_width_inches: float = 13.333,
        slide_height_inches: float = 7.5,
        font_name: str = "맑은 고딕",
        default_font_size: int = 12,
        include_background: bool = True,
        background_opacity: int = 100
    ):
        """
        PPTX 빌더 초기화

        Args:
            slide_width_inches: 슬라이드 너비 (인치)
            slide_height_inches: 슬라이드 높이 (인치)
            font_name: 기본 폰트 이름
            default_font_size: 기본 폰트 크기 (pt)
            include_background: 배경 이미지 포함 여부
            background_opacity: 배경 투명도 (0-100%)
        """
        from pptx import Presentation
        from pptx.util import Inches

        self.slide_width_inches = slide_width_inches
        self.slide_height_inches = slide_height_inches
        self.font_name = font_name
        self.default_font_size = default_font_size
        self.include_background = include_background
        self.background_opacity = background_opacity

        # Presentation 생성
        self._prs = Presentation()
        self._prs.slide_width = Inches(slide_width_inches)
        self._prs.slide_height = Inches(slide_height_inches)

        # 빈 레이아웃
        self._blank_layout = self._prs.slide_layouts[6]

    def _bbox_to_pptx_coords(
        self,
        bbox: List[List[float]],
        page_size: Tuple[float, float]
    ) -> Tuple[float, float, float, float]:
        """
        OCR bbox를 PPTX 좌표(인치)로 변환

        Args:
            bbox: [[x1,y1], [x2,y2], [x3,y3], [x4,y4]]
            page_size: (너비, 높이) 픽셀

        Returns:
            (x, y, width, height) in inches
        """
        x_min = min(p[0] for p in bbox)
        y_min = min(p[1] for p in bbox)
        x_max = max(p[0] for p in bbox)
        y_max = max(p[1] for p in bbox)

        page_w, page_h = page_size

        # 비율 계산하여 인치로 변환
        x = (x_min / page_w) * self.slide_width_inches
        y = (y_min / page_h) * self.slide_height_inches
        w = ((x_max - x_min) / page_w) * self.slide_width_inches
        h = ((y_max - y_min) / page_h) * self.slide_height_inches

        # 약간의 여백 추가 (텍스트가 잘리지 않도록)
        padding = 0.05
        w += padding
        h += padding

        return x, y, w, h

    def _adjust_background_opacity(self, image: Image.Image) -> Image.Image:
        """
        배경 이미지 투명도 조절

        Args:
            image: 원본 이미지

        Returns:
            투명도 조절된 이미지
        """
        if self.background_opacity >= 100:
            return image

        # 흰색 배경과 합성
        alpha = self.background_opacity / 100.0

        # RGBA로 변환
        if image.mode != 'RGBA':
            image = image.convert('RGBA')

        # 흰색 배경 생성
        white_bg = Image.new('RGBA', image.size, (255, 255, 255, 255))

        # 블렌딩
        blended = Image.blend(white_bg, image, alpha)

        return blended.convert('RGB')

    def add_slide(
        self,
        background_image: Optional[Image.Image],
        text_blocks: List[TextBlock],
        page_size: Optional[Tuple[float, float]] = None,
        use_auto_font_size: bool = True,
        fixed_font_size: Optional[int] = None
    ):
        """
        슬라이드 추가

        Args:
            background_image: 배경 이미지 (None이면 배경 없음)
            text_blocks: 텍스트 블록 리스트
            page_size: 페이지 크기 (픽셀) - None이면 배경 이미지 크기 사용
            use_auto_font_size: OCR 추정 폰트 크기 사용 여부
            fixed_font_size: 고정 폰트 크기 (use_auto_font_size=False일 때)
        """
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        slide = self._prs.slides.add_slide(self._blank_layout)

        # 페이지 크기 결정
        if page_size is None and background_image is not None:
            page_size = background_image.size
        elif page_size is None:
            # 기본 크기 (150 DPI 기준)
            page_size = (
                int(self.slide_width_inches * 150),
                int(self.slide_height_inches * 150)
            )

        # 배경 이미지 추가
        if self.include_background and background_image is not None:
            # 투명도 조절
            bg_image = self._adjust_background_opacity(background_image)

            # 이미지를 바이트로 변환
            img_buffer = io.BytesIO()
            bg_image.save(img_buffer, format='PNG')
            img_buffer.seek(0)

            # 배경으로 추가
            slide.shapes.add_picture(
                img_buffer,
                Inches(0), Inches(0),
                width=self._prs.slide_width,
                height=self._prs.slide_height
            )

        # 텍스트 블록 추가
        for block in text_blocks:
            x, y, w, h = self._bbox_to_pptx_coords(block.bbox, page_size)

            # 폰트 크기 결정
            if use_auto_font_size:
                font_size = block.estimated_font_size
            else:
                font_size = fixed_font_size or self.default_font_size

            # 텍스트박스 생성
            textbox = slide.shapes.add_textbox(
                Inches(x), Inches(y),
                Inches(w), Inches(h)
            )

            tf = textbox.text_frame
            tf.word_wrap = False  # 자동 줄바꿈 비활성화 (원본 위치 유지)
            tf.auto_size = None

            p = tf.paragraphs[0]
            p.text = block.text
            p.font.name = self.font_name
            p.font.size = Pt(font_size)
            p.font.color.rgb = RGBColor(0, 0, 0)  # 검정색

    def add_slides_from_ocr_results(
        self,
        images: List[Image.Image],
        ocr_results: List[List[TextBlock]],
        use_auto_font_size: bool = True,
        fixed_font_size: Optional[int] = None,
        progress_callback=None
    ):
        """
        OCR 결과로 여러 슬라이드 추가

        Args:
            images: 배경 이미지 리스트
            ocr_results: 페이지별 OCR 결과
            use_auto_font_size: 자동 폰트 크기 사용 여부
            fixed_font_size: 고정 폰트 크기
            progress_callback: 진행률 콜백 (page_num, total)
        """
        total = len(images)

        for i, (image, blocks) in enumerate(zip(images, ocr_results)):
            if progress_callback:
                progress_callback(i + 1, total)

            self.add_slide(
                background_image=image,
                text_blocks=blocks,
                use_auto_font_size=use_auto_font_size,
                fixed_font_size=fixed_font_size
            )

    def save(self, output_path: Union[Path, str]) -> Path:
        """
        PPTX 파일 저장

        Args:
            output_path: 출력 파일 경로

        Returns:
            저장된 파일 경로
        """
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        self._prs.save(output_path)
        return output_path

    def save_to_bytes(self) -> bytes:
        """
        PPTX를 바이트로 반환

        Returns:
            PPTX 바이트 데이터
        """
        buffer = io.BytesIO()
        self._prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

    @property
    def slide_count(self) -> int:
        """현재 슬라이드 수"""
        return len(self._prs.slides)


def create_editable_pptx(
    pdf_images: List[Image.Image],
    ocr_results: List[List[TextBlock]],
    output_path: Union[Path, str],
    font_name: str = "맑은 고딕",
    use_auto_font_size: bool = True,
    fixed_font_size: int = 12,
    include_background: bool = True,
    background_opacity: int = 100,
    progress_callback=None
) -> Path:
    """
    편집 가능한 PPTX 생성 (편의 함수)

    Args:
        pdf_images: PDF 페이지 이미지 리스트
        ocr_results: 페이지별 OCR 결과
        output_path: 출력 파일 경로
        font_name: 폰트 이름
        use_auto_font_size: 자동 폰트 크기 사용
        fixed_font_size: 고정 폰트 크기
        include_background: 배경 이미지 포함
        background_opacity: 배경 투명도 (0-100%)
        progress_callback: 진행률 콜백

    Returns:
        저장된 파일 경로
    """
    builder = EditablePPTXBuilder(
        font_name=font_name,
        default_font_size=fixed_font_size,
        include_background=include_background,
        background_opacity=background_opacity
    )

    builder.add_slides_from_ocr_results(
        images=pdf_images,
        ocr_results=ocr_results,
        use_auto_font_size=use_auto_font_size,
        fixed_font_size=fixed_font_size,
        progress_callback=progress_callback
    )

    return builder.save(output_path)
